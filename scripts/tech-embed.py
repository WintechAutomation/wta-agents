"""tech-embed.py — 프로젝트 기술문서 파싱 + 임베딩 파이프라인.

data/tech/{project_id}/{category}/ 하위의 PDF, XLSX, DOCX, PPTX 파일을
파싱 → 청킹 → Qwen3-Embedding-8B 임베딩 → manual.tech_documents 테이블에 저장.

실행:
  python tech-embed.py                          # 전체 증분 처리
  python tech-embed.py --dry-run                # 파싱만 (임베딩/DB 저장 X)
  python tech-embed.py --no-embed               # DB 저장만 (임베딩 X)
  python tech-embed.py --project 30             # 특정 프로젝트만
  python tech-embed.py --category specification # 특정 카테고리만
  python tech-embed.py --full                   # 전체 재처리 (해시 무시)

필수 패키지:
  pip install pymupdf openpyxl python-docx python-pptx psycopg2-binary requests
"""

import argparse
import hashlib
import json
import logging
import os
import re
import sys
import time
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

import psycopg2
from psycopg2.extras import execute_values
import requests

# ── 설정 ──
BASE_DIR = os.path.normpath(os.path.join(os.path.dirname(__file__), ".."))
TECH_DIR = os.path.join(BASE_DIR, "data", "tech")
MANIFEST_FILE = os.path.join(TECH_DIR, "file-manifest.json")

EMBED_URL = "http://182.224.6.147:11434/api/embed"
EMBED_MODEL = "qwen3-embedding:8b"
EMBED_DIM = 2000
EMBED_BATCH = 16
EMBED_DELAY = 0.5

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100

DB_TABLE = "manual.tech_documents"
DB_CONFIG = {
    "host": "localhost",
    "port": 55432,
    "user": "postgres",
    "password": "your-super-secret-and-long-postgres-password",
    "dbname": "postgres",
}

SUPPORTED_EXTS = {".pdf", ".xlsx", ".docx", ".pptx"}

logging.basicConfig(
    level=logging.INFO,
    format="[tech-embed] %(asctime)s %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("tech-embed")


# ── manifest에서 project_code 매핑 로드 ──

def load_manifest() -> dict[str, dict]:
    """file-manifest.json에서 project_id → project_code/project_name 매핑 생성."""
    mapping: dict[str, dict] = {}
    if not os.path.isfile(MANIFEST_FILE):
        log.warning("file-manifest.json 없음 — project_code 매핑 불가")
        return mapping

    with open(MANIFEST_FILE, "r", encoding="utf-8") as f:
        manifest = json.load(f)

    for entry in manifest.get("files", []):
        pid = str(entry.get("project_id", ""))
        if pid and pid not in mapping:
            mapping[pid] = {
                "project_code": entry.get("project_code", ""),
                "project_name": entry.get("project_name", ""),
            }
    return mapping


# ── 파일 해시 ──

def compute_file_hash(file_path: str) -> str:
    """파일 SHA-256 해시 앞 16자리."""
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for block in iter(lambda: f.read(8192), b""):
            h.update(block)
    return h.hexdigest()[:16]


# ── 텍스트 청킹 ──

def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """단락 기반 텍스트 청킹."""
    paragraphs = re.split(r"\n{2,}", text)
    chunks: list[str] = []
    current = ""
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if len(current) + len(para) + 1 <= size:
            current = (current + "\n\n" + para).strip()
        else:
            if current:
                chunks.append(current)
            if len(para) > size:
                for j in range(0, len(para), size - overlap):
                    chunks.append(para[j:j + size])
                current = ""
            else:
                current = para
    if current:
        chunks.append(current)
    return chunks


# ── 파일 유형별 파싱 ──

def parse_pdf(file_path: str) -> list[dict]:
    """PyMuPDF로 PDF 텍스트+표 추출. 이미지 전용은 경고만."""
    try:
        import fitz
    except ImportError:
        log.error("pymupdf 미설치. 실행: pip install pymupdf")
        return []

    try:
        doc = fitz.open(file_path)
    except Exception as e:
        log.warning(f"  PDF 열기 실패: {e}")
        return []

    all_chunks: list[dict] = []
    total_text = ""

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_num = page_idx + 1

        # 텍스트 추출
        text = page.get_text("text").replace("\x00", "").strip()
        if text:
            total_text += text + "\n\n"

        # 표 추출 (PyMuPDF 내장)
        try:
            tables = page.find_tables()
            for table in tables:
                table_data = table.extract()
                if table_data:
                    # 테이블을 마크다운 형태로 변환
                    rows = []
                    for row in table_data:
                        cells = [str(c).strip() if c else "" for c in row]
                        rows.append(" | ".join(cells))
                    table_text = "\n".join(rows)
                    if len(table_text.strip()) > 20:
                        all_chunks.append({
                            "content": table_text,
                            "chunk_type": "table",
                            "page_number": page_num,
                        })
        except Exception:
            pass  # 표 추출 실패 시 무시

    doc.close()

    # 텍스트가 거의 없는 PDF (이미지 전용)
    if len(total_text.strip()) < 100:
        log.warning(f"  이미지 전용 PDF — 텍스트 부족 ({len(total_text.strip())}자), 스킵")
        return []

    # 텍스트 청킹
    text_chunks = chunk_text(total_text)
    for chunk in text_chunks:
        all_chunks.append({
            "content": chunk,
            "chunk_type": "text",
            "page_number": None,
        })

    return all_chunks


def parse_xlsx(file_path: str) -> list[dict]:
    """openpyxl로 시트별 셀 데이터 → 테이블 형태 텍스트."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        log.error("openpyxl 미설치. 실행: pip install openpyxl")
        return []

    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
    except Exception as e:
        log.warning(f"  XLSX 열기 실패: {e}")
        return []

    all_chunks: list[dict] = []

    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
        ws = wb[sheet_name]
        rows: list[str] = []

        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            # 빈 행 스킵
            if not any(cells):
                continue
            rows.append(" | ".join(cells))

        if not rows:
            continue

        sheet_text = f"[시트: {sheet_name}]\n" + "\n".join(rows)

        # 청킹
        chunks = chunk_text(sheet_text)
        for chunk in chunks:
            all_chunks.append({
                "content": chunk,
                "chunk_type": "table",
                "page_number": sheet_idx + 1,
            })

    wb.close()
    return all_chunks


def parse_docx(file_path: str) -> list[dict]:
    """python-docx로 텍스트+테이블 추출."""
    try:
        from docx import Document
    except ImportError:
        log.error("python-docx 미설치. 실행: pip install python-docx")
        return []

    try:
        doc = Document(file_path)
    except Exception as e:
        log.warning(f"  DOCX 열기 실패: {e}")
        return []

    all_chunks: list[dict] = []

    # 본문 텍스트
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    full_text = "\n\n".join(paragraphs)

    if full_text:
        for chunk in chunk_text(full_text):
            all_chunks.append({
                "content": chunk,
                "chunk_type": "text",
                "page_number": None,
            })

    # 테이블
    for table_idx, table in enumerate(doc.tables):
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        table_text = "\n".join(rows)
        if len(table_text.strip()) > 20:
            all_chunks.append({
                "content": table_text,
                "chunk_type": "table",
                "page_number": None,
            })

    return all_chunks


def parse_pptx(file_path: str) -> list[dict]:
    """python-pptx로 슬라이드별 텍스트+표 추출."""
    try:
        from pptx import Presentation
    except ImportError:
        log.error("python-pptx 미설치. 실행: pip install python-pptx")
        return []

    try:
        prs = Presentation(file_path)
    except Exception as e:
        log.warning(f"  PPTX 열기 실패: {e}")
        return []

    all_chunks: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        texts: list[str] = []

        for shape in slide.shapes:
            # 텍스트 프레임
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

            # 테이블
            if shape.has_table:
                rows: list[str] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                table_text = "\n".join(rows)
                if len(table_text.strip()) > 20:
                    all_chunks.append({
                        "content": table_text,
                        "chunk_type": "table",
                        "page_number": slide_num,
                    })

        slide_text = "\n".join(texts)
        if slide_text.strip():
            for chunk in chunk_text(slide_text):
                all_chunks.append({
                    "content": chunk,
                    "chunk_type": "text",
                    "page_number": slide_num,
                })

    return all_chunks


PARSERS = {
    ".pdf": parse_pdf,
    ".xlsx": parse_xlsx,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
}


# ── 임베딩 ──

def embed_texts(texts: list[str]) -> list[list[float]]:
    """Ollama Qwen3-Embedding 서버 배치 임베딩."""
    payload = {"model": EMBED_MODEL, "input": texts}
    resp = requests.post(EMBED_URL, json=payload, timeout=300)
    resp.raise_for_status()
    data = resp.json()
    if "embeddings" not in data:
        raise ValueError(f"임베딩 응답 오류: {data}")
    # Matryoshka: 4096→2000 차원 잘라서 반환
    return [v[:EMBED_DIM] for v in data["embeddings"]]


# ── DB ──

def ensure_schema(conn) -> None:
    """manual.tech_documents 테이블 및 인덱스 생성."""
    with conn.cursor() as cur:
        cur.execute("CREATE SCHEMA IF NOT EXISTS manual")
        cur.execute(f"""CREATE TABLE IF NOT EXISTS {DB_TABLE} (
            id SERIAL PRIMARY KEY,
            source_file TEXT NOT NULL,
            file_hash TEXT NOT NULL,
            project_id TEXT,
            project_code TEXT,
            category TEXT NOT NULL DEFAULT '',
            chunk_index INTEGER NOT NULL,
            chunk_type TEXT NOT NULL DEFAULT 'text',
            page_number INTEGER,
            content TEXT NOT NULL,
            metadata JSONB DEFAULT '{{}}'::jsonb,
            embedding vector({EMBED_DIM}),
            created_at TIMESTAMPTZ DEFAULT now(),
            updated_at TIMESTAMPTZ DEFAULT now(),
            UNIQUE (source_file, chunk_index)
        )""")
        cur.execute(f"""CREATE INDEX IF NOT EXISTS tech_documents_embedding_idx
            ON {DB_TABLE} USING hnsw (embedding vector_cosine_ops)""")
    conn.commit()


def get_existing_hashes(conn) -> dict[str, str]:
    """기존 파일 해시 조회 (증분 처리용)."""
    with conn.cursor() as cur:
        cur.execute(f"SELECT DISTINCT source_file, file_hash FROM {DB_TABLE}")
        return {row[0]: row[1] for row in cur.fetchall()}


def delete_file_chunks(conn, source_file: str) -> None:
    """특정 파일의 기존 청크 삭제 (재처리용)."""
    with conn.cursor() as cur:
        cur.execute(f"DELETE FROM {DB_TABLE} WHERE source_file = %s", (source_file,))
    conn.commit()


def upsert_chunks(conn, chunks_data: list[dict]) -> None:
    """청크 데이터 DB upsert."""
    if not chunks_data:
        return
    sql = f"""INSERT INTO {DB_TABLE}
        (source_file, file_hash, project_id, project_code, category,
         chunk_index, chunk_type, page_number, content, metadata, embedding, updated_at)
        VALUES %s
        ON CONFLICT (source_file, chunk_index) DO UPDATE SET
            file_hash=EXCLUDED.file_hash, project_id=EXCLUDED.project_id,
            project_code=EXCLUDED.project_code, category=EXCLUDED.category,
            chunk_type=EXCLUDED.chunk_type, page_number=EXCLUDED.page_number,
            content=EXCLUDED.content, metadata=EXCLUDED.metadata,
            embedding=EXCLUDED.embedding, updated_at=now()"""
    template = (
        "(%(source_file)s, %(file_hash)s, %(project_id)s, %(project_code)s, %(category)s, "
        "%(chunk_index)s, %(chunk_type)s, %(page_number)s, %(content)s, "
        "%(metadata)s::jsonb, %(embedding)s::vector, now())"
    )
    with conn.cursor() as cur:
        execute_values(cur, sql, chunks_data, template=template, page_size=50)
    conn.commit()


# ── 파일 수집 ──

def collect_files(
    tech_dir: str,
    project_filter: str | None = None,
    category_filter: str | None = None,
) -> list[dict]:
    """data/tech/{project_id}/{category}/ 하위 파일 수집."""
    files: list[dict] = []

    if not os.path.isdir(tech_dir):
        log.error(f"디렉토리 없음: {tech_dir}")
        return files

    for project_id in sorted(os.listdir(tech_dir)):
        project_path = os.path.join(tech_dir, project_id)
        if not os.path.isdir(project_path):
            continue
        if project_filter and project_id != project_filter:
            continue

        for category in sorted(os.listdir(project_path)):
            category_path = os.path.join(project_path, category)
            if not os.path.isdir(category_path):
                continue
            if category_filter and category != category_filter:
                continue

            for filename in sorted(os.listdir(category_path)):
                ext = os.path.splitext(filename)[1].lower()
                if ext not in SUPPORTED_EXTS:
                    continue
                file_path = os.path.join(category_path, filename)
                # 상대 경로 (source_file 키)
                rel_path = f"{project_id}/{category}/{filename}"
                files.append({
                    "file_path": file_path,
                    "rel_path": rel_path,
                    "project_id": project_id,
                    "category": category,
                    "filename": filename,
                    "ext": ext,
                })

    return files


# ── 메인 ──

def main():
    parser = argparse.ArgumentParser(description="프로젝트 기술문서 파싱 + 임베딩")
    parser.add_argument("--dry-run", action="store_true", help="파싱만 (임베딩/DB 저장 X)")
    parser.add_argument("--no-embed", action="store_true", help="DB 저장만 (임베딩 X)")
    parser.add_argument("--project", type=str, help="특정 프로젝트 ID만 처리")
    parser.add_argument("--category", type=str, help="특정 카테고리만 처리")
    parser.add_argument("--full", action="store_true", help="전체 재처리 (해시 무시)")
    args = parser.parse_args()

    mode_label = "전체" if args.full else "증분"
    if args.dry_run:
        mode_label += " (dry-run)"
    elif args.no_embed:
        mode_label += " (no-embed)"
    log.info(f"시작 (모드: {mode_label})")

    # manifest 로드
    manifest = load_manifest()
    log.info(f"manifest: {len(manifest)}개 프로젝트 매핑")

    # 파일 수집
    files = collect_files(TECH_DIR, args.project, args.category)
    log.info(f"대상 파일: {len(files)}개")

    if not files:
        log.info("처리할 파일 없음 — 종료")
        return

    # DB 연결 (dry-run이 아닌 경우)
    conn = None
    existing_hashes: dict[str, str] = {}
    if not args.dry_run:
        conn = psycopg2.connect(**DB_CONFIG)
        ensure_schema(conn)
        if not args.full:
            existing_hashes = get_existing_hashes(conn)
            log.info(f"기존 DB 레코드: {len(existing_hashes)}개 파일")

    # 통계
    stats = {"processed": 0, "skipped": 0, "chunks": 0, "embedded": 0, "errors": 0}

    for file_idx, file_info in enumerate(files):
        file_path = file_info["file_path"]
        rel_path = file_info["rel_path"]
        ext = file_info["ext"]
        project_id = file_info["project_id"]
        category = file_info["category"]

        # project_code 매핑
        proj_info = manifest.get(project_id, {})
        project_code = proj_info.get("project_code", "")

        # 해시 계산 및 증분 체크
        file_hash = compute_file_hash(file_path)

        if not args.full and not args.dry_run:
            if rel_path in existing_hashes and existing_hashes[rel_path] == file_hash:
                stats["skipped"] += 1
                continue

        log.info(f"[{file_idx + 1}/{len(files)}] {rel_path}")

        # 파싱
        parser_fn = PARSERS.get(ext)
        if not parser_fn:
            log.warning(f"  미지원 확장자: {ext}")
            stats["errors"] += 1
            continue

        try:
            raw_chunks = parser_fn(file_path)
        except Exception as e:
            log.error(f"  파싱 실패: {e}")
            stats["errors"] += 1
            continue

        if not raw_chunks:
            log.info(f"  청크 없음 (빈 파일 또는 이미지 전용)")
            stats["skipped"] += 1
            continue

        log.info(f"  청크 {len(raw_chunks)}개 추출")
        stats["chunks"] += len(raw_chunks)

        if args.dry_run:
            stats["processed"] += 1
            continue

        # 기존 데이터 삭제 (재처리 시)
        if rel_path in existing_hashes:
            delete_file_chunks(conn, rel_path)

        # 메타데이터 구성
        meta_base = {
            "project_id": project_id,
            "project_code": project_code,
            "category": category,
            "file_extension": ext,
        }

        # 임베딩
        db_rows: list[dict] = []
        if args.no_embed:
            # 임베딩 없이 DB 저장
            for ci, chunk in enumerate(raw_chunks):
                meta = {**meta_base, "page_number": chunk.get("page_number"), "chunk_type": chunk["chunk_type"]}
                db_rows.append({
                    "source_file": rel_path,
                    "file_hash": file_hash,
                    "project_id": project_id,
                    "project_code": project_code,
                    "category": category,
                    "chunk_index": ci,
                    "chunk_type": chunk["chunk_type"],
                    "page_number": chunk.get("page_number"),
                    "content": chunk["content"],
                    "metadata": json.dumps(meta, ensure_ascii=False),
                    "embedding": None,
                })
        else:
            # 배치 임베딩
            for batch_start in range(0, len(raw_chunks), EMBED_BATCH):
                batch = raw_chunks[batch_start:batch_start + EMBED_BATCH]
                texts = [c["content"] for c in batch]

                try:
                    embeddings = embed_texts(texts)
                except Exception as e:
                    log.error(f"  임베딩 실패 (배치 {batch_start}~{batch_start + len(batch)}): {e}")
                    stats["errors"] += 1
                    # 실패한 배치는 임베딩 없이 저장
                    embeddings = [None] * len(batch)

                for ci_offset, (chunk, emb) in enumerate(zip(batch, embeddings)):
                    ci = batch_start + ci_offset
                    meta = {**meta_base, "page_number": chunk.get("page_number"), "chunk_type": chunk["chunk_type"]}
                    db_rows.append({
                        "source_file": rel_path,
                        "file_hash": file_hash,
                        "project_id": project_id,
                        "project_code": project_code,
                        "category": category,
                        "chunk_index": ci,
                        "chunk_type": chunk["chunk_type"],
                        "page_number": chunk.get("page_number"),
                        "content": chunk["content"],
                        "metadata": json.dumps(meta, ensure_ascii=False),
                        "embedding": str(emb) if emb else None,
                    })
                    if emb:
                        stats["embedded"] += 1

                if batch_start + EMBED_BATCH < len(raw_chunks):
                    time.sleep(EMBED_DELAY)

        # DB upsert
        try:
            upsert_chunks(conn, db_rows)
            stats["processed"] += 1
        except Exception as e:
            log.error(f"  DB 저장 실패: {e}")
            conn.rollback()
            stats["errors"] += 1

    # 결과 출력
    log.info("=" * 50)
    log.info(f"완료: 처리 {stats['processed']}개, 스킵 {stats['skipped']}개, "
             f"청크 {stats['chunks']}개, 임베딩 {stats['embedded']}개, 에러 {stats['errors']}개")

    if conn:
        conn.close()


if __name__ == "__main__":
    main()
