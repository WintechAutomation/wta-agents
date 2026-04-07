"""
프레임그레버 이탈 이슈 PPT 생성 스크립트
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

OUTPUT = r"C:\MES\wta-agents\reports\control-agent\framegrabber-issue.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# --- Helper ---
def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color=RGBColor(0x44, 0x44, 0x44), alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_shape_fill(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def emu(inches):
    return int(inches * 914400)

# --- Constants ---
SLIDE_W = emu(13.333)
SLIDE_H = emu(7.5)
MARGIN = emu(0.6)
RED = RGBColor(0xCC, 0x00, 0x00)
DARK = RGBColor(0x22, 0x22, 0x22)
BODY = RGBColor(0x44, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_YELLOW = RGBColor(0xFF, 0xF8, 0xE1)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
CRED = RGBColor(0xC6, 0x28, 0x28)
GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)

# --- Title Bar ---
title_bar = add_shape_fill(slide, 0, 0, SLIDE_W, emu(0.85), RED)
add_textbox(slide, emu(0.7), emu(0.15), emu(12), emu(0.55),
           "검사기 프레임그레버(Frame Grabber) 이탈 이슈",
           font_size=28, bold=True, color=WHITE)

# --- Info Tags ---
y_info = emu(1.1)
info_items = [
    ("장비: HIM-F 검사기", 0.6),
    ("부품: Frame Grabber (PCIe)", 3.1),
    ("현상: PCIe 슬롯 이탈", 5.9),
    ("담당: 지건승", 8.4),
    ("일자: 2026. 04.", 10.2),
]
for text, x in info_items:
    bg = add_shape_fill(slide, emu(x), y_info, emu(2.3), emu(0.35), GRAY_BG)
    add_textbox(slide, emu(x + 0.1), y_info + emu(0.02), emu(2.1), emu(0.3),
               text, font_size=11, color=BODY)

# --- Cause Analysis Box ---
y_cause = emu(1.65)
cause_bg = add_shape_fill(slide, MARGIN, y_cause, SLIDE_W - 2*MARGIN, emu(0.95), LIGHT_YELLOW)
add_textbox(slide, emu(0.8), y_cause + emu(0.05), emu(3), emu(0.25),
           "원인 분석", font_size=13, bold=True, color=ORANGE)

cause_texts = [
    "- 장비 운송/가동 중 진동에 의한 PCIe 슬롯 접촉 불량 발생",
    "- 프레임그레버 카드 고정 브라켓 체결 미흡 또는 나사 풀림",
    "- 반복 진동으로 카드가 점진적으로 슬롯에서 이탈 → 카메라 통신 두절",
]
cause_text = "\n".join(cause_texts)
txBox = slide.shapes.add_textbox(Emu(emu(0.8)), Emu(y_cause + emu(0.3)), Emu(emu(11.5)), Emu(emu(0.6)))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(cause_texts):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p.font.name = "맑은 고딕"
    p.space_after = Pt(2)

# --- Before / After Comparison ---
y_vs = emu(2.8)
col_w = (SLIDE_W - 2*MARGIN - emu(0.2)) // 2
col_h = emu(3.8)

# Before (Left)
before_bg = add_shape_fill(slide, MARGIN, y_vs, col_w, col_h, RGBColor(0xFF, 0xF5, 0xF5))
add_textbox(slide, emu(0.8), y_vs + emu(0.1), emu(5), emu(0.35),
           "⚠ 조치 전 (Before)", font_size=16, bold=True, color=CRED)

# Video placeholder - Before
vid_before = add_shape_fill(slide, emu(0.9), y_vs + emu(0.55), col_w - emu(0.6), emu(2.0), RGBColor(0xE0, 0xE0, 0xE0))
add_textbox(slide, emu(0.9), y_vs + emu(1.3), col_w - emu(0.6), emu(0.3),
           "🎬 조치 전 동영상 삽입", font_size=12, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

# Before descriptions
before_items = [
    "• 프레임그레버 PCIe 슬롯에서 이탈 확인",
    "• 카메라 이미지 취득 실패 / 간헐적 통신 끊김",
    "• 장비 가동 중 검사 중단 발생 (BSOD 또는 프리징)",
    "• 고정 브라켓 나사 풀림 / 미체결 상태",
]
txBox = slide.shapes.add_textbox(Emu(emu(0.9)), Emu(y_vs + emu(2.65)), Emu(col_w - emu(0.6)), Emu(emu(1.1)))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(before_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p.font.name = "맑은 고딕"
    p.space_after = Pt(2)

# After (Right)
after_x = MARGIN + col_w + emu(0.2)
after_bg = add_shape_fill(slide, after_x, y_vs, col_w, col_h, RGBColor(0xF0, 0xFF, 0xF0))
add_textbox(slide, after_x + emu(0.2), y_vs + emu(0.1), emu(5), emu(0.35),
           "✅ 조치 후 (After)", font_size=16, bold=True, color=GREEN)

# Video placeholder - After
vid_after = add_shape_fill(slide, after_x + emu(0.3), y_vs + emu(0.55), col_w - emu(0.6), emu(2.0), RGBColor(0xE0, 0xE0, 0xE0))
add_textbox(slide, after_x + emu(0.3), y_vs + emu(1.3), col_w - emu(0.6), emu(0.3),
           "🎬 조치 후 동영상 삽입", font_size=12, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

# After descriptions
after_items = [
    "• 프레임그레버 카드 재장착 및 슬롯 완전 삽입 확인",
    "• 고정 브라켓 나사 체결 + 풀림 방지 조치 적용",
    "• 카메라 이미지 취득 정상 동작 확인",
    "• 연속 가동 테스트 완료 (이탈 재발 없음)",
]
txBox = slide.shapes.add_textbox(Emu(after_x + emu(0.3)), Emu(y_vs + emu(2.65)), Emu(col_w - emu(0.6)), Emu(emu(1.1)))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(after_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p.font.name = "맑은 고딕"
    p.space_after = Pt(2)

# --- Result Bar ---
y_result = y_vs + col_h + emu(0.15)
result_bg = add_shape_fill(slide, MARGIN, y_result, SLIDE_W - 2*MARGIN, emu(0.4), LIGHT_GREEN)
add_textbox(slide, emu(0.8), y_result + emu(0.05), emu(11.5), emu(0.3),
           "✔ 조치 결과: 프레임그레버 재장착 및 고정 강화 → 카메라 통신 정상 복구, 연속 가동 안정성 확인 완료",
           font_size=12, bold=False, color=GREEN)

# --- Footer ---
add_textbox(slide, 0, SLIDE_H - emu(0.35), SLIDE_W, emu(0.3),
           "WTA Automation | Control Design Team",
           font_size=9, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)

# Save
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"PPT saved: {OUTPUT}")
