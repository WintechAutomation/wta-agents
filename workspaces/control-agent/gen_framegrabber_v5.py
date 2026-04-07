# -*- coding: utf-8 -*-
"""
프레임그레버 이탈 이슈 PPT v5
기본양식의 슬라이드 4(1_Title Only)를 유지, 한 장으로 구성
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

TEMPLATE = r"C:\MES\wta-agents\workspaces\control-agent\template_base.pptx"
OUTPUT = r"C:\MES\wta-agents\reports\control-agent\framegrabber-issue-v2.pptx"

prs = Presentation(TEMPLATE)

# 슬라이드 1,2,3 삭제 → 슬라이드 4만 남김
for _ in range(3):
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

slide = prs.slides[0]  # 원래 슬라이드 4

# 기존 TextBox(Servo, RTX 등) 제거
sp_tree = slide.shapes._spTree
for shape in list(slide.shapes):
    if not shape.is_placeholder:
        sp_tree.remove(shape._element)

# placeholder 내용 교체
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:  # 제목
        ph.text = ""
        p = ph.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = "검사기 프레임그레버(Frame Grabber) 이탈 이슈"
    elif idx == 14:  # 카테고리
        ph.text = ""
        p = ph.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = "[제어설계]"
        r.font.bold = True
    elif idx == 12:  # 페이지 번호
        ph.text = ""
        p = ph.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = "1"

# --- 헬퍼 ---
def emu(inches):
    return int(inches * 914400)

def add_tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame

def add_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    return sh

def ar(p, text, size=10, bold=False, color=None, font="맑은 고딕"):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return r

C_RED = RGBColor(0xC6, 0x28, 0x28)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_ORANGE = RGBColor(0xE6, 0x51, 0x00)
C_DARK = RGBColor(0x22, 0x22, 0x22)
C_BODY = RGBColor(0x44, 0x44, 0x44)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_LIGHT = RGBColor(0x88, 0x88, 0x88)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

LM = emu(1.75)
CW = emu(11.6) - LM  # ~9.85 inches

# =====================================================
# 한 장에 모든 내용 구성
# =====================================================

# --- 정보 바 ---
y = emu(1.3)
add_rect(slide, LM, y, CW, emu(0.28), RGBColor(0xF5, 0xF5, 0xF5), RGBColor(0xE0, 0xE0, 0xE0))
tf = add_tb(slide, LM + emu(0.1), y + emu(0.02), CW - emu(0.2), emu(0.24))
p = tf.paragraphs[0]
for t, b in [("장비: ",1),("HIM-F 검사기   ",0),("부품: ",1),("Frame Grabber (PCIe)   ",0),
             ("현상: ",1),("PCIe 슬롯 이탈   ",0),("담당: ",1),("지건승   ",0),("일자: ",1),("2026. 04.",0)]:
    ar(p, t, size=9, bold=bool(b), color=C_DARK if b else C_BODY)

# --- 원인 분석 ---
y = emu(1.68)
add_rect(slide, LM, y, CW, emu(0.65), RGBColor(0xFF, 0xF8, 0xE1), RGBColor(0xFF, 0xE0, 0x82))

tf = add_tb(slide, LM + emu(0.1), y + emu(0.03), emu(2), emu(0.18))
p = tf.paragraphs[0]
ar(p, "원인 분석", size=10, bold=True, color=C_ORANGE)

tf = add_tb(slide, LM + emu(0.1), y + emu(0.22), CW - emu(0.2), emu(0.4))
for i, text in enumerate([
    "• 장비 운송/가동 중 진동에 의한 PCIe 슬롯 접촉 불량 발생",
    "• 프레임그레버 카드 고정 브라켓 체결 미흡 또는 나사 풀림",
    "• 반복 진동으로 카드가 점진적으로 슬롯에서 이탈 → 카메라 통신 두절",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    ar(p, text, size=8.5, color=RGBColor(0x55, 0x55, 0x55))
    p.space_after = Pt(0)

# --- Before / After ---
y_vs = emu(2.45)
gap = emu(0.1)
col_w = (CW - gap) // 2
col_h = emu(3.55)

# -- Before --
add_rect(slide, LM, y_vs, col_w, col_h, RGBColor(0xFF, 0xF5, 0xF5), RGBColor(0xEF, 0x9A, 0x9A))

# 제목 바
add_rect(slide, LM, y_vs, col_w, emu(0.3), C_RED)
tf = add_tb(slide, LM + emu(0.1), y_vs + emu(0.03), col_w - emu(0.2), emu(0.24))
p = tf.paragraphs[0]
ar(p, "⚠  조치 전 (Before)", size=12, bold=True, color=C_WHITE)

# 동영상
vy = y_vs + emu(0.38)
vh = emu(1.55)
add_rect(slide, LM + emu(0.1), vy, col_w - emu(0.2), vh,
         RGBColor(0xF0, 0xF0, 0xF0), RGBColor(0xCC, 0xCC, 0xCC))
tf = add_tb(slide, LM + emu(0.1), vy + emu(0.5), col_w - emu(0.2), emu(0.45))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
ar(p, "🎬", size=22, color=C_LIGHT)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
ar(p2, "조치 전 동영상 삽입", size=9, color=C_LIGHT)

# 상세
dy = vy + vh + emu(0.08)
for i, (main, sub) in enumerate([
    ("프레임그레버 PCIe 슬롯에서 이탈", "확인"),
    ("카메라 이미지 취득 실패", "간헐적 통신 끊김"),
    ("장비 가동 중 검사 중단", "BSOD 또는 프리징"),
    ("고정 브라켓 나사 풀림", "미체결 상태"),
]):
    tf = add_tb(slide, LM + emu(0.1), dy + i * emu(0.24), col_w - emu(0.2), emu(0.22))
    p = tf.paragraphs[0]
    ar(p, "✕ ", size=8, bold=True, color=C_RED)
    ar(p, main, size=8.5, bold=True, color=C_DARK)
    ar(p, f" ({sub})", size=8, color=C_GRAY)

# -- After --
ax = LM + col_w + gap
add_rect(slide, ax, y_vs, col_w, col_h, RGBColor(0xF0, 0xFF, 0xF0), RGBColor(0xA5, 0xD6, 0xA7))

add_rect(slide, ax, y_vs, col_w, emu(0.3), C_GREEN)
tf = add_tb(slide, ax + emu(0.1), y_vs + emu(0.03), col_w - emu(0.2), emu(0.24))
p = tf.paragraphs[0]
ar(p, "✅  조치 후 (After)", size=12, bold=True, color=C_WHITE)

add_rect(slide, ax + emu(0.1), vy, col_w - emu(0.2), vh,
         RGBColor(0xF0, 0xF0, 0xF0), RGBColor(0xCC, 0xCC, 0xCC))
tf = add_tb(slide, ax + emu(0.1), vy + emu(0.5), col_w - emu(0.2), emu(0.45))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
ar(p, "🎬", size=22, color=C_LIGHT)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
ar(p2, "조치 후 동영상 삽입", size=9, color=C_LIGHT)

for i, (main, sub) in enumerate([
    ("프레임그레버 카드 재장착", "슬롯 완전 삽입 확인"),
    ("고정 브라켓 나사 체결", "풀림 방지 조치 적용"),
    ("카메라 이미지 취득", "정상 동작 확인"),
    ("연속 가동 테스트 완료", "이탈 재발 없음"),
]):
    tf = add_tb(slide, ax + emu(0.1), dy + i * emu(0.24), col_w - emu(0.2), emu(0.22))
    p = tf.paragraphs[0]
    ar(p, "✔ ", size=8, bold=True, color=C_GREEN)
    ar(p, main, size=8.5, bold=True, color=C_DARK)
    ar(p, f" ({sub})", size=8, color=C_GRAY)

# --- 조치 결과 ---
y_r = y_vs + col_h + emu(0.08)
add_rect(slide, LM, y_r, CW, emu(0.35), RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xA5, 0xD6, 0xA7))
tf = add_tb(slide, LM + emu(0.1), y_r + emu(0.05), CW - emu(0.2), emu(0.25))
p = tf.paragraphs[0]
ar(p, "✔ 조치 결과:  ", size=10, bold=True, color=C_GREEN)
ar(p, "프레임그레버 재장착 및 고정 강화 → 카메라 통신 정상 복구, 연속 가동 안정성 확인 완료", size=9.5, color=C_BODY)

# 저장
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"OK: {OUTPUT} ({os.path.getsize(OUTPUT)} bytes)")
