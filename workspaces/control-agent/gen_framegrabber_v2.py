# -*- coding: utf-8 -*-
"""
프레임그레버 이탈 이슈 PPT - 지건승 기본양식 기반
템플릿: 기본양식_Copied[00].pptx (13.333x7.5, 1_Title Only 레이아웃)
"""
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

TEMPLATE = r"C:\MES\wta-agents\workspaces\control-agent\template_base.pptx"
OUTPUT = r"C:\MES\wta-agents\reports\control-agent\framegrabber-issue-v2.pptx"

prs = Presentation(TEMPLATE)

# 기존 슬라이드 모두 삭제
from lxml import etree
prs_elem = prs.part._element
ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
sldIdLst = prs_elem.find(f'{ns}sldIdLst')
for sldId in list(sldIdLst):
    sldIdLst.remove(sldId)

# --- 슬라이드 1: 표지 (Title Slide 레이아웃) ---
title_layout = prs.slide_layouts[0]  # "Title Slide"
slide1 = prs.slides.add_slide(title_layout)

# 제목
title_ph = slide1.placeholders[0]
title_ph.text = ""
p = title_ph.text_frame.paragraphs[0]
run = p.add_run()
run.text = "검사기 프레임그레버"
run.font.bold = True
run.font.size = Pt(36)
p2 = title_ph.text_frame.add_paragraph()
run2 = p2.add_run()
run2.text = "(Frame Grabber) 이탈 이슈"
run2.font.bold = True
run2.font.size = Pt(36)

# 부제목
subtitle_ph = slide1.placeholders[1]
subtitle_ph.text = ""
p = subtitle_ph.text_frame.paragraphs[0]
run = p.add_run()
run.text = "작성자 : 지건승"
run.font.size = Pt(16)
p2 = subtitle_ph.text_frame.add_paragraph()
run2 = p2.add_run()
run2.text = "작성일 : 2026. 04."
run2.font.size = Pt(16)

# --- 슬라이드 2: 본문 (1_Title Only 레이아웃) ---
content_layout = prs.slide_layouts[11]  # "1_Title Only"
slide2 = prs.slides.add_slide(content_layout)

# ph[14] - 상단 카테고리 [제어설계]
slide2.placeholders[14].text = ""
p = slide2.placeholders[14].text_frame.paragraphs[0]
run = p.add_run()
run.text = "[제어설계]"
run.font.bold = True
run.font.size = Pt(12)

# ph[0] - 제목
slide2.placeholders[0].text = ""
p = slide2.placeholders[0].text_frame.paragraphs[0]
run = p.add_run()
run.text = "1. 검사기 프레임그레버(Frame Grabber) 이탈 이슈"
run.font.bold = True
run.font.size = Pt(22)

# ph[12] - 페이지 번호 (있으면 설정)
if 12 in slide2.placeholders:
    slide2.placeholders[12].text = ""
    p = slide2.placeholders[12].text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "2"
    run.font.size = Pt(10)

# --- 본문 내용 추가 (TextBox로) ---
def emu(inches):
    return int(inches * 914400)

def add_textbox(slide, left, top, width, height, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

# 상수
CONTENT_LEFT = emu(1.75)  # 레이아웃 좌측 여백에 맞춤
CONTENT_RIGHT = emu(11.6)
CONTENT_W = CONTENT_RIGHT - CONTENT_LEFT

# --- 정보 바 ---
y_info = emu(1.45)
info_bg = add_rounded_rect(slide2, CONTENT_LEFT, y_info, CONTENT_W, emu(0.35),
                            RGBColor(0xF5, 0xF5, 0xF5), RGBColor(0xE0, 0xE0, 0xE0))

tf = add_textbox(slide2, CONTENT_LEFT + emu(0.15), y_info + emu(0.03), CONTENT_W - emu(0.3), emu(0.3))
p = tf.paragraphs[0]
info_parts = [
    ("장비: ", True), ("HIM-F 검사기   ", False),
    ("부품: ", True), ("Frame Grabber (PCIe)   ", False),
    ("현상: ", True), ("PCIe 슬롯 이탈   ", False),
    ("담당: ", True), ("지건승   ", False),
    ("일자: ", True), ("2026. 04.", False),
]
for text, is_bold in info_parts:
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(10)
    run.font.bold = is_bold
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33) if is_bold else RGBColor(0x55, 0x55, 0x55)

# --- 원인 분석 박스 ---
y_cause = emu(1.95)
cause_bg = add_rounded_rect(slide2, CONTENT_LEFT, y_cause, CONTENT_W, emu(0.9),
                             RGBColor(0xFF, 0xF8, 0xE1), RGBColor(0xFF, 0xE0, 0x82))

# 원인 분석 제목
tf = add_textbox(slide2, CONTENT_LEFT + emu(0.15), y_cause + emu(0.05), emu(2), emu(0.25))
p = tf.paragraphs[0]
run = p.add_run()
run.text = "원인 분석"
run.font.name = "맑은 고딕"
run.font.size = Pt(12)
run.font.bold = True
run.font.color.rgb = RGBColor(0xE6, 0x51, 0x00)

# 원인 분석 내용
cause_items = [
    "장비 운송/가동 중 진동에 의한 PCIe 슬롯 접촉 불량 발생",
    "프레임그레버 카드 고정 브라켓 체결 미흡 또는 나사 풀림",
    "반복 진동으로 카드가 점진적으로 슬롯에서 이탈 → 카메라 통신 두절",
]
tf = add_textbox(slide2, CONTENT_LEFT + emu(0.15), y_cause + emu(0.3), CONTENT_W - emu(0.3), emu(0.55))
for i, item in enumerate(cause_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = f"• {item}"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p.space_after = Pt(1)

# --- Before / After 비교 영역 ---
y_vs = emu(3.0)
col_gap = emu(0.15)
col_w = (CONTENT_W - col_gap) // 2
col_h = emu(3.3)

# Before (좌측)
before_bg = add_rounded_rect(slide2, CONTENT_LEFT, y_vs, col_w, col_h,
                              RGBColor(0xFF, 0xF5, 0xF5), RGBColor(0xE0, 0xC0, 0xC0))

# Before 제목
tf = add_textbox(slide2, CONTENT_LEFT + emu(0.15), y_vs + emu(0.08), col_w - emu(0.3), emu(0.3))
p = tf.paragraphs[0]
run = p.add_run()
run.text = "⚠ 조치 전 (Before)"
run.font.name = "맑은 고딕"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = RGBColor(0xC6, 0x28, 0x28)

# Before 동영상 플레이스홀더
vid_y = y_vs + emu(0.45)
vid_h = emu(1.5)
vid_before = add_rounded_rect(slide2, CONTENT_LEFT + emu(0.15), vid_y,
                                col_w - emu(0.3), vid_h,
                                RGBColor(0xE8, 0xE8, 0xE8), RGBColor(0xCC, 0xCC, 0xCC))

tf = add_textbox(slide2, CONTENT_LEFT + emu(0.15), vid_y + emu(0.5), col_w - emu(0.3), emu(0.4))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "🎬 조치 전 동영상 삽입"
run.font.name = "맑은 고딕"
run.font.size = Pt(11)
run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

# Before 상세 내용
before_items = [
    "프레임그레버 PCIe 슬롯에서 이탈 확인",
    "카메라 이미지 취득 실패 / 간헐적 통신 끊김",
    "장비 가동 중 검사 중단 발생 (BSOD 또는 프리징)",
    "고정 브라켓 나사 풀림 / 미체결 상태",
]
desc_y = vid_y + vid_h + emu(0.1)
tf = add_textbox(slide2, CONTENT_LEFT + emu(0.15), desc_y, col_w - emu(0.3), emu(1.1))
for i, item in enumerate(before_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = f"• {item}"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p.space_after = Pt(2)

# After (우측)
after_x = CONTENT_LEFT + col_w + col_gap
after_bg = add_rounded_rect(slide2, after_x, y_vs, col_w, col_h,
                              RGBColor(0xF0, 0xFF, 0xF0), RGBColor(0xC0, 0xE0, 0xC0))

# After 제목
tf = add_textbox(slide2, after_x + emu(0.15), y_vs + emu(0.08), col_w - emu(0.3), emu(0.3))
p = tf.paragraphs[0]
run = p.add_run()
run.text = "✅ 조치 후 (After)"
run.font.name = "맑은 고딕"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

# After 동영상 플레이스홀더
vid_after = add_rounded_rect(slide2, after_x + emu(0.15), vid_y,
                              col_w - emu(0.3), vid_h,
                              RGBColor(0xE8, 0xE8, 0xE8), RGBColor(0xCC, 0xCC, 0xCC))

tf = add_textbox(slide2, after_x + emu(0.15), vid_y + emu(0.5), col_w - emu(0.3), emu(0.4))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "🎬 조치 후 동영상 삽입"
run.font.name = "맑은 고딕"
run.font.size = Pt(11)
run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

# After 상세 내용
after_items = [
    "프레임그레버 카드 재장착 및 슬롯 완전 삽입 확인",
    "고정 브라켓 나사 체결 + 풀림 방지 조치 적용",
    "카메라 이미지 취득 정상 동작 확인",
    "연속 가동 테스트 완료 (이탈 재발 없음)",
]
tf = add_textbox(slide2, after_x + emu(0.15), desc_y, col_w - emu(0.3), emu(1.1))
for i, item in enumerate(after_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = f"• {item}"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p.space_after = Pt(2)

# --- 조치 결과 바 ---
y_result = y_vs + col_h + emu(0.15)
result_bg = add_rounded_rect(slide2, CONTENT_LEFT, y_result, CONTENT_W, emu(0.35),
                              RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xA5, 0xD6, 0xA7))

tf = add_textbox(slide2, CONTENT_LEFT + emu(0.15), y_result + emu(0.03), CONTENT_W - emu(0.3), emu(0.3))
p = tf.paragraphs[0]
parts = [
    ("✔ 조치 결과: ", True, RGBColor(0x2E, 0x7D, 0x32)),
    ("프레임그레버 재장착 및 고정 강화 → 카메라 통신 정상 복구, 연속 가동 안정성 확인 완료", False, RGBColor(0x44, 0x44, 0x44)),
]
for text, is_bold, color in parts:
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(10)
    run.font.bold = is_bold
    run.font.color.rgb = color

# 저장
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"PPT saved: {OUTPUT}")
print(f"Size: {os.path.getsize(OUTPUT)} bytes")
