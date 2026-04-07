# -*- coding: utf-8 -*-
"""
핵심 부품 및 시스템 안정화 전략 PPT
기본양식의 슬라이드 레이아웃 유지, 2장 구성 (사진 1 + 사진 2)
"""
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os

TEMPLATE = r"C:\MES\wta-agents\workspaces\control-agent\template_base.pptx"
OUTPUT = r"C:\MES\wta-agents\reports\control-agent\parts-stability.pptx"

prs = Presentation(TEMPLATE)

# 기존 슬라이드 삭제
for _ in range(len(prs.slides)):
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# 헬퍼
def emu(inches):
    return int(inches * 914400)

def add_tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame

def add_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    return sh

def add_rounded(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    return sh

def ar(p, text, size=10, bold=False, color=None, font="맑은 고딕"):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return r

# 색상
C_DARK = RGBColor(0x22, 0x22, 0x22)
C_BODY = RGBColor(0x44, 0x44, 0x44)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_LIGHT = RGBColor(0x88, 0x88, 0x88)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_NAVY = RGBColor(0x1A, 0x23, 0x7E)
C_BLUE = RGBColor(0x42, 0x72, 0xC4)
C_LBLUE = RGBColor(0xD6, 0xE4, 0xF0)
C_RED = RGBColor(0xC6, 0x28, 0x28)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_ORANGE = RGBColor(0xE6, 0x51, 0x00)

LM = emu(1.75)
CW = emu(11.6) - LM

# =====================================================
# 슬라이드 1: 핵심 부품 및 시스템 안정화 전략 (사진 1 - 요약 대시보드)
# =====================================================
s1 = prs.slides.add_slide(prs.slide_layouts[11])  # 1_Title Only

# 카테고리
s1.placeholders[14].text = ""
p = s1.placeholders[14].text_frame.paragraphs[0]
ar(p, "[제어설계]", bold=True)

# 제목
s1.placeholders[0].text = ""
p = s1.placeholders[0].text_frame.paragraphs[0]
ar(p, "1. 핵심 부품 및 시스템 안정화 전략")

if 12 in s1.placeholders:
    s1.placeholders[12].text = ""
    p = s1.placeholders[12].text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    ar(p, "2")

# --- 섹션 제목: 1) 핵심 부품 안정화 ---
y = emu(1.25)

# 제목 바 (왼쪽 굵은 선 + 텍스트)
add_rect(s1, LM, y, emu(0.05), emu(0.28), C_NAVY)
tf = add_tb(s1, LM + emu(0.12), y, emu(4), emu(0.28))
p = tf.paragraphs[0]
ar(p, "1) 핵심 부품 안정화", size=13, bold=True, color=C_DARK)

# 오른쪽: 개발일정 상시
tf = add_tb(s1, LM + CW - emu(1.8), y, emu(1.8), emu(0.28))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
ar(p, "개발 일정: ", size=9, color=C_GRAY)
ar(p, "상시", size=9, bold=True, color=C_RED)

# --- ❶ 추진 방향 & 목표 ---
y1 = y + emu(0.4)
add_rounded(s1, LM, y1, CW * 0.48, emu(0.7), RGBColor(0xEE, 0xF2, 0xFF), RGBColor(0xB0, 0xC4, 0xDE))

# 추진 방향 제목
tf = add_tb(s1, LM + emu(0.1), y1 + emu(0.03), emu(2), emu(0.2))
p = tf.paragraphs[0]
ar(p, "❶  추진 방향 & 목표", size=10, bold=True, color=C_NAVY)

# 추진 방향 내용
tf = add_tb(s1, LM + emu(0.15), y1 + emu(0.24), CW * 0.48 - emu(0.25), emu(0.45))
p = tf.paragraphs[0]
ar(p, "선제적 장애 방지", size=9, bold=True, color=C_DARK)
p2 = tf.add_paragraph()
ar(p2, "핵심 부품 등은 표준화로 정하고, 신규/대체/교체는 사전 검증으로 Trouble 최소화", size=8, color=C_BODY)
p3 = tf.add_paragraph()
ar(p3, "기술 품질 내재화", size=9, bold=True, color=C_DARK)
p4 = tf.add_paragraph()
ar(p4, "현장 Trouble을 사전 차단하여 반복 이슈와 셋업 지연 최소화, 부품 변경 시 사전 테스트 후 적용", size=8, color=C_BODY)

# --- ❷ 핵심 요소 관리 매트릭스 ---
mx = LM + CW * 0.48 + emu(0.1)
mw = CW * 0.52 - emu(0.1)
add_rounded(s1, mx, y1, mw, emu(0.7), RGBColor(0xF5, 0xF5, 0xF5), RGBColor(0xD0, 0xD0, 0xD0))

tf = add_tb(s1, mx + emu(0.1), y1 + emu(0.03), emu(3), emu(0.2))
p = tf.paragraphs[0]
ar(p, "❷  핵심 요소 관리 매트릭스", size=10, bold=True, color=C_NAVY)

# 4개 박스 (2x2)
boxes = [
    ("SYSTEM & OS", "• PC 환경 표준화 및 이미지 관리\n• Win11 드라이버 안정 버전 고정"),
    ("INTERFACE", "• 통신 안정성 및 카드 관리 강화\n• PCIe, 케이블 규격 표준화"),
    ("NETWORK", "• 산업용 네트워크 안정성 확보\n• EtherCAT Drop 최소화"),
    ("VISION & TEACHING", "• 비전 시스템 구성 표준화\n• 노이즈 대책 수립"),
]

bw = (mw - emu(0.3)) / 2
bh = emu(0.2)
for i, (title, desc) in enumerate(boxes):
    bx = mx + emu(0.1) + (i % 2) * (bw + emu(0.1))
    by = y1 + emu(0.25) + (i // 2) * emu(0.23)
    tf = add_tb(s1, bx, by, bw, bh)
    p = tf.paragraphs[0]
    ar(p, f"{title}: ", size=8, bold=True, color=C_BLUE)
    ar(p, desc.split('\n')[0].replace('• ', ''), size=7.5, color=C_BODY)

# --- ❸ 리스크 대응 & 주요 개선 사례 ---
y2 = y1 + emu(0.8)

# 공급 및 단종 리스크 관리
rw = CW * 0.48
add_rounded(s1, LM, y2, rw, emu(0.5), RGBColor(0xFF, 0xF0, 0xE0), RGBColor(0xFF, 0xCC, 0x80))

tf = add_tb(s1, LM + emu(0.1), y2 + emu(0.03), rw - emu(0.2), emu(0.18))
p = tf.paragraphs[0]
ar(p, "⚠  공급 및 단종 리스크 관리", size=9, bold=True, color=C_ORANGE)

tf = add_tb(s1, LM + emu(0.1), y2 + emu(0.2), rw - emu(0.2), emu(0.28))
p = tf.paragraphs[0]
ar(p, "• 핵심 부품 대체 후보군 상시 확보", size=8, color=C_BODY)
p2 = tf.add_paragraph()
ar(p2, "• 공급 리스크 모니터링 및 선제 대응", size=8, color=C_BODY)

# 주요 개선 사례
sx = LM + rw + emu(0.1)
sw = CW - rw - emu(0.1)
add_rounded(s1, sx, y2, sw, emu(0.5), RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xA5, 0xD6, 0xA7))

tf = add_tb(s1, sx + emu(0.1), y2 + emu(0.03), sw - emu(0.2), emu(0.18))
p = tf.paragraphs[0]
ar(p, "✔  주요 개선 사례 (Spacer 규격/지그)", size=9, bold=True, color=C_GREEN)

tf = add_tb(s1, sx + emu(0.1), y2 + emu(0.2), sw - emu(0.2), emu(0.28))
p = tf.paragraphs[0]
ar(p, "• 종류 6가지→3가지 단순화 및 표준화", size=8, color=C_BODY)
p2 = tf.add_paragraph()
ar(p2, "• 보충은 전담 인원이 관리, 목록 갱신 능동적 관리 추진", size=8, color=C_BODY)

# 하단 안내문
tf = add_tb(s1, LM, y2 + emu(0.55), CW, emu(0.2))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
ar(p, "☑ 안정화 표준안이 완비된 신뢰성을 보장합니다. (목록 기능 업데이트 중..)", size=8, color=C_LIGHT)

# =====================================================
# 슬라이드 2: 핵심 부품 안정화 상세 테이블 (사진 2)
# =====================================================
s2 = prs.slides.add_slide(prs.slide_layouts[11])

s2.placeholders[14].text = ""
p = s2.placeholders[14].text_frame.paragraphs[0]
ar(p, "[제어설계]", bold=True)

s2.placeholders[0].text = ""
p = s2.placeholders[0].text_frame.paragraphs[0]
ar(p, "2. 핵심 부품 안정화")

if 12 in s2.placeholders:
    s2.placeholders[12].text = ""
    p = s2.placeholders[12].text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    ar(p, "3")

# --- 섹션 제목 ---
y = emu(1.25)
add_rect(s2, LM, y, emu(0.05), emu(0.28), C_NAVY)
tf = add_tb(s2, LM + emu(0.12), y, emu(4), emu(0.28))
p = tf.paragraphs[0]
ar(p, "1) 핵심 부품 안정화", size=13, bold=True, color=C_DARK)

tf = add_tb(s2, LM + CW - emu(1.8), y, emu(1.8), emu(0.28))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
ar(p, "개발 일정: ", size=9, color=C_GRAY)
ar(p, "상시", size=9, bold=True, color=C_RED)

# --- 1. 추진 방향 ---
y_s = y + emu(0.38)
tf = add_tb(s2, LM, y_s, CW, emu(0.22))
p = tf.paragraphs[0]
ar(p, "1.  추진 방향", size=10, bold=True, color=C_DARK)

tf = add_tb(s2, LM + emu(0.25), y_s + emu(0.2), CW - emu(0.3), emu(0.18))
p = tf.paragraphs[0]
ar(p, "▪  핵심 부품 등은 표준화로 정하고, 신규/대체/교체는 사전 검증으로 Trouble 최소화", size=9, color=C_BODY)

# --- 2. 목표 ---
y_m = y_s + emu(0.42)
tf = add_tb(s2, LM, y_m, CW, emu(0.22))
p = tf.paragraphs[0]
ar(p, "2.  목표", size=10, bold=True, color=C_DARK)

tf = add_tb(s2, LM + emu(0.25), y_m + emu(0.2), CW - emu(0.3), emu(0.18))
p = tf.paragraphs[0]
ar(p, "▪  현장 Trouble을 사전에 차단하여 반복적인 이슈와 셋업 지연을 최소화", size=9, color=C_BODY)

tf = add_tb(s2, LM + emu(0.25), y_m + emu(0.36), CW - emu(0.3), emu(0.18))
p = tf.paragraphs[0]
ar(p, "▪  부품 변경 시 그에 맞는 제품을 사전 선정하여 테스트 후 적용하는 절차 확립", size=9, color=C_BODY)

# --- 3. 진행 사항 (테이블) ---
y_t = y_m + emu(0.6)
tf = add_tb(s2, LM, y_t, CW, emu(0.22))
p = tf.paragraphs[0]
ar(p, "3.  진행 사항", size=10, bold=True, color=C_DARK)

# 테이블 데이터
headers = ["안정화 요소", "관리 자재 (부품/S/W)", "대응 방안 (Trouble 대응)", "추진 방향 (표준화)"]
rows = [
    ["PC 환경 불일치", "PC, Win11, WMX3, 드라이버", "표준 이미지 복제 및 변경 시 검증 시행", "환경 표준화 및 사양 고정"],
    ["블루스크린/프리징", "OS 정책, 칩셋/GPU/LAN 드라이버", "발생 원인별 이력 관리 및 추적", "누적 데이터를 통한 예방 관리"],
    ["PCIe 카드 에러", "PCIe 카드, 슬롯, BIOS, 프로그램 셋팅", "점검 순서 고정 및 호환성 리스트 관리", "부품 간 호환성 검증 체계 구축"],
    ["Ethernet 불안정", "통신 케이블", "케이블 규격 고정 및 정기 점검 루틴화", "물리적 인터페이스 사양 고정"],
    ["EtherCAT Drop", "Master, Slave, ESI, 케이블", "사전 검증 및 장시간 부하 테스트 실시", "장비 출고 전 사전 검증 의무화"],
    ["USB 끊김", "USB 포트/허브, 케이블", "권장 포트 지정 및 표준 허브 사용", "연결 방식 및 자재 표준화"],
    ["Serial 장애", "Serial 카드/변환기, 통신 설정", "파라미터 표준화 및 교차 테스트", "장비별 통신 설정 표준화"],
    ["하네스 노이즈/단선", "Power/IO/ENC 하네스, 실드/접지", "규격·실드 고정 및 임의 현장 수정 제한", "하네스 제작 및 배선 기준 고정"],
    ["Vision 오동작", "카메라, 조명, 트리거", "전원/통신 분리 배선 및 노이즈 대책 수립", "Vision 시스템 구성 표준화"],
    ["외부 장비 연동", "TP, 저울, 라벨 프린터, 레이저 마킹기 등", "교체 시 사전 검토 및 사용 매뉴얼화", "임의 사양 변경 통제 및 관리"],
    ["공급/단종 리스크", "핵심 부품 전반", "리스트 관리 및 대체 후보군 상시 확보", "부품 단종 대비 리스크 관리 체계"],
]

# 테이블 위치/크기
tx = LM
ty = y_t + emu(0.25)
tw = CW
col_widths = [emu(1.4), emu(2.7), emu(3.0), emu(2.75)]
row_h = emu(0.22)
header_h = emu(0.25)

# 헤더 행
for j, (header, cw_val) in enumerate(zip(headers, col_widths)):
    cx = tx + sum(col_widths[:j])
    add_rect(s2, cx, ty, cw_val, header_h, C_NAVY)
    tf = add_tb(s2, cx + emu(0.05), ty + emu(0.02), cw_val - emu(0.1), header_h - emu(0.04))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    ar(p, header, size=8, bold=True, color=C_WHITE)

# 데이터 행
for i, row in enumerate(rows):
    ry = ty + header_h + i * row_h
    bg = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xF8, 0xF8, 0xF8)
    line_color = RGBColor(0xE0, 0xC0, 0x80)

    for j, (cell, cw_val) in enumerate(zip(row, col_widths)):
        cx = tx + sum(col_widths[:j])
        add_rect(s2, cx, ry, cw_val, row_h, bg, line_color)
        tf = add_tb(s2, cx + emu(0.05), ry + emu(0.01), cw_val - emu(0.1), row_h - emu(0.02))
        p = tf.paragraphs[0]
        if j == 0:
            ar(p, cell, size=7.5, bold=True, color=C_DARK)
        else:
            ar(p, cell, size=7.5, color=C_BODY)

# 저장
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"OK: {OUTPUT} ({os.path.getsize(OUTPUT)} bytes)")
