# -*- coding: utf-8 -*-
"""
프레임그레버 이탈 이슈 PPT v4 - 기본양식에 맞춰 정돈된 디자인
표지(Title Slide) + 이슈 요약(1_Title Only) + Before/After(1_Title Only)
"""
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

TEMPLATE = r"C:\MES\wta-agents\workspaces\control-agent\template_base.pptx"
OUTPUT = r"C:\MES\wta-agents\reports\control-agent\framegrabber-issue-v2.pptx"

prs = Presentation(TEMPLATE)

# 기존 슬라이드 삭제 (안전한 방법)
def delete_slide(prs, index):
    rId = prs.slides._sldIdLst[index].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]

for _ in range(len(prs.slides)):
    delete_slide(prs, 0)

# === 헬퍼 ===
def emu(inches):
    return int(inches * 914400)

def add_tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame

def add_rect(slide, l, t, w, h, fill, line=None, radius=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh

def add_run(p, text, size=10, bold=False, color=None, font="맑은 고딕"):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return r

# 색상
C_RED = RGBColor(0xC6, 0x28, 0x28)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_ORANGE = RGBColor(0xE6, 0x51, 0x00)
C_DARK = RGBColor(0x22, 0x22, 0x22)
C_BODY = RGBColor(0x44, 0x44, 0x44)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_LIGHT = RGBColor(0x88, 0x88, 0x88)

# 레이아웃 영역 (양식 기준)
LM = emu(1.75)   # 왼쪽 여백 (레이아웃 사이드바 뒤)
RM = emu(11.6)   # 오른쪽 끝
CW = RM - LM     # 컨텐츠 폭

# =====================================================
# 슬라이드 1: 표지
# =====================================================
s1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide

# 제목
ph_title = s1.placeholders[0]
ph_title.text = ""
tf = ph_title.text_frame
p = tf.paragraphs[0]
add_run(p, "검사기 프레임그레버", size=36, bold=True)
p2 = tf.add_paragraph()
add_run(p2, "(Frame Grabber) 이탈 이슈", size=36, bold=True)

# 부제목
ph_sub = s1.placeholders[1]
ph_sub.text = ""
tf = ph_sub.text_frame
p = tf.paragraphs[0]
add_run(p, "장비: HIM-F 검사기  |  부품: Frame Grabber (PCIe)", size=14)
p2 = tf.add_paragraph()
add_run(p2, "", size=6)
p3 = tf.add_paragraph()
add_run(p3, "작성자 : 지건승", size=14)
p4 = tf.add_paragraph()
add_run(p4, "작성일 : 2026. 04.", size=14)

# =====================================================
# 슬라이드 2: 이슈 개요 + 원인 분석
# =====================================================
s2 = prs.slides.add_slide(prs.slide_layouts[11])  # 1_Title Only

# 카테고리
s2.placeholders[14].text = ""
p = s2.placeholders[14].text_frame.paragraphs[0]
add_run(p, "[제어설계]", bold=True)

# 제목
s2.placeholders[0].text = ""
p = s2.placeholders[0].text_frame.paragraphs[0]
add_run(p, "1. 이슈 개요 및 원인 분석", bold=True)

# 페이지 번호
if 12 in s2.placeholders:
    s2.placeholders[12].text = ""
    p = s2.placeholders[12].text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    add_run(p, "2")

# --- 이슈 개요 테이블 형태 ---
y = emu(1.35)

# 개요 영역 배경
add_rect(s2, LM, y, CW, emu(1.6), RGBColor(0xFA, 0xFA, 0xFA), RGBColor(0xE0, 0xE0, 0xE0))

# 개요 제목
tf = add_tb(s2, LM + emu(0.15), y + emu(0.08), emu(2), emu(0.25))
p = tf.paragraphs[0]
add_run(p, "이슈 개요", size=13, bold=True, color=C_DARK)

# 개요 항목들 - 테이블 스타일
items = [
    ("장 비", "HIM-F 검사기"),
    ("부 품", "Frame Grabber (PCIe 카드)"),
    ("현 상", "PCIe 슬롯에서 카드 이탈 → 카메라 통신 두절 → 검사 중단"),
    ("담 당", "지건승"),
    ("발생일", "2026년 4월"),
]

label_x = LM + emu(0.15)
value_x = LM + emu(1.35)
item_h = emu(0.22)

for i, (label, value) in enumerate(items):
    iy = y + emu(0.4) + i * item_h

    # 라벨 배경
    add_rect(s2, label_x, iy, emu(1.05), item_h - emu(0.02),
             RGBColor(0xE8, 0xE8, 0xE8))
    tf = add_tb(s2, label_x + emu(0.05), iy, emu(0.95), item_h)
    p = tf.paragraphs[0]
    add_run(p, label, size=9, bold=True, color=C_DARK)

    # 값
    tf = add_tb(s2, value_x + emu(0.1), iy, CW - emu(1.45), item_h)
    p = tf.paragraphs[0]
    add_run(p, value, size=9, color=C_BODY)

# --- 원인 분석 ---
y2 = y + emu(1.75)
add_rect(s2, LM, y2, CW, emu(1.55), RGBColor(0xFF, 0xF8, 0xE1), RGBColor(0xFF, 0xE0, 0x82))

tf = add_tb(s2, LM + emu(0.15), y2 + emu(0.08), emu(3), emu(0.25))
p = tf.paragraphs[0]
add_run(p, "⚠  원인 분석", size=13, bold=True, color=C_ORANGE)

causes = [
    ("직접 원인", "장비 운송 및 가동 중 발생하는 진동에 의해 PCIe 슬롯 접촉 불량 발생"),
    ("간접 원인", "프레임그레버 카드 고정 브라켓 체결 미흡 또는 나사 풀림"),
    ("진행 경과", "반복 진동 → 카드 점진적 이탈 → 카메라 통신 두절 → 검사 중단 (BSOD/프리징)"),
]

for i, (title, desc) in enumerate(causes):
    cy = y2 + emu(0.4) + i * emu(0.35)

    # 번호 원
    add_rect(s2, LM + emu(0.2), cy + emu(0.02), emu(0.22), emu(0.22),
             RGBColor(0xE6, 0x51, 0x00))
    tf = add_tb(s2, LM + emu(0.2), cy + emu(0.02), emu(0.22), emu(0.22))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, str(i + 1), size=9, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # 제목 + 설명
    tf = add_tb(s2, LM + emu(0.5), cy, CW - emu(0.6), emu(0.3))
    p = tf.paragraphs[0]
    add_run(p, f"{title}:  ", size=10, bold=True, color=C_DARK)
    add_run(p, desc, size=10, color=C_BODY)

# --- 영향 분석 ---
y3 = y2 + emu(1.7)
add_rect(s2, LM, y3, CW, emu(0.9), RGBColor(0xFF, 0xEB, 0xEE), RGBColor(0xEF, 0x9A, 0x9A))

tf = add_tb(s2, LM + emu(0.15), y3 + emu(0.08), emu(3), emu(0.25))
p = tf.paragraphs[0]
add_run(p, "📋  영향 범위", size=13, bold=True, color=C_RED)

impacts = [
    "검사 라인 가동 중단 → 생산 지연 발생",
    "카메라 통신 두절로 검사 데이터 유실 위험",
    "동일 장비군(HIM-F) 전체에 동일 이슈 잠재",
]

for i, text in enumerate(impacts):
    tf = add_tb(s2, LM + emu(0.35), y3 + emu(0.38) + i * emu(0.18), CW - emu(0.5), emu(0.18))
    p = tf.paragraphs[0]
    add_run(p, f"•  {text}", size=9.5, color=RGBColor(0x55, 0x33, 0x33))

# =====================================================
# 슬라이드 3: Before / After 비교
# =====================================================
s3 = prs.slides.add_slide(prs.slide_layouts[11])  # 1_Title Only

s3.placeholders[14].text = ""
p = s3.placeholders[14].text_frame.paragraphs[0]
add_run(p, "[제어설계]", bold=True)

s3.placeholders[0].text = ""
p = s3.placeholders[0].text_frame.paragraphs[0]
add_run(p, "2. 조치 내용 및 결과 (Before / After)", bold=True)

if 12 in s3.placeholders:
    s3.placeholders[12].text = ""
    p = s3.placeholders[12].text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    add_run(p, "3")

# Before / After 영역
y_vs = emu(1.35)
gap = emu(0.15)
col_w = (CW - gap) // 2
col_h = emu(3.9)

# ---- Before (좌측) ----
add_rect(s3, LM, y_vs, col_w, col_h, RGBColor(0xFF, 0xF5, 0xF5), RGBColor(0xEF, 0x9A, 0x9A))

# Before 제목 바
add_rect(s3, LM, y_vs, col_w, emu(0.38), C_RED)
tf = add_tb(s3, LM + emu(0.15), y_vs + emu(0.05), col_w - emu(0.3), emu(0.28))
p = tf.paragraphs[0]
add_run(p, "⚠  조치 전 (Before)", size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# 동영상 플레이스홀더
vy = y_vs + emu(0.5)
vh = emu(1.7)
vid = add_rect(s3, LM + emu(0.15), vy, col_w - emu(0.3), vh,
               RGBColor(0xF0, 0xF0, 0xF0), RGBColor(0xCC, 0xCC, 0xCC))
tf = add_tb(s3, LM + emu(0.15), vy + emu(0.6), col_w - emu(0.3), emu(0.4))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "🎬", size=24, color=C_LIGHT)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
add_run(p2, "조치 전 동영상 삽입", size=10, color=C_LIGHT)

# Before 상세
before_items = [
    ("프레임그레버 PCIe 슬롯에서 이탈", "확인"),
    ("카메라 이미지 취득 실패", "간헐적 통신 끊김"),
    ("장비 가동 중 검사 중단 발생", "BSOD 또는 프리징"),
    ("고정 브라켓 나사 풀림", "미체결 상태"),
]
dy = vy + vh + emu(0.12)
for i, (main, sub) in enumerate(before_items):
    tf = add_tb(s3, LM + emu(0.15), dy + i * emu(0.28), col_w - emu(0.3), emu(0.28))
    p = tf.paragraphs[0]
    add_run(p, "✕  ", size=9, bold=True, color=C_RED)
    add_run(p, main, size=9.5, bold=True, color=C_DARK)
    add_run(p, f"  ({sub})", size=9, color=C_GRAY)

# ---- After (우측) ----
ax = LM + col_w + gap
add_rect(s3, ax, y_vs, col_w, col_h, RGBColor(0xF0, 0xFF, 0xF0), RGBColor(0xA5, 0xD6, 0xA7))

# After 제목 바
add_rect(s3, ax, y_vs, col_w, emu(0.38), C_GREEN)
tf = add_tb(s3, ax + emu(0.15), y_vs + emu(0.05), col_w - emu(0.3), emu(0.28))
p = tf.paragraphs[0]
add_run(p, "✅  조치 후 (After)", size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# 동영상 플레이스홀더
vid2 = add_rect(s3, ax + emu(0.15), vy, col_w - emu(0.3), vh,
                RGBColor(0xF0, 0xF0, 0xF0), RGBColor(0xCC, 0xCC, 0xCC))
tf = add_tb(s3, ax + emu(0.15), vy + emu(0.6), col_w - emu(0.3), emu(0.4))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "🎬", size=24, color=C_LIGHT)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
add_run(p2, "조치 후 동영상 삽입", size=10, color=C_LIGHT)

# After 상세
after_items = [
    ("프레임그레버 카드 재장착", "슬롯 완전 삽입 확인"),
    ("고정 브라켓 나사 체결", "풀림 방지 조치 적용"),
    ("카메라 이미지 취득", "정상 동작 확인"),
    ("연속 가동 테스트 완료", "이탈 재발 없음"),
]
for i, (main, sub) in enumerate(after_items):
    tf = add_tb(s3, ax + emu(0.15), dy + i * emu(0.28), col_w - emu(0.3), emu(0.28))
    p = tf.paragraphs[0]
    add_run(p, "✔  ", size=9, bold=True, color=C_GREEN)
    add_run(p, main, size=9.5, bold=True, color=C_DARK)
    add_run(p, f"  ({sub})", size=9, color=C_GRAY)

# --- 조치 결과 요약 바 ---
y_r = y_vs + col_h + emu(0.12)
add_rect(s3, LM, y_r, CW, emu(0.55), RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xA5, 0xD6, 0xA7))

tf = add_tb(s3, LM + emu(0.15), y_r + emu(0.05), CW - emu(0.3), emu(0.2))
p = tf.paragraphs[0]
add_run(p, "✔ 조치 결과", size=12, bold=True, color=C_GREEN)

tf = add_tb(s3, LM + emu(0.15), y_r + emu(0.25), CW - emu(0.3), emu(0.25))
p = tf.paragraphs[0]
add_run(p, "프레임그레버 재장착 및 고정 브라켓 강화 완료 → ", size=10, color=C_BODY)
add_run(p, "카메라 통신 정상 복구", size=10, bold=True, color=C_GREEN)
add_run(p, ", 연속 가동 안정성 확인 완료. ", size=10, color=C_BODY)
add_run(p, "동일 장비군 전수 점검 및 예방 조치 시행.", size=10, bold=True, color=C_DARK)

# 저장
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"OK: {OUTPUT} ({os.path.getsize(OUTPUT)} bytes)")
