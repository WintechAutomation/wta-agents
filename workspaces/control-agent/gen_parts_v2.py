# -*- coding: utf-8 -*-
"""
핵심 부품 및 시스템 안정화 전략 PPT v2
사진 레이아웃 충실 재현 + 기본양식 사용
슬라이드2는 실제 Table 객체 사용
"""
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

TEMPLATE = r"C:\MES\wta-agents\workspaces\control-agent\template_base.pptx"
OUTPUT = r"C:\MES\wta-agents\reports\control-agent\parts-stability.pptx"

prs = Presentation(TEMPLATE)

for _ in range(len(prs.slides)):
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

def emu(inches):
    return int(inches * 914400)

def add_tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame

def add_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    return sh

def add_rr(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    return sh

def ar(p, text, sz=10, b=False, c=None, fn="맑은 고딕"):
    r = p.add_run()
    r.text = text
    r.font.name = fn
    r.font.size = Pt(sz)
    r.font.bold = b
    if c: r.font.color.rgb = c
    return r

# Colors
NAVY = RGBColor(0x1A, 0x23, 0x7E)
DBLUE = RGBColor(0x2C, 0x3E, 0x6B)
BLUE = RGBColor(0x44, 0x72, 0xC4)
DARK = RGBColor(0x22, 0x22, 0x22)
BODY = RGBColor(0x44, 0x44, 0x44)
GRAY = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC6, 0x28, 0x28)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
BG_DARK = RGBColor(0x2D, 0x2D, 0x2D)
BG_CARD = RGBColor(0x3A, 0x3A, 0x3A)

LM = emu(1.75)
CW = emu(9.85)

# =====================================================
# 슬라이드 1: 요약 대시보드 (경영 1.png)
# =====================================================
s1 = prs.slides.add_slide(prs.slide_layouts[11])

s1.placeholders[14].text = ""
p = s1.placeholders[14].text_frame.paragraphs[0]
ar(p, "[제어설계]", b=True)

s1.placeholders[0].text = ""
p = s1.placeholders[0].text_frame.paragraphs[0]
ar(p, "1. 핵심 부품 및 시스템 안정화 전략")

if 12 in s1.placeholders:
    s1.placeholders[12].text = ""
    p = s1.placeholders[12].text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    ar(p, "2")

# 전체 컨텐츠 영역에 어두운 배경
y0 = emu(1.2)
add_rr(s1, LM, y0, CW, emu(5.5), BG_DARK)

# 섹션 제목 바
add_rect(s1, LM + emu(0.15), y0 + emu(0.1), CW - emu(0.3), emu(0.3), RGBColor(0x45, 0x45, 0x45))
tf = add_tb(s1, LM + emu(0.25), y0 + emu(0.12), emu(5), emu(0.26))
p = tf.paragraphs[0]
ar(p, "1) 핵심 부품 및 시스템 안정화 전략", sz=12, b=True, c=WHITE)
tf = add_tb(s1, LM + CW - emu(2.2), y0 + emu(0.12), emu(2), emu(0.26))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
ar(p, "개발일정: ", sz=9, c=LGRAY)
ar(p, "상시", sz=9, b=True, c=RGBColor(0xFF, 0x80, 0x80))

# --- ❶ 추진 방향 & 목표 ---
y1 = y0 + emu(0.5)
tf = add_tb(s1, LM + emu(0.2), y1, emu(3), emu(0.22))
p = tf.paragraphs[0]
ar(p, "❶  추진 방향 & 목표", sz=10, b=True, c=RGBColor(0xBB, 0xCC, 0xFF))

# 좌: 선제적 장애 방지
bw = (CW - emu(0.5)) / 2
y1c = y1 + emu(0.25)
add_rr(s1, LM + emu(0.2), y1c, bw, emu(0.75), BG_CARD, RGBColor(0x55, 0x55, 0x55))

# 헤더 바
add_rect(s1, LM + emu(0.2), y1c, bw, emu(0.22), DBLUE)
tf = add_tb(s1, LM + emu(0.3), y1c + emu(0.02), bw - emu(0.2), emu(0.18))
p = tf.paragraphs[0]
ar(p, "선제적 장애 방지", sz=9, b=True, c=WHITE)

tf = add_tb(s1, LM + emu(0.3), y1c + emu(0.25), bw - emu(0.2), emu(0.45))
p = tf.paragraphs[0]
ar(p, "핵심 부품 등은 표준화로 정하고, 신규/대체/교체는 사전 검증으로 Trouble 최소화", sz=8, c=RGBColor(0xCC, 0xCC, 0xCC))
p2 = tf.add_paragraph()
ar(p2, "목표", sz=8, b=True, c=RGBColor(0xDD, 0xDD, 0xDD))
p3 = tf.add_paragraph()
ar(p3, "현장 Trouble을 사전에 차단하여 반복 이슈와 셋업 지연 최소화", sz=8, c=RGBColor(0xCC, 0xCC, 0xCC))

# 우: 기술 품질 내재화
bx2 = LM + emu(0.2) + bw + emu(0.1)
add_rr(s1, bx2, y1c, bw, emu(0.75), BG_CARD, RGBColor(0x55, 0x55, 0x55))
add_rect(s1, bx2, y1c, bw, emu(0.22), DBLUE)
tf = add_tb(s1, bx2 + emu(0.1), y1c + emu(0.02), bw - emu(0.2), emu(0.18))
p = tf.paragraphs[0]
ar(p, "기술 품질 내재화", sz=9, b=True, c=WHITE)

tf = add_tb(s1, bx2 + emu(0.1), y1c + emu(0.25), bw - emu(0.2), emu(0.45))
p = tf.paragraphs[0]
ar(p, "부품 변경 시 그에 맞는 제품을 사전 선정하여 테스트 후 적용하는 절차 확립", sz=8, c=RGBColor(0xCC, 0xCC, 0xCC))
p2 = tf.add_paragraph()
ar(p2, "목표", sz=8, b=True, c=RGBColor(0xDD, 0xDD, 0xDD))
p3 = tf.add_paragraph()
ar(p3, "부품 변경 시 사전 선정 → 테스트 → 적용 절차 표준화", sz=8, c=RGBColor(0xCC, 0xCC, 0xCC))

# --- ❷ 핵심 요소 관리 매트릭스 ---
y2 = y1c + emu(0.85)
tf = add_tb(s1, LM + emu(0.2), y2, emu(3), emu(0.22))
p = tf.paragraphs[0]
ar(p, "❷  핵심 요소 관리 매트릭스", sz=10, b=True, c=RGBColor(0xBB, 0xCC, 0xFF))

matrix = [
    ("SYSTEM & OS", ["• PC 환경 표준화 및 이미지 관리", "• Win11 드라이버 안정 버전(Realtek Image 사용)", "• OS 정책/업데이트 관리 표준화"]),
    ("INTERFACE", ["• 통신 안정성 및 카드 관리 강화", "• PCIe 카드 체결 및 호환성 관리", "• 케이블 규격 표준화 및 정기 점검"]),
    ("NETWORK", ["• 산업용 네트워크 안정성 확보", "• EtherCAT Drop 원인 분석 및 대책", "• 장비 출고 전 사전 검증 의무화"]),
    ("VISION & TEACHING", ["• 비전 시스템 구성 표준화", "• 전원/통신 분리 배선 및 노이즈 대책", "• 카메라/조명/트리거 설정 관리"]),
]

y2c = y2 + emu(0.25)
mw = (CW - emu(0.5)) / 2
mh = emu(0.7)

for i, (title, items) in enumerate(matrix):
    mx = LM + emu(0.2) + (i % 2) * (mw + emu(0.1))
    my = y2c + (i // 2) * (mh + emu(0.08))

    add_rr(s1, mx, my, mw, mh, BG_CARD, RGBColor(0x55, 0x55, 0x55))

    # 제목
    tf = add_tb(s1, mx + emu(0.1), my + emu(0.04), mw - emu(0.2), emu(0.18))
    p = tf.paragraphs[0]
    ar(p, title, sz=9, b=True, c=BLUE)

    # 내용
    tf = add_tb(s1, mx + emu(0.1), my + emu(0.22), mw - emu(0.2), mh - emu(0.26))
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        ar(p, item, sz=7, c=RGBColor(0xBB, 0xBB, 0xBB))
        p.space_after = Pt(0)

# --- ❸ 리스크 대응 & 주요 개선 사례 ---
y3 = y2c + 2 * (mh + emu(0.08)) + emu(0.05)
tf = add_tb(s1, LM + emu(0.2), y3, emu(4), emu(0.22))
p = tf.paragraphs[0]
ar(p, "❸  리스크 대응 및 주요 개선 사례", sz=10, b=True, c=RGBColor(0xBB, 0xCC, 0xFF))

y3c = y3 + emu(0.25)
rw = (CW - emu(0.5)) / 2

# 좌: 공급 및 단종 리스크 관리
add_rr(s1, LM + emu(0.2), y3c, rw, emu(0.65), RGBColor(0x4A, 0x3A, 0x2A), RGBColor(0x80, 0x60, 0x30))

tf = add_tb(s1, LM + emu(0.3), y3c + emu(0.04), rw - emu(0.2), emu(0.18))
p = tf.paragraphs[0]
ar(p, "⚠  공급 및 단종 리스크 관리", sz=9, b=True, c=RGBColor(0xFF, 0xCC, 0x80))

tf = add_tb(s1, LM + emu(0.3), y3c + emu(0.22), rw - emu(0.2), emu(0.4))
p = tf.paragraphs[0]
ar(p, "• 핵심 부품에 대한 대체 후보군 및 호환 테스트 결과를 관리", sz=7.5, c=RGBColor(0xCC, 0xCC, 0xCC))
p2 = tf.add_paragraph()
ar(p2, "• EOL(End of Channel Support) 선제 대응 체계 구축", sz=7.5, c=RGBColor(0xCC, 0xCC, 0xCC))
p3 = tf.add_paragraph()
ar(p3, "• 공급 리스크 모니터링 및 대체품 사전 확보", sz=7.5, c=RGBColor(0xCC, 0xCC, 0xCC))

# 우: 주요 개선 사례
rx2 = LM + emu(0.2) + rw + emu(0.1)
add_rr(s1, rx2, y3c, rw, emu(0.65), RGBColor(0x2A, 0x3A, 0x2A), RGBColor(0x40, 0x70, 0x40))

tf = add_tb(s1, rx2 + emu(0.1), y3c + emu(0.04), rw - emu(0.2), emu(0.18))
p = tf.paragraphs[0]
ar(p, "✔  주요 개선 사례 (Spacer 규격/지그)", sz=9, b=True, c=RGBColor(0x80, 0xFF, 0x80))

tf = add_tb(s1, rx2 + emu(0.1), y3c + emu(0.22), rw - emu(0.2), emu(0.4))
p = tf.paragraphs[0]
ar(p, "• 종류 6가지 → 3가지로 단순화 및 표준화 완료", sz=7.5, c=RGBColor(0xCC, 0xCC, 0xCC))
p2 = tf.add_paragraph()
ar(p2, "• 보충은 전담 인원이 관리, 목록 갱신 능동적 관리 추진", sz=7.5, c=RGBColor(0xCC, 0xCC, 0xCC))
p3 = tf.add_paragraph()
ar(p3, "• 사전 검증 프로세스 도입 후 현장 이슈 감소", sz=7.5, c=RGBColor(0xCC, 0xCC, 0xCC))

# 하단 안내
tf = add_tb(s1, LM + emu(0.2), y3c + emu(0.72), CW - emu(0.4), emu(0.18))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
ar(p, "☑ 안정화 표준안이 완비된 신뢰성을 보장합니다. (목록 기능 업데이트 중..)", sz=7.5, c=RGBColor(0x88, 0x88, 0x88))

# =====================================================
# 슬라이드 2: 상세 테이블 (2.png) - 실제 Table 사용
# =====================================================
s2 = prs.slides.add_slide(prs.slide_layouts[11])

s2.placeholders[14].text = ""
p = s2.placeholders[14].text_frame.paragraphs[0]
ar(p, "[제어설계]", b=True)

s2.placeholders[0].text = ""
p = s2.placeholders[0].text_frame.paragraphs[0]
ar(p, "2. 핵심 부품 안정화")

if 12 in s2.placeholders:
    s2.placeholders[12].text = ""
    p = s2.placeholders[12].text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    ar(p, "3")

# 섹션 제목
y = emu(1.2)
add_rect(s2, LM, y, CW, emu(0.28), RGBColor(0xF0, 0xF0, 0xF0), RGBColor(0xCC, 0xCC, 0xCC))
add_rect(s2, LM, y, emu(0.04), emu(0.28), NAVY)
tf = add_tb(s2, LM + emu(0.12), y + emu(0.01), emu(4), emu(0.26))
p = tf.paragraphs[0]
ar(p, "1) 핵심 부품 안정화", sz=12, b=True, c=DARK)
tf = add_tb(s2, LM + CW - emu(1.8), y + emu(0.01), emu(1.7), emu(0.26))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
ar(p, "개발 일정: ", sz=9, c=GRAY)
ar(p, "상시", sz=9, b=True, c=RED)

# 1. 추진 방향
ys = y + emu(0.38)
tf = add_tb(s2, LM, ys, CW, emu(0.2))
p = tf.paragraphs[0]
ar(p, "1.  추진 방향", sz=10, b=True, c=DARK)

tf = add_tb(s2, LM + emu(0.3), ys + emu(0.2), CW - emu(0.35), emu(0.16))
p = tf.paragraphs[0]
ar(p, "▪  핵심 부품 등은 표준화로 정하고, 신규/대체/교체는 사전 검증으로 Trouble 최소화", sz=9, c=BODY)

# 2. 목표
ym = ys + emu(0.42)
tf = add_tb(s2, LM, ym, CW, emu(0.2))
p = tf.paragraphs[0]
ar(p, "2.  목표", sz=10, b=True, c=DARK)

tf = add_tb(s2, LM + emu(0.3), ym + emu(0.2), CW - emu(0.35), emu(0.16))
p = tf.paragraphs[0]
ar(p, "▪  현장 Trouble을 사전에 차단하여 반복적인 이슈와 셋업 지연을 최소화", sz=9, c=BODY)

tf = add_tb(s2, LM + emu(0.3), ym + emu(0.36), CW - emu(0.35), emu(0.16))
p = tf.paragraphs[0]
ar(p, "▪  부품 변경 시 그에 맞는 제품을 사전 선정하여 테스트 후 적용하는 절차 확립", sz=9, c=BODY)

# 3. 진행 사항
yt = ym + emu(0.58)
tf = add_tb(s2, LM, yt, CW, emu(0.2))
p = tf.paragraphs[0]
ar(p, "3.  진행 사항", sz=10, b=True, c=DARK)

# 실제 Table
headers = ["안정화 요소", "관리 자재 (부품/S/W)", "대응 방안 (Trouble 대응)", "추진 방향 (표준화)"]
rows_data = [
    ["PC 환경 불일치", "PC, Win11, WMX3, 드라이버", "표준 이미지 복제 및 변경 시 검증 시행", "환경 표준화 및 사양 고정"],
    ["블루스크린/프리징", "OS 정책, 칩셋/GPU/LAN 드라이버", "발생 원인별 이력 관리 및 추적", "누적 데이터를 통한 예방 관리"],
    ["PCIe 카드 에러", "PCIe 카드, 슬롯, BIOS, 프로그램 셋팅", "점검 순서 고정 및 호환성 리스트 관리", "부품 간 호환성 검증 체계 구축"],
    ["Ethernet 불안정", "통신 케이블", "케이블 규격 고정 및 정기 점검 루틴화", "물리적 인터페이스 사양 고정"],
    ["EtherCAT Drop", "Master, Slave, ESI, 케이블", "사전 검증 및 장시간 부하 테스트 실시", "장비 출고 전 사전 검증 의무화"],
    ["USB 끊김", "USB 포트/허브, 케이블", "권장 포트 지정 및 표준 허브 사용", "연결 방식 및 자재 표준화"],
    ["Serial 장애", "Serial 카드/변환기, 통신 설정", "파라미터 표준화 및 교차 테스트", "장비별 통신 설정 표준화"],
    ["하네스 노이즈/단선", "Power/IO/ENC 하네스, 실드/접지", "규격·실드 고정 및 임의 현장 수정 제한", "하네스 제작 및 배선 기준 고정"],
    ["Vision 오동작", "카메라, 조명, 트리거", "전원/통신 분리 배선 및 노이즈 대책 수립", "Vision 시스템 구성 표준화"],
    ["외부 장비 연동", "TP, 저울, 라벨 프린터, 레이저 마킹기 등", "교체 시 사전 검토 및 사용 매뉴얼화", "임의 사양 변경 통제 및 관리"],
    ["공급/단종 리스크", "핵심 부품 전반", "리스트 관리 및 대체 후보군 상시 확보", "부품 단종 대비 리스크 관리 체계"],
]

n_rows = len(rows_data) + 1  # +1 for header
n_cols = 4
tbl_y = yt + emu(0.22)
tbl_w = CW
tbl_h = emu(0.22) * n_rows + emu(0.03)

col_widths = [Emu(emu(1.4)), Emu(emu(2.65)), Emu(emu(3.05)), Emu(emu(2.75))]

table_shape = s2.shapes.add_table(n_rows, n_cols, Emu(LM), Emu(tbl_y), Emu(tbl_w), Emu(tbl_h))
table = table_shape.table

# 열 너비
for j, w in enumerate(col_widths):
    table.columns[j].width = w

# 헤더 스타일
HEADER_BG = RGBColor(0xE6, 0x8A, 0x00)  # 오렌지/금색 (사진과 유사)
BORDER_COLOR = RGBColor(0xE0, 0xA0, 0x40)

for j, header in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = header
    r.font.name = "맑은 고딕"
    r.font.size = Pt(8)
    r.font.bold = True
    r.font.color.rgb = WHITE
    # 배경
    cell.fill.solid()
    cell.fill.fore_color.rgb = HEADER_BG
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# 데이터 행
for i, row in enumerate(rows_data):
    for j, text in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = "맑은 고딕"
        r.font.size = Pt(7.5)
        r.font.color.rgb = DARK if j == 0 else BODY
        r.font.bold = (j == 0)
        # 배경
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(0xFD, 0xF8, 0xF0)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# 테이블 테두리 설정 (오렌지)
from pptx.oxml.ns import qn
from lxml import etree

def set_cell_border(cell, color_hex="E0A040"):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = etree.SubElement(tcPr, qn(f'a:{edge}'))
        ln.set('w', '6350')  # 0.5pt
        solidFill = etree.SubElement(ln, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', color_hex)

for i in range(n_rows):
    for j in range(n_cols):
        set_cell_border(table.cell(i, j))

# 저장
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"OK: {OUTPUT} ({os.path.getsize(OUTPUT)} bytes)")
