# -*- coding: utf-8 -*-
"""
프레임그레버 이탈 이슈 PPT v3 - 슬라이드 삭제 대신 내용 교체 방식
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os, copy

TEMPLATE = r"C:\MES\wta-agents\workspaces\control-agent\template_base.pptx"
OUTPUT = r"C:\MES\wta-agents\reports\control-agent\framegrabber-issue-v2.pptx"

prs = Presentation(TEMPLATE)

# 방법: 슬라이드 4(1_Title Only)의 내용을 교체, 나머지 슬라이드 삭제
# python-pptx에서 안전한 슬라이드 삭제 방법
def delete_slide(prs, slide_index):
    """안전하게 슬라이드 삭제"""
    rId = prs.slides._sldIdLst[slide_index].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]

# 슬라이드 4(index 3)를 작업 대상으로 사용
# 슬라이드 1,2,3 삭제 (역순)
for idx in [2, 1, 0]:
    delete_slide(prs, idx)

# 이제 원래 슬라이드 4가 슬라이드 1이 됨
slide = prs.slides[0]

# 기존 placeholder 내용 교체
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 14:  # 카테고리
        ph.text = ""
        p = ph.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "[제어설계]"
        run.font.bold = True
    elif idx == 0:  # 제목
        ph.text = ""
        p = ph.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "검사기 프레임그레버(Frame Grabber) 이탈 이슈"
    elif idx == 12:  # 페이지 번호
        ph.text = ""
        p = ph.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = "1"

# 기존 TextBox들 제거 (placeholder가 아닌 것들)
shapes_to_remove = []
for shape in slide.shapes:
    if not shape.is_placeholder:
        shapes_to_remove.append(shape)

sp_tree = slide.shapes._spTree
for shape in shapes_to_remove:
    sp_tree.remove(shape._element)

# --- 헬퍼 ---
def emu(inches):
    return int(inches * 914400)

def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_rect(slide, left, top, width, height, fill_color, line_color=None, corner_size=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    # 모서리 둥글기 조절
    if corner_size is not None:
        shape.adjustments[0] = corner_size
    return shape

# --- 레이아웃 상수 ---
CONTENT_L = emu(1.75)
CONTENT_R = emu(11.6)
CONTENT_W = CONTENT_R - CONTENT_L

# === 정보 바 ===
y = emu(1.4)
add_rect(slide, CONTENT_L, y, CONTENT_W, emu(0.32), RGBColor(0xF5, 0xF5, 0xF5), RGBColor(0xE0, 0xE0, 0xE0))
tf = add_textbox(slide, CONTENT_L + emu(0.12), y + emu(0.03), CONTENT_W - emu(0.24), emu(0.26))
p = tf.paragraphs[0]
for text, bold in [("장비: ", True), ("HIM-F 검사기    ", False),
                   ("부품: ", True), ("Frame Grabber (PCIe)    ", False),
                   ("현상: ", True), ("PCIe 슬롯 이탈    ", False),
                   ("담당: ", True), ("지건승    ", False),
                   ("일자: ", True), ("2026. 04.", False)]:
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(10)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33) if bold else RGBColor(0x55, 0x55, 0x55)

# === 원인 분석 ===
y = emu(1.85)
add_rect(slide, CONTENT_L, y, CONTENT_W, emu(0.85), RGBColor(0xFF, 0xF8, 0xE1), RGBColor(0xFF, 0xE0, 0x82))

tf = add_textbox(slide, CONTENT_L + emu(0.12), y + emu(0.05), emu(2), emu(0.22))
p = tf.paragraphs[0]
run = p.add_run()
run.text = "원인 분석"
run.font.name = "맑은 고딕"
run.font.size = Pt(11)
run.font.bold = True
run.font.color.rgb = RGBColor(0xE6, 0x51, 0x00)

tf = add_textbox(slide, CONTENT_L + emu(0.12), y + emu(0.27), CONTENT_W - emu(0.24), emu(0.55))
for i, text in enumerate([
    "• 장비 운송/가동 중 진동에 의한 PCIe 슬롯 접촉 불량 발생",
    "• 프레임그레버 카드 고정 브라켓 체결 미흡 또는 나사 풀림",
    "• 반복 진동으로 카드가 점진적으로 슬롯에서 이탈 → 카메라 통신 두절",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(9.5)
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p.space_after = Pt(1)

# === Before / After ===
y_vs = emu(2.85)
gap = emu(0.12)
col_w = (CONTENT_W - gap) // 2
col_h = emu(3.35)

# --- Before ---
add_rect(slide, CONTENT_L, y_vs, col_w, col_h, RGBColor(0xFF, 0xF5, 0xF5), RGBColor(0xE8, 0xC0, 0xC0))

tf = add_textbox(slide, CONTENT_L + emu(0.12), y_vs + emu(0.06), col_w - emu(0.24), emu(0.28))
p = tf.paragraphs[0]
run = p.add_run()
run.text = "⚠ 조치 전 (Before)"
run.font.name = "맑은 고딕"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = RGBColor(0xC6, 0x28, 0x28)

# 동영상 플레이스홀더
vy = y_vs + emu(0.4)
vh = emu(1.6)
add_rect(slide, CONTENT_L + emu(0.12), vy, col_w - emu(0.24), vh,
         RGBColor(0xE8, 0xE8, 0xE8), RGBColor(0xCC, 0xCC, 0xCC))
tf = add_textbox(slide, CONTENT_L + emu(0.12), vy + emu(0.55), col_w - emu(0.24), emu(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "🎬 조치 전 동영상 삽입"
run.font.name = "맑은 고딕"
run.font.size = Pt(11)
run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

# 상세
dy = vy + vh + emu(0.08)
tf = add_textbox(slide, CONTENT_L + emu(0.12), dy, col_w - emu(0.24), emu(1.1))
for i, text in enumerate([
    "• 프레임그레버 PCIe 슬롯에서 이탈 확인",
    "• 카메라 이미지 취득 실패 / 간헐적 통신 끊김",
    "• 장비 가동 중 검사 중단 발생 (BSOD/프리징)",
    "• 고정 브라켓 나사 풀림 / 미체결 상태",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(9.5)
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p.space_after = Pt(2)

# --- After ---
ax = CONTENT_L + col_w + gap
add_rect(slide, ax, y_vs, col_w, col_h, RGBColor(0xF0, 0xFF, 0xF0), RGBColor(0xC0, 0xE0, 0xC0))

tf = add_textbox(slide, ax + emu(0.12), y_vs + emu(0.06), col_w - emu(0.24), emu(0.28))
p = tf.paragraphs[0]
run = p.add_run()
run.text = "✅ 조치 후 (After)"
run.font.name = "맑은 고딕"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

add_rect(slide, ax + emu(0.12), vy, col_w - emu(0.24), vh,
         RGBColor(0xE8, 0xE8, 0xE8), RGBColor(0xCC, 0xCC, 0xCC))
tf = add_textbox(slide, ax + emu(0.12), vy + emu(0.55), col_w - emu(0.24), emu(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "🎬 조치 후 동영상 삽입"
run.font.name = "맑은 고딕"
run.font.size = Pt(11)
run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

tf = add_textbox(slide, ax + emu(0.12), dy, col_w - emu(0.24), emu(1.1))
for i, text in enumerate([
    "• 프레임그레버 카드 재장착 및 슬롯 완전 삽입 확인",
    "• 고정 브라켓 나사 체결 + 풀림 방지 조치 적용",
    "• 카메라 이미지 취득 정상 동작 확인",
    "• 연속 가동 테스트 완료 (이탈 재발 없음)",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(9.5)
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p.space_after = Pt(2)

# === 조치 결과 ===
y_r = y_vs + col_h + emu(0.1)
add_rect(slide, CONTENT_L, y_r, CONTENT_W, emu(0.32), RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xA5, 0xD6, 0xA7))
tf = add_textbox(slide, CONTENT_L + emu(0.12), y_r + emu(0.03), CONTENT_W - emu(0.24), emu(0.26))
p = tf.paragraphs[0]
for text, bold, color in [
    ("✔ 조치 결과: ", True, RGBColor(0x2E, 0x7D, 0x32)),
    ("프레임그레버 재장착 및 고정 강화 → 카메라 통신 정상 복구, 연속 가동 안정성 확인 완료", False, RGBColor(0x44, 0x44, 0x44)),
]:
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(10)
    run.font.bold = bold
    run.font.color.rgb = color

# 저장
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"OK: {OUTPUT} ({os.path.getsize(OUTPUT)} bytes)")
