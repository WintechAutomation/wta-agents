"""문서4: AI 시스템 보안 관리 현황 및 계획 (7슬라이드) → PPTX"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_PATH = 'C:/MES/wta-agents/data/uploads/fa85f7cd-7ce6-4fe8-9a13-c9fc4d77b56d/Template.pptx'
OUTPUT_PATH = 'C:/MES/wta-agents/reports/slide-ai-security.pptx'

prs = Presentation(TEMPLATE_PATH)

RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x15, 0x65, 0xC0)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
DARK = RGBColor(0x22, 0x22, 0x22)
BODY = RGBColor(0x44, 0x44, 0x44)
GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = '맑은 고딕'

def find_layout(name_hint):
    for layout in prs.slide_layouts:
        if name_hint in layout.name:
            return layout
    return prs.slide_layouts[0]

cover_layout = find_layout('제목')
content_layout = find_layout('제목 및 내용')
blank_layout = find_layout('빈')

def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=BODY, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.font.bold = bold
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_table(slide, rows_data, left, top, width, col_widths, row_h=Cm(0.6)):
    nrows, ncols = len(rows_data), len(rows_data[0])
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, row_h * nrows)
    tbl = shape.table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x44, 0x54, 0x6A)
                else:
                    p.font.color.rgb = BODY
    return shape

def slide_title(slide, text):
    add_text(slide, Cm(5), Cm(1), Cm(20), Cm(1.5), text,
             font_size=28, bold=True, color=RED)

def slide_num(slide, num):
    add_text(slide, Cm(25), Cm(14.5), Cm(2), Cm(0.5), f'{num:02d}',
             font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

# 기존 슬라이드 삭제
while len(prs.slides) > 0:
    sld_id = prs.slides._sldIdLst[0]
    rId = sld_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') or sld_id.get('r:id')
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sld_id)


# ═══════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════
s = prs.slides.add_slide(cover_layout)
add_text(s, Cm(2), Cm(3.5), Cm(20), Cm(3),
         'AI 시스템 보안 관리\n현황 및 계획',
         font_size=36, bold=True, color=DARK)
add_text(s, Cm(2), Cm(7), Cm(20), Cm(1),
         '자가 검진 결과 및 강화 로드맵',
         font_size=18, color=GRAY)
add_text(s, Cm(2), Cm(9), Cm(20), Cm(1.5),
         '(주)윈텍오토메이션 생산관리팀 (AI운영팀)\n2026년 4월 2일',
         font_size=13, color=GRAY)


# ═══════════════════════════════════════════
# SLIDE 2: 현재 보안 체계
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '현재 보안 체계')

security_areas = [
    ('시크릿 관리', [
        '• .env 중앙관리 (Git 커밋 금지)',
        '• 커맨드라인 노출 금지',
        '• db-query.py 통한 안전한 접근',
    ], GREEN),
    ('DB 접근통제', [
        '• SQL 검증: SELECT/WITH만 허용',
        '• readonly 트랜잭션 강제',
        '• 9개 역할별 분리 계정',
    ], BLUE),
    ('권한분리', [
        '• MAX: 시스템 권한 (포트, 프로세스)',
        '• 팀원: 역할 한정 접근',
        '• PreToolUse 훅 기술적 차단',
    ], ORANGE),
    ('감사 로그', [
        '• audit_log 테이블 (사용자 활동)',
        '• agent_audit_log (에이전트 활동)',
        '• JSONL 세션 로깅 (CS 대화)',
    ], RED),
]
for i, (title, items, clr) in enumerate(security_areas):
    col = i % 2
    row = i // 2
    x = Cm(2) + col * Cm(12.5)
    y = Cm(3) + row * Cm(5.2)
    add_rect(s, x, y, Cm(11.5), Cm(4.8), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.5), y + Cm(0.2), Cm(10.5), Cm(0.5), title,
             font_size=13, bold=True, color=clr)
    add_multiline(s, x + Cm(0.5), y + Cm(0.9), Cm(10.5), Cm(3.5), items, font_size=10)
slide_num(s, 1)


# ═══════════════════════════════════════════
# SLIDE 3: 네트워크 보안
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '네트워크 보안')

# 내부 영역
add_rect(s, Cm(2), Cm(3), Cm(23), Cm(4.5), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(2.5), Cm(3.2), Cm(22), Cm(0.5), '사내 네트워크 (보호 영역)', font_size=14, bold=True, color=GREEN)

internal_table = [
    ['구성요소', '주소', '접근 범위', '암호화'],
    ['에이전트 통신', 'localhost:5600~5616', '사내 전용', 'P2P HTTP'],
    ['MES DB (PostgreSQL)', 'localhost:55432', '사내 전용', 'TLS'],
    ['ERP DB (SQL Server)', '192.168.1.201:1433', '사내망 전용', 'TLS'],
    ['MCP 메시지 큐', 'localhost JSONL', '사내 전용', '로컬 디스크'],
    ['대시보드', 'localhost:5555', 'MAX 전용', 'HTTP'],
]
add_table(s, internal_table, Cm(2.5), Cm(4), Cm(22), [Cm(5), Cm(5.5), Cm(4.5), Cm(7)], row_h=Cm(0.5))

# 외부
add_rect(s, Cm(2), Cm(8), Cm(23), Cm(5.5), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(2.5), Cm(8.2), Cm(22), Cm(0.5), '외부 접점', font_size=14, bold=True, color=ORANGE)

ext_table = [
    ['서비스', '프로토콜', '데이터 범위', '보존 기간'],
    ['Claude API (Anthropic)', 'TLS 1.3', '질문/답변 텍스트만', '30일 자동 삭제'],
    ['Slack (Socket Mode)', 'TLS 1.2+', '채널 대화 메시지', '사내 백업 유지'],
    ['Jira/Confluence', 'TLS 1.2+', '이슈/문서', '직접 관리'],
    ['MES 외부 접속', 'Cloudflare Tunnel + HTTPS', '웹 UI', '세션 기반'],
]
add_table(s, ext_table, Cm(2.5), Cm(9), Cm(22), [Cm(5.5), Cm(5), Cm(5.5), Cm(6)], row_h=Cm(0.5))

add_text(s, Cm(2), Cm(12), Cm(23), Cm(0.8),
         '핵심: DB 원본·도면·CAD 파일은 절대 외부 전송 없음. AI에는 텍스트 요약본만 전달.',
         font_size=12, bold=True, color=RED)
slide_num(s, 2)


# ═══════════════════════════════════════════
# SLIDE 4: 자가 검진 결과
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '자가 검진 결과')

# 양호 항목
add_text(s, Cm(2), Cm(3), Cm(10), Cm(0.5), '양호 항목', font_size=14, bold=True, color=GREEN)
good_items = [
    ['점검 항목', '현황', '등급'],
    ['DB 계정 분리', '9개 역할별 분리', '양호'],
    ['감사 로그', 'audit_log + agent_audit_log 운영', '양호'],
    ['ERP 접근', '읽기전용 전용 계정, .env 보관', '양호'],
    ['전송 암호화', 'Supabase TLS + Cloudflare HTTPS', '양호'],
    ['시크릿 관리', '.env 중앙관리, CLI 노출 금지', '양호'],
    ['에이전트 권한', 'PreToolUse 훅 기술적 차단', '양호'],
    ['코드 보안', 'Git 시크릿 커밋 금지 규칙', '양호'],
]
add_table(s, good_items, Cm(2), Cm(3.7), Cm(23), [Cm(6), Cm(11), Cm(6)], row_h=Cm(0.5))

# 개선 필요
add_text(s, Cm(2), Cm(8.5), Cm(10), Cm(0.5), '개선 필요 항목', font_size=14, bold=True, color=ORANGE)
warn_items = [
    ['점검 항목', '현황', '위험도', '조치 계획'],
    ['RLS 미적용', 'anon/authenticated 전 스키마 CRUD', '중간', 'Q2 RLS 적용'],
    ['벡터DB 접근', '9.3GB 외부 역할 접근 범위', '낮음', '접근 권한 세분화'],
    ['토큰 로테이션', '수동 갱신', '낮음', '자동화 예정'],
    ['DPA 미체결', 'Anthropic 법적 계약 미체결', '낮음', 'Q2 체결 추진'],
]
add_table(s, warn_items, Cm(2), Cm(9.2), Cm(23), [Cm(5), Cm(8), Cm(3), Cm(7)], row_h=Cm(0.55))
slide_num(s, 3)


# ═══════════════════════════════════════════
# SLIDE 5: 추가 보안 계획
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '추가 보안 강화 계획')

plans = [
    ('AWS VPN', [
        '• Site-to-Site IPsec IKEv2',
        '• 사내↔AWS 전용 암호화 터널',
        '• 외부 접근 원천 차단',
    ], BLUE),
    ('WAF 적용', [
        '• Cloudflare OWASP Core Ruleset',
        '• SQL Injection / XSS 차단',
        '• Rate Limiting + Bot 관리',
    ], GREEN),
    ('RLS 적용', [
        '• 테이블별 행 수준 접근제어',
        '• 에이전트별 데이터 범위 제한',
        '• anon/authenticated 권한 분리',
    ], ORANGE),
    ('토큰 자동 로테이션', [
        '• 슬랙/Atlassian 토큰 자동 갱신',
        '• 만료 알림 + 자동 교체',
        '• 유출 시 즉시 무효화',
    ], RED),
    ('벡터DB 권한 세분화', [
        '• 스키마별 접근 권한 분리',
        '• cs-agent: CS 벡터만 접근',
        '• 교차 스키마 접근 차단',
    ], BLUE),
    ('감사 로그 고도화', [
        '• 대시보드 실시간 모니터링',
        '• 이상 행동 자동 탐지',
        '• 분기별 감사 보고서 자동화',
    ], GREEN),
]
for i, (title, items, clr) in enumerate(plans):
    col = i % 3
    row = i // 3
    x = Cm(1.5) + col * Cm(8.3)
    y = Cm(3) + row * Cm(5.2)
    add_rect(s, x, y, Cm(7.5), Cm(4.8), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.3), y + Cm(0.2), Cm(6.9), Cm(0.5), title,
             font_size=12, bold=True, color=clr)
    add_multiline(s, x + Cm(0.3), y + Cm(0.9), Cm(6.9), Cm(3.5), items, font_size=10)
slide_num(s, 4)


# ═══════════════════════════════════════════
# SLIDE 6: 보안 로드맵
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '보안 로드맵')

# 즉시
add_rect(s, Cm(2), Cm(3), Cm(23), Cm(2.5), RGBColor(0xFF, 0xEB, 0xEE), RED)
add_text(s, Cm(2.5), Cm(3.2), Cm(3), Cm(0.5), '즉시', font_size=14, bold=True, color=RED)
add_multiline(s, Cm(6), Cm(3.2), Cm(18), Cm(2), [
    '• PostgreSQL RLS 적용 — 테이블별 행 수준 접근제어',
    '• 벡터DB 접근 제한 — 스키마별 에이전트 권한 분리',
    '• 민감 데이터 자동 필터링 — 도면번호·고객코드 마스킹',
], font_size=11)

# Q2
add_rect(s, Cm(2), Cm(6), Cm(23), Cm(2.5), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(6.2), Cm(3), Cm(0.5), 'Q2 2026', font_size=14, bold=True, color=BLUE)
add_multiline(s, Cm(6), Cm(6.2), Cm(18), Cm(2), [
    '• AWS Site-to-Site VPN 구축 — IPsec IKEv2 전용 터널',
    '• Cloudflare WAF 적용 — OWASP Core Ruleset + Rate Limiting',
    '• GeoIP 제한 — 허용 국가: KR/CN/JP/AT/US',
    '• Anthropic DPA 체결 — 법적 데이터 보호 계약',
], font_size=11)

# Q3
add_rect(s, Cm(2), Cm(9), Cm(23), Cm(2.2), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(2.5), Cm(9.2), Cm(3), Cm(0.5), 'Q3 2026', font_size=14, bold=True, color=ORANGE)
add_multiline(s, Cm(6), Cm(9.2), Cm(18), Cm(1.5), [
    '• 토큰 자동 로테이션 — 슬랙/Atlassian 자동 갱신',
    '• 감사 로그 대시보드 — 실시간 이상 행동 탐지',
], font_size=11)

# 지속
add_rect(s, Cm(2), Cm(11.7), Cm(23), Cm(2), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(2.5), Cm(11.9), Cm(3), Cm(0.5), '지속', font_size=14, bold=True, color=GREEN)
add_multiline(s, Cm(6), Cm(11.9), Cm(18), Cm(1.5), [
    '• 분기별 자가 검진 + 침투 테스트',
    '• ISO 27001 준비 (2026 Q4 목표)',
], font_size=11)
slide_num(s, 5)


# ═══════════════════════════════════════════
# SLIDE 7: 엔딩
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)


# ═══ 저장 ═══
prs.save(OUTPUT_PATH)
print(f'PPTX 저장 완료: {OUTPUT_PATH}')
print(f'슬라이드 수: {len(prs.slides)}')
