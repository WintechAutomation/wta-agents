"""문서2: CS 업무용 웹 시스템 개설 보고 (7슬라이드) → PPTX"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_PATH = 'C:/MES/wta-agents/data/uploads/fa85f7cd-7ce6-4fe8-9a13-c9fc4d77b56d/Template.pptx'
OUTPUT_PATH = 'C:/MES/wta-agents/reports/slide-cs-webpage.pptx'

prs = Presentation(TEMPLATE_PATH)

RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x15, 0x65, 0xC0)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
DARK = RGBColor(0x22, 0x22, 0x22)
BODY = RGBColor(0x44, 0x44, 0x44)
GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = '맑은 고딕'

def find_layout(name_hint):
    for layout in prs.slide_layouts:
        if name_hint in layout.name:
            return layout
    return prs.slide_layouts[0]

cover_layout = find_layout('제목')
content_layout = find_layout('제목 및 내용')
blank_layout = find_layout('빈')

def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=BODY, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.font.bold = bold
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_table(slide, rows_data, left, top, width, col_widths, row_h=Cm(0.6)):
    nrows, ncols = len(rows_data), len(rows_data[0])
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, row_h * nrows)
    tbl = shape.table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x44, 0x54, 0x6A)
                else:
                    p.font.color.rgb = BODY
    return shape

def slide_title(slide, text):
    add_text(slide, Cm(5), Cm(1), Cm(20), Cm(1.5), text,
             font_size=28, bold=True, color=RED)

def slide_num(slide, num):
    add_text(slide, Cm(25), Cm(14.5), Cm(2), Cm(0.5), f'{num:02d}',
             font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

# 기존 슬라이드 삭제
while len(prs.slides) > 0:
    sld_id = prs.slides._sldIdLst[0]
    rId = sld_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') or sld_id.get('r:id')
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sld_id)


# ═══════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════
s = prs.slides.add_slide(cover_layout)
add_text(s, Cm(2), Cm(3.5), Cm(20), Cm(3),
         'CS 업무용 웹 시스템\n개설 보고',
         font_size=36, bold=True, color=DARK)
add_text(s, Cm(2), Cm(7), Cm(20), Cm(1),
         'cs-wta.com / MES 통합 플랫폼',
         font_size=18, color=GRAY)
add_text(s, Cm(2), Cm(9), Cm(20), Cm(1.5),
         '(주)윈텍오토메이션 생산관리팀 (AI운영팀)\n2026년 4월 2일',
         font_size=13, color=GRAY)


# ═══════════════════════════════════════════
# SLIDE 2: 시스템 개요
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '시스템 개요')

# 기술 스택
stack_table = [
    ['구분', '기술', '설명'],
    ['프론트엔드', 'React 19 + TypeScript 5.9', 'SPA, shadcn/ui + Tailwind CSS 4'],
    ['백엔드', 'Go 1.26 + Gin 1.12', 'REST API + WebSocket'],
    ['데이터베이스', 'PostgreSQL (Supabase)', '관계형 + pgvector (벡터 검색)'],
    ['ERP 연동', 'SQL Server (읽기전용)', 'mirae 스키마, 사내망 전용'],
    ['데이터 그리드', 'AG Grid', '대량 데이터 가상 스크롤'],
    ['실시간', 'gorilla/websocket', '실시간 알림·상태 업데이트'],
    ['파일 저장', 'Supabase Storage', '문서·이미지 관리'],
    ['접속', 'Cloudflare Tunnel', 'https://mes-wta.com'],
]
add_table(s, stack_table, Cm(2), Cm(3), Cm(23), [Cm(4), Cm(7.5), Cm(11.5)], row_h=Cm(0.55))

add_rect(s, Cm(2), Cm(9.5), Cm(23), Cm(4), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(9.7), Cm(22), Cm(0.5), '시스템 특징', font_size=12, bold=True, color=BLUE)
add_multiline(s, Cm(2.5), Cm(10.4), Cm(22), Cm(2.5), [
    '• Django+JS 레거시 wMES를 Go+TypeScript로 완전 재구축 (기능·UX 동일 유지)',
    '• 사내 MES와 CS 웹을 하나의 코드베이스로 통합 — 유지보수 효율 극대화',
    '• AI 에이전트가 직접 활용하는 API·DB 구조로 설계 (ai_agent_action_log 등)',
], font_size=10)
slide_num(s, 1)


# ═══════════════════════════════════════════
# SLIDE 3: 구현 규모
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '구현 규모')

# 큰 숫자
nums = [
    ('119', RED, '구현 페이지'),
    ('73', BLUE, 'API 핸들러'),
    ('47', GREEN, 'CS 테이블'),
]
for i, (val, clr, lbl) in enumerate(nums):
    x = Cm(2) + i * Cm(8.3)
    add_rect(s, x, Cm(3), Cm(7.3), Cm(2.5), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x, Cm(3.2), Cm(7.3), Cm(1.2), val, font_size=40, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_text(s, x, Cm(4.5), Cm(7.3), Cm(0.5), lbl, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# 모듈별 페이지 수
module_table = [
    ['모듈', '페이지 수', '주요 기능'],
    ['생산관리', '14', '작업지시, 생산실적, 공정관리, LOT 추적'],
    ['원가관리', '18', '원가계산, 자재비, 인건비, 경비, 손익'],
    ['자재·재고', '11', '재고 현황, 입출고, 수불, BOM'],
    ['품질관리', '10', '검사, 부적합, 불량분석, 체크리스트'],
    ['ERP 연동', '10', '수주, 거래처, 제품, 재고, 원가'],
    ['설계', '9', '도면관리, BOM, ECN, 3D뷰어'],
    ['시스템', '19', '사용자, 권한, 로그, 대시보드, AI관리'],
    ['CS', '6', 'CS이력, 장비, 견적, 에러코드'],
    ['일정', '6', '프로젝트, 마일스톤, 간트차트'],
    ['영업·출하', '7', '수주관리, 견적, 출하검사'],
    ['유틸리티', '4', '파일관리, 라벨인쇄, 알림'],
]
add_table(s, module_table, Cm(2), Cm(6.5), Cm(23), [Cm(4), Cm(3), Cm(16)], row_h=Cm(0.5))
slide_num(s, 2)


# ═══════════════════════════════════════════
# SLIDE 4: CS 전용 기능
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, 'CS 전용 기능')

add_rect(s, Cm(2), Cm(3), Cm(23), Cm(1.5), RGBColor(0xFF, 0xEB, 0xEE), RED)
add_text(s, Cm(2.5), Cm(3.2), Cm(22), Cm(1),
         'csagent 스키마: 47개 테이블 / CS이력 5,160건 축적 / 견적 243건',
         font_size=13, bold=True, color=RED)

# 기능 카드
features = [
    ('CS 이력관리', 'cs-agent 연동, 자재/인건비/경비\n자동 연결, 상세 이력 추적', BLUE),
    ('장비관리', '고객사별 장비 현황, 설치일·\n보증기간·정비 이력 관리', GREEN),
    ('재고/부품', '부품 재고 실시간 조회,\n사용 이력·단가 관리', ORANGE),
    ('견적관리', '243건 견적 데이터, 자재·\n인건비·경비 항목별 산출', RED),
    ('AI Q&A', 'qa_cache + qa_feedback\nRAG 기반 자동 응답·피드백', BLUE),
    ('에러코드 DB', '장비별 에러코드 사전,\n증상·원인·조치 방법', GREEN),
    ('라벨 인쇄', '부품·장비 라벨 자동 생성,\nQR코드 포함', ORANGE),
    ('감사 로그', 'audit_log + agent_audit_log\n모든 작업 기록 추적', RED),
]
for i, (title, desc, clr) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Cm(1.5) + col * Cm(6.2)
    y = Cm(5) + row * Cm(4)
    add_rect(s, x, y, Cm(5.8), Cm(3.5), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.3), y + Cm(0.2), Cm(5.2), Cm(0.5), title,
             font_size=11, bold=True, color=clr)
    add_text(s, x + Cm(0.3), y + Cm(0.9), Cm(5.2), Cm(2.2), desc,
             font_size=9, color=BODY)
slide_num(s, 3)


# ═══════════════════════════════════════════
# SLIDE 5: ERP 연동
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, 'ERP 연동')

add_text(s, Cm(2), Cm(3.2), Cm(23), Cm(0.8),
         'SQL Server 192.168.1.201 (mirae 스키마) — 사내망 전용, 읽기전용 접근',
         font_size=13, bold=True, color=BODY)

# 연동 구조
add_rect(s, Cm(2), Cm(4.5), Cm(7), Cm(3), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(4.7), Cm(6), Cm(0.5), 'MES 백엔드 (Go)', font_size=12, bold=True, color=BLUE)
add_multiline(s, Cm(2.5), Cm(5.4), Cm(6), Cm(1.5), [
    '• sqlx + pgx v5',
    '• 타입 안전 쿼리',
    '• 연결 풀 관리',
], font_size=10)

add_text(s, Cm(9.5), Cm(5.5), Cm(1.5), Cm(0.8), '↔', font_size=20, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)

add_rect(s, Cm(11.5), Cm(4.5), Cm(7), Cm(3), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(12), Cm(4.7), Cm(6), Cm(0.5), 'ERP DB (SQL Server)', font_size=12, bold=True, color=ORANGE)
add_multiline(s, Cm(12), Cm(5.4), Cm(6), Cm(1.5), [
    '• 읽기전용 전용 계정',
    '• SELECT만 허용',
    '• 사내망 전용 접근',
], font_size=10)

add_rect(s, Cm(19), Cm(4.5), Cm(6), Cm(3), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(19.5), Cm(4.7), Cm(5), Cm(0.5), 'PostgreSQL', font_size=12, bold=True, color=GREEN)
add_multiline(s, Cm(19.5), Cm(5.4), Cm(5), Cm(1.5), [
    '• MES 자체 데이터',
    '• pgvector (RAG)',
    '• Supabase 호스팅',
], font_size=10)

# ERP 데이터 활용 영역
erp_table = [
    ['ERP 데이터', 'MES 활용', '담당 에이전트'],
    ['수주 데이터', '일정관리·납기 추적·생산계획', 'schedule-agent'],
    ['재고 현황', '자재소요 분석·부족 알림', 'db-manager'],
    ['원가 정보', '원가계산 자동화·손익분석', 'dev-agent (MES)'],
    ['거래처 정보', '고객사별 장비·CS이력 연결', 'cs-agent'],
    ['제품 마스터', 'BOM·도면·사양 통합 관리', 'design-agent'],
]
add_table(s, erp_table, Cm(2), Cm(8.5), Cm(23), [Cm(5), Cm(10), Cm(8)])

add_rect(s, Cm(2), Cm(12.5), Cm(23), Cm(1.2), RGBColor(0xFF, 0xEB, 0xEE), RED)
add_text(s, Cm(2.5), Cm(12.6), Cm(22), Cm(0.8),
         '보안 원칙: ERP DB는 절대 쓰기 불가 — AI 에이전트도 SELECT만 실행 (db-query.py 검증)',
         font_size=11, bold=True, color=RED)
slide_num(s, 4)


# ═══════════════════════════════════════════
# SLIDE 6: 향후 확장
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '향후 확장 계획')

# AWS 이전
add_rect(s, Cm(2), Cm(3), Cm(11), Cm(5), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(3.2), Cm(10), Cm(0.5), 'AWS 클라우드 이전', font_size=13, bold=True, color=BLUE)
add_multiline(s, Cm(2.5), Cm(4), Cm(10), Cm(3.5), [
    '• ECS Fargate — 서버리스 컨테이너',
    '• RDS PostgreSQL — 관리형 DB',
    '• Site-to-Site VPN — 사내↔AWS 전용 터널',
    '• CloudFront + WAF — CDN + 보안',
    '• 가용성 99.9%+ 보장',
], font_size=10)

# 모바일
add_rect(s, Cm(14), Cm(3), Cm(11), Cm(5), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(14.5), Cm(3.2), Cm(10), Cm(0.5), '모바일 대응', font_size=13, bold=True, color=GREEN)
add_multiline(s, Cm(14.5), Cm(4), Cm(10), Cm(3.5), [
    '• 반응형 UI (태블릿/모바일)',
    '• PWA — 앱 설치 없이 모바일 접근',
    '• 오프라인 지원 (Service Worker)',
    '• 현장 작업자 태블릿 활용',
    '• QR 스캔 → 장비 정보 즉시 조회',
], font_size=10)

# 대시보드 고도화
add_rect(s, Cm(2), Cm(8.5), Cm(11), Cm(5), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(2.5), Cm(8.7), Cm(10), Cm(0.5), '대시보드 고도화', font_size=13, bold=True, color=ORANGE)
add_multiline(s, Cm(2.5), Cm(9.4), Cm(10), Cm(3.5), [
    '• 실시간 KPI 대시보드',
    '• 생산 현황 모니터링',
    '• 품질 트렌드 시각화',
    '• AI 에이전트 활동 현황',
    '• 경영진 전용 요약 뷰',
], font_size=10)

# 확장 로드맵
add_rect(s, Cm(14), Cm(8.5), Cm(11), Cm(5), RGBColor(0xFF, 0xEB, 0xEE), RED)
add_text(s, Cm(14.5), Cm(8.7), Cm(10), Cm(0.5), '확장 로드맵', font_size=13, bold=True, color=RED)
add_multiline(s, Cm(14.5), Cm(9.4), Cm(10), Cm(3.5), [
    '• Q2 2026: 모바일 반응형 완성',
    '• Q3 2026: AWS 이전 + 고가용성',
    '• Q4 2026: 실시간 KPI + 경영 대시보드',
    '• 2027: 협력업체 포탈 오픈',
], font_size=10)
slide_num(s, 5)


# ═══════════════════════════════════════════
# SLIDE 7: 엔딩
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)


# ═══ 저장 ═══
prs.save(OUTPUT_PATH)
print(f'PPTX 저장 완료: {OUTPUT_PATH}')
print(f'슬라이드 수: {len(prs.slides)}')
