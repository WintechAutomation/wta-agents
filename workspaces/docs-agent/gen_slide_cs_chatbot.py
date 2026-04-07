"""문서1: 사내 CS AI 챗봇 시스템 구축 보고 (8슬라이드) → PPTX"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_PATH = 'C:/MES/wta-agents/data/uploads/fa85f7cd-7ce6-4fe8-9a13-c9fc4d77b56d/Template.pptx'
OUTPUT_PATH = 'C:/MES/wta-agents/reports/slide-cs-chatbot.pptx'

prs = Presentation(TEMPLATE_PATH)

# 색상
RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x15, 0x65, 0xC0)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
DARK = RGBColor(0x22, 0x22, 0x22)
BODY = RGBColor(0x44, 0x44, 0x44)
GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = '맑은 고딕'

def find_layout(name_hint):
    for layout in prs.slide_layouts:
        if name_hint in layout.name:
            return layout
    return prs.slide_layouts[0]

cover_layout = find_layout('제목')
content_layout = find_layout('제목 및 내용')
blank_layout = find_layout('빈')

def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=BODY, bold=False, spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.font.bold = bold
        if spacing:
            p.space_after = Pt(spacing)
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_table(slide, rows_data, left, top, width, col_widths, row_h=Cm(0.6)):
    nrows, ncols = len(rows_data), len(rows_data[0])
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, row_h * nrows)
    tbl = shape.table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x44, 0x54, 0x6A)
                else:
                    p.font.color.rgb = BODY
    return shape

def slide_title(slide, text):
    add_text(slide, Cm(5), Cm(1), Cm(20), Cm(1.5), text,
             font_size=28, bold=True, color=RED)

def slide_num(slide, num):
    add_text(slide, Cm(25), Cm(14.5), Cm(2), Cm(0.5), f'{num:02d}',
             font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

# 기존 슬라이드 삭제
while len(prs.slides) > 0:
    sld_id = prs.slides._sldIdLst[0]
    rId = sld_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') or sld_id.get('r:id')
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sld_id)


# ═══════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════
s = prs.slides.add_slide(cover_layout)
add_text(s, Cm(2), Cm(3.5), Cm(20), Cm(3),
         '사내 CS AI 챗봇 시스템\n구축 보고',
         font_size=36, bold=True, color=DARK)
add_text(s, Cm(2), Cm(7), Cm(20), Cm(1),
         '데이터 학습 현황 및 운영 성과',
         font_size=18, color=GRAY)
add_text(s, Cm(2), Cm(9), Cm(20), Cm(1.5),
         '(주)윈텍오토메이션 생산관리팀 (AI운영팀)\n2026년 4월 2일',
         font_size=13, color=GRAY)


# ═══════════════════════════════════════════
# SLIDE 2: 시스템 개요
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '시스템 개요')

add_rect(s, Cm(2), Cm(3), Cm(23), Cm(2.5), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(3.2), Cm(22), Cm(0.5), '목적', font_size=14, bold=True, color=BLUE)
add_multiline(s, Cm(2.5), Cm(3.8), Cm(22), Cm(1.5), [
    '• CS 기술지원 응답 자동화 — 축적된 장비 노하우 데이터를 AI가 즉시 검색·응답',
    '• 슬랙 기반 실시간 소통 — 직원이 평소 쓰는 채널에서 바로 질문, AI가 바로 답변',
], font_size=11)

# 아키텍처
add_text(s, Cm(2), Cm(6), Cm(23), Cm(0.5), '시스템 아키텍처', font_size=14, bold=True, color=DARK)

flow_boxes = [
    ('슬랙\n(Socket Mode)', BLUE),
    ('slack-bot\n(Haiku)', GREEN),
    ('cs-agent\n(Opus)', RED),
    ('RAG 벡터검색\n(pgvector)', ORANGE),
    ('슬랙 응답\n(자동)', BLUE),
]
for i, (label, clr) in enumerate(flow_boxes):
    x = Cm(1.5) + i * Cm(5)
    add_rect(s, x, Cm(7), Cm(4.2), Cm(1.8), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.2), Cm(7.1), Cm(3.8), Cm(1.6), label,
             font_size=9, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    if i < len(flow_boxes) - 1:
        add_text(s, x + Cm(4.2), Cm(7.5), Cm(0.8), Cm(0.8), '→',
                 font_size=16, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)

# 특징
add_rect(s, Cm(2), Cm(9.5), Cm(11), Cm(4), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(2.5), Cm(9.7), Cm(10), Cm(0.5), '통신 방식', font_size=12, bold=True, color=GREEN)
add_multiline(s, Cm(2.5), Cm(10.4), Cm(10), Cm(2.5), [
    '• P2P HTTP 직접통신 (중앙 서버 불필요)',
    '• 오프라인 메시지 큐 (자동 재전달)',
    '• 에이전트 간 MCP 채널 연동',
], font_size=10)

add_rect(s, Cm(14), Cm(9.5), Cm(11), Cm(4), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(14.5), Cm(9.7), Cm(10), Cm(0.5), '멀티모델 파이프라인', font_size=12, bold=True, color=ORANGE)
add_multiline(s, Cm(14.5), Cm(10.4), Cm(10), Cm(2.5), [
    '• Haiku — 접수·라우팅 (빠른 분류)',
    '• Opus — 고급 기술응답 (이미지 분석)',
    '• Sonnet — 일반 업무처리',
], font_size=10)
slide_num(s, 1)


# ═══════════════════════════════════════════
# SLIDE 3: 데이터 학습 현황
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '데이터 학습 현황')

# 큰 숫자 — 총 규모
nums = [
    ('9.3GB+', RED, '벡터 DB 총 용량'),
    ('391,000+', BLUE, '총 벡터 청크 수'),
    ('99.7%', GREEN, '임베딩 완료율'),
]
for i, (val, clr, lbl) in enumerate(nums):
    x = Cm(2) + i * Cm(8.3)
    add_rect(s, x, Cm(3), Cm(7.3), Cm(2.5), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x, Cm(3.2), Cm(7.3), Cm(1.2), val, font_size=32, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_text(s, x, Cm(4.5), Cm(7.3), Cm(0.5), lbl, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

# 상세 테이블
data_table = [
    ['데이터셋', '청크 수', '용량', '상태'],
    ['CS이력 RAG', '265,635', '6.1GB', '✅ 완료'],
    ['WTA 자체 매뉴얼', '120,492', '3.1GB', '✅ 완료 (656/656)'],
    ['기술문서', '2,095', '-', '✅ 완료'],
    ['부품매뉴얼 (892개)', '633', '-', '✅ 완료'],
    ['CS이력 원본 임베딩', '3,318건', '-', '✅ 완료'],
]
add_table(s, data_table, Cm(2), Cm(6.5), Cm(23), [Cm(7), Cm(5), Cm(4), Cm(7)])

# 임베딩 모델 정보
add_rect(s, Cm(2), Cm(11.5), Cm(23), Cm(2), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(11.7), Cm(22), Cm(0.5), '임베딩 모델', font_size=12, bold=True, color=BLUE)
add_multiline(s, Cm(2.5), Cm(12.3), Cm(22), Cm(1), [
    '• Qwen3-Embedding-8B (2000차원) — 한국어·영어·중국어 고성능 임베딩',
    '• PostgreSQL pgvector 기반 코사인 유사도 Top-K 검색 — 벡터 인덱스 최적화',
], font_size=10)
slide_num(s, 2)


# ═══════════════════════════════════════════
# SLIDE 4: 슬랙 채널별 AI 라우팅
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '슬랙 채널별 AI 자동 라우팅')

add_text(s, Cm(2), Cm(3.2), Cm(23), Cm(0.8),
         '각 슬랙 채널에 전담 AI 에이전트를 배치하여, 직원 대화가 자동으로 업무 데이터로 전환됩니다.',
         font_size=13, color=BODY)

routing_table = [
    ['슬랙 채널', 'AI 에이전트', '주요 기능'],
    ['#cs', '🛠️ cs-agent (Opus)', 'CS 기술지원 — RAG 검색 + 이미지 분석'],
    ['#부적합', '🔍 nc-manager', 'LLM 자동판단 → MES 부적합 등록'],
    ['#영업', '💰 sales-agent', '수주·견적·매출 조회 및 분석'],
    ['#개발', '💻 dev-agent', 'MES 버그 리포트·기능 요청 처리'],
    ['#docs', '📝 docs-agent', '문서 작성·번역·매뉴얼 검색'],
    ['#제어-*', '⚡ control-agent', '개인별 제어설계 기술지원'],
    ['#admin', '📋 admin-agent', '인프라 모니터링·계정 관리'],
]
add_table(s, routing_table, Cm(2), Cm(4.5), Cm(23), [Cm(4), Cm(7), Cm(12)])

# 비전 박스
add_rect(s, Cm(2), Cm(10.5), Cm(23), Cm(3), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(2.5), Cm(10.7), Cm(22), Cm(0.5), '핵심 비전: 직원 대화 → AI 자동 수집 → 정형 데이터 → MES DB', font_size=13, bold=True, color=ORANGE)
add_multiline(s, Cm(2.5), Cm(11.5), Cm(22), Cm(1.5), [
    '• 직원이 슬랙에서 자연스럽게 대화하면, AI가 자동으로 구조화된 데이터를 생성',
    '• 부적합 보고서, CS 이력, 장비 이슈 등이 수동 입력 없이 MES에 축적',
    '• 현실 직원 → AI 데이터 채널 확보 (전사 핵심 과제)',
], font_size=10)
slide_num(s, 3)


# ═══════════════════════════════════════════
# SLIDE 5: 운영 성과
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '운영 성과')

# 핵심 수치
metrics = [
    ('154건', RED, 'CS 챗봇 세션'),
    ('614건', BLUE, '처리 메시지'),
    ('100%', GREEN, '긍정 피드백'),
]
for i, (val, clr, lbl) in enumerate(metrics):
    x = Cm(2) + i * Cm(8.3)
    add_rect(s, x, Cm(3), Cm(7.3), Cm(2.5), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x, Cm(3.2), Cm(7.3), Cm(1.2), val, font_size=36, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_text(s, x, Cm(4.5), Cm(7.3), Cm(0.5), lbl, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

# Before / After 비교
add_rect(s, Cm(2), Cm(6.5), Cm(11), Cm(5), RGBColor(0xFF, 0xF5, 0xF5), RED)
add_text(s, Cm(2.5), Cm(6.7), Cm(10), Cm(0.5), 'Before (기존 방식)', font_size=13, bold=True, color=RED)
add_multiline(s, Cm(2.5), Cm(7.5), Cm(10), Cm(3.5), [
    '✗ CS 응답: 수 시간 ~ 1일 소요',
    '✗ 매뉴얼 검색: 수십 개 PDF 수동 탐색',
    '✗ 부적합 등록: 수동 양식 작성 (5~10분)',
    '✗ 기술 노하우: 특정 담당자 의존',
    '✗ 야간/주말: 응답 불가',
], font_size=10, color=RED)

add_rect(s, Cm(14), Cm(6.5), Cm(11), Cm(5), RGBColor(0xF0, 0xFF, 0xF0), GREEN)
add_text(s, Cm(14.5), Cm(6.7), Cm(10), Cm(0.5), 'After (AI 챗봇)', font_size=13, bold=True, color=GREEN)
add_multiline(s, Cm(14.5), Cm(7.5), Cm(10), Cm(3.5), [
    '✓ CS 응답: 수 분 이내 자동 응답',
    '✓ RAG 검색: 39만 청크 즉시 탐색',
    '✓ 부적합: 슬랙 대화에서 자동 파싱·등록',
    '✓ 노하우: AI에 축적, 누구나 접근 가능',
    '✓ 24/7: 야간·주말 무인 응답',
], font_size=10, color=GREEN)

# 기간
add_text(s, Cm(2), Cm(12), Cm(23), Cm(0.5),
         '운영 기간: 2025년 5월 ~ 2026년 3월 (11개월)',
         font_size=11, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)
slide_num(s, 4)


# ═══════════════════════════════════════════
# SLIDE 6: 부적합 자동 등록 상세
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '부적합 자동 등록 프로세스')

add_text(s, Cm(2), Cm(3.2), Cm(23), Cm(0.8),
         '슬랙 #부적합 채널에서 직원이 보고하면, AI가 자동으로 구조화하여 MES에 등록합니다.',
         font_size=13, color=BODY)

# 프로세스 흐름
steps = [
    ('① 직원 보고', '슬랙 #부적합에\n자유 형식 작성', BLUE),
    ('② AI 파싱', 'LLM이 장비·모델·\n증상·조치 자동 분류', ORANGE),
    ('③ MES 등록', '부적합 테이블에\n자동 INSERT', GREEN),
    ('④ 알림', '관련 담당자에게\n자동 알림', RED),
]
for i, (title, desc, clr) in enumerate(steps):
    x = Cm(1.5) + i * Cm(6.2)
    add_rect(s, x, Cm(4.5), Cm(5.5), Cm(3.5), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.3), Cm(4.7), Cm(4.9), Cm(0.6), title, font_size=12, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Cm(0.3), Cm(5.5), Cm(4.9), Cm(2), desc, font_size=10, color=BODY, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(s, x + Cm(5.5), Cm(5.8), Cm(0.7), Cm(0.8), '→',
                 font_size=16, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)

# 효과
add_rect(s, Cm(2), Cm(9), Cm(23), Cm(4.5), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(2.5), Cm(9.2), Cm(22), Cm(0.5), '도입 효과', font_size=14, bold=True, color=GREEN)

effects = [
    ['항목', '기존', 'AI 적용 후', '개선율'],
    ['등록 소요시간', '5~10분 (수동)', '30초 이내 (자동)', '95% 단축'],
    ['누락률', '20~30% 추정', '0% (전수 등록)', '100% 개선'],
    ['데이터 표준화', '담당자별 상이', '100% 표준 양식', '완전 표준화'],
    ['축적 데이터', '684건 (부적합 보고서)', '지속 자동 누적', '실시간 축적'],
]
add_table(s, effects, Cm(2.5), Cm(10), Cm(22), [Cm(5), Cm(5.5), Cm(5.5), Cm(6)])
slide_num(s, 5)


# ═══════════════════════════════════════════
# SLIDE 7: 향후 계획
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '향후 계획')

# 멀티모델 파이프라인 고도화
add_rect(s, Cm(2), Cm(3), Cm(11), Cm(5), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(3.2), Cm(10), Cm(0.5), '멀티모델 파이프라인 고도화', font_size=12, bold=True, color=BLUE)
add_multiline(s, Cm(2.5), Cm(4), Cm(10), Cm(3.5), [
    '• Haiku: 접수·라우팅 (빠른 분류)',
    '• Opus: 이미지 분석 (장비 사진·도면)',
    '• Sonnet: 일반 기술응답',
    '→ 비용 최적화 + 품질 극대화',
], font_size=10)

# 자가학습
add_rect(s, Cm(14), Cm(3), Cm(11), Cm(5), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(14.5), Cm(3.2), Cm(10), Cm(0.5), 'JSONL 축적 → 자가학습', font_size=12, bold=True, color=ORANGE)
add_multiline(s, Cm(14.5), Cm(4), Cm(10), Cm(3.5), [
    '• 모든 CS 세션을 JSONL 로그로 축적',
    '• 반복 패턴 자동 추출 → 스킬 생성',
    '• 시간이 갈수록 응답 정확도 향상',
    '→ 운영할수록 똑똑해지는 시스템',
], font_size=10)

# 다국어
add_rect(s, Cm(2), Cm(8.5), Cm(11), Cm(4.5), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(2.5), Cm(8.7), Cm(10), Cm(0.5), '다국어 CS 대응 강화', font_size=12, bold=True, color=GREEN)
add_multiline(s, Cm(2.5), Cm(9.4), Cm(10), Cm(3), [
    '• 7개국 60+ 고객사 기술지원',
    '• 중국/일본/유럽 고객사 직접 응답',
    '• 매뉴얼 자동 번역 (한↔영↔중)',
    '→ 해외 CS 대응시간 혁신적 단축',
], font_size=10)

# 데이터 현황 요약
add_rect(s, Cm(14), Cm(8.5), Cm(11), Cm(4.5), RGBColor(0xFF, 0xEB, 0xEE), RED)
add_text(s, Cm(14.5), Cm(8.7), Cm(10), Cm(0.5), '학습 데이터 완성도', font_size=12, bold=True, color=RED)
add_multiline(s, Cm(14.5), Cm(9.4), Cm(10), Cm(3), [
    '• 부품매뉴얼: 892/892개 (100%)',
    '• WTA매뉴얼: 656/656파일 (99.7%)',
    '• CS이력: 3,318건 임베딩 완료',
    '→ 핵심 데이터 학습 거의 완료',
], font_size=10)
slide_num(s, 6)


# ═══════════════════════════════════════════
# SLIDE 8: 엔딩
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)


# ═══ 저장 ═══
prs.save(OUTPUT_PATH)
print(f'PPTX 저장 완료: {OUTPUT_PATH}')
print(f'슬라이드 수: {len(prs.slides)}')
