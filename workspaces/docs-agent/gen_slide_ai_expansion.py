"""문서3: 제조 업무 AI 확장 계획 보고 (8슬라이드) → PPTX"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_PATH = 'C:/MES/wta-agents/data/uploads/fa85f7cd-7ce6-4fe8-9a13-c9fc4d77b56d/Template.pptx'
OUTPUT_PATH = 'C:/MES/wta-agents/reports/slide-ai-expansion.pptx'

prs = Presentation(TEMPLATE_PATH)

RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x15, 0x65, 0xC0)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
DARK = RGBColor(0x22, 0x22, 0x22)
BODY = RGBColor(0x44, 0x44, 0x44)
GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = '맑은 고딕'

def find_layout(name_hint):
    for layout in prs.slide_layouts:
        if name_hint in layout.name:
            return layout
    return prs.slide_layouts[0]

cover_layout = find_layout('제목')
content_layout = find_layout('제목 및 내용')
blank_layout = find_layout('빈')

def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=BODY, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.font.bold = bold
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_table(slide, rows_data, left, top, width, col_widths, row_h=Cm(0.6)):
    nrows, ncols = len(rows_data), len(rows_data[0])
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, row_h * nrows)
    tbl = shape.table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x44, 0x54, 0x6A)
                else:
                    p.font.color.rgb = BODY
    return shape

def slide_title(slide, text):
    add_text(slide, Cm(5), Cm(1), Cm(20), Cm(1.5), text,
             font_size=28, bold=True, color=RED)

def slide_num(slide, num):
    add_text(slide, Cm(25), Cm(14.5), Cm(2), Cm(0.5), f'{num:02d}',
             font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

# 기존 슬라이드 삭제
while len(prs.slides) > 0:
    sld_id = prs.slides._sldIdLst[0]
    rId = sld_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') or sld_id.get('r:id')
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sld_id)


# ═══════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════
s = prs.slides.add_slide(cover_layout)
add_text(s, Cm(2), Cm(3.5), Cm(20), Cm(3),
         '제조 업무 AI 확장 계획',
         font_size=36, bold=True, color=DARK)
add_text(s, Cm(2), Cm(7), Cm(20), Cm(1),
         'AI 에이전트 기반 스마트 팩토리 로드맵',
         font_size=18, color=GRAY)
add_text(s, Cm(2), Cm(9), Cm(20), Cm(1.5),
         '(주)윈텍오토메이션 생산관리팀 (AI운영팀)\n2026년 4월 2일',
         font_size=13, color=GRAY)


# ═══════════════════════════════════════════
# SLIDE 2: 현재 AI 팀 구성
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '현재 AI 팀 구성 (15명)')

# 운영
add_text(s, Cm(2), Cm(3), Cm(23), Cm(0.5), '운영 그룹', font_size=13, bold=True, color=RED)
ops = [
    ('👑 MAX', '총괄 오케스트레이터', RED),
    ('📋 admin', '인프라·모니터링', BLUE),
    ('💬 slack-bot', '메시징·라우팅', GREEN),
    ('📝 docs', '문서·번역·RAG', ORANGE),
]
for i, (name, role, clr) in enumerate(ops):
    x = Cm(2) + i * Cm(6.2)
    add_rect(s, x, Cm(3.7), Cm(5.6), Cm(1.5), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.3), Cm(3.8), Cm(5), Cm(0.5), name, font_size=10, bold=True, color=clr)
    add_text(s, x + Cm(0.3), Cm(4.4), Cm(5), Cm(0.5), role, font_size=8, color=GRAY)

# CS/영업
add_text(s, Cm(2), Cm(5.7), Cm(23), Cm(0.5), 'CS·영업 그룹', font_size=13, bold=True, color=RED)
cs = [
    ('🛠️ cs-agent', 'CS 기술지원', BLUE),
    ('💰 sales', '수주·견적·매출', GREEN),
    ('📅 schedule', '일정·마일스톤', ORANGE),
]
for i, (name, role, clr) in enumerate(cs):
    x = Cm(2) + i * Cm(6.2)
    add_rect(s, x, Cm(6.4), Cm(5.6), Cm(1.5), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.3), Cm(6.5), Cm(5), Cm(0.5), name, font_size=10, bold=True, color=clr)
    add_text(s, x + Cm(0.3), Cm(7.1), Cm(5), Cm(0.5), role, font_size=8, color=GRAY)

# 개발/설계
add_text(s, Cm(2), Cm(8.3), Cm(23), Cm(0.5), '개발·설계 그룹', font_size=13, bold=True, color=RED)
dev = [
    ('🔧 crafter', '시스템 구축·배포', BLUE),
    ('💻 dev', 'MES 프론트/백엔드', GREEN),
    ('📐 design', '기구설계·도면', ORANGE),
    ('⚡ control', '제어설계·모션', RED),
]
for i, (name, role, clr) in enumerate(dev):
    x = Cm(2) + i * Cm(6.2)
    add_rect(s, x, Cm(9), Cm(5.6), Cm(1.5), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.3), Cm(9.1), Cm(5), Cm(0.5), name, font_size=10, bold=True, color=clr)
    add_text(s, x + Cm(0.3), Cm(9.7), Cm(5), Cm(0.5), role, font_size=8, color=GRAY)

# 품질/데이터
add_text(s, Cm(2), Cm(11), Cm(23), Cm(0.5), '품질·데이터 그룹', font_size=13, bold=True, color=RED)
qa = [
    ('🔍 nc-manager', '부적합 관리', BLUE),
    ('🔬 qa', '출하 품질검사', GREEN),
    ('🚨 issue', '제품개선 트래킹', ORANGE),
    ('📊 db-manager', 'DB·ERP·벡터DB', RED),
]
for i, (name, role, clr) in enumerate(qa):
    x = Cm(2) + i * Cm(6.2)
    add_rect(s, x, Cm(11.7), Cm(5.6), Cm(1.5), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.3), Cm(11.8), Cm(5), Cm(0.5), name, font_size=10, bold=True, color=clr)
    add_text(s, x + Cm(0.3), Cm(12.4), Cm(5), Cm(0.5), role, font_size=8, color=GRAY)
slide_num(s, 1)


# ═══════════════════════════════════════════
# SLIDE 3: MES 시스템 현황
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, 'MES 시스템 현황')

nums = [
    ('119', RED, '구현 페이지'), ('73', BLUE, 'API 핸들러'),
    ('47', GREEN, 'CS 테이블'), ('15명', ORANGE, 'AI 에이전트'),
]
for i, (val, clr, lbl) in enumerate(nums):
    x = Cm(1.5) + i * Cm(6.2)
    add_rect(s, x, Cm(3), Cm(5.6), Cm(2.2), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x, Cm(3.1), Cm(5.6), Cm(1), val, font_size=32, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_text(s, x, Cm(4.3), Cm(5.6), Cm(0.5), lbl, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

# 기술 스택
add_rect(s, Cm(2), Cm(5.8), Cm(11), Cm(4), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(6), Cm(10), Cm(0.5), '기술 스택', font_size=12, bold=True, color=BLUE)
add_multiline(s, Cm(2.5), Cm(6.7), Cm(10), Cm(2.5), [
    '• 백엔드: Go 1.26 + Gin 1.12',
    '• 프론트: React 19 + TypeScript 5.9',
    '• DB: PostgreSQL (Supabase) + pgvector',
    '• ERP: SQL Server (읽기전용)',
], font_size=10)

# DB 구조
add_rect(s, Cm(14), Cm(5.8), Cm(11), Cm(4), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(14.5), Cm(6), Cm(10), Cm(0.5), 'AI 통합 데이터 구조', font_size=12, bold=True, color=GREEN)
add_multiline(s, Cm(14.5), Cm(6.7), Cm(10), Cm(2.5), [
    '• ai_agent_action_log — 에이전트 활동 기록',
    '• ai_agent_session — 세션 관리',
    '• vector_embeddings — RAG 벡터 저장',
    '• 이미 구축된 AI-MES 연동 인프라',
], font_size=10)

# 이중 DB 연동
add_rect(s, Cm(2), Cm(10.3), Cm(23), Cm(3.2), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(2.5), Cm(10.5), Cm(22), Cm(0.5), '이중 DB 연동 아키텍처', font_size=12, bold=True, color=ORANGE)
add_multiline(s, Cm(2.5), Cm(11.2), Cm(22), Cm(2), [
    '• PostgreSQL (MES 자체 데이터) + SQL Server (ERP 데이터) 실시간 조회',
    '• Go 백엔드에서 sqlx로 안전한 타입 검증 쿼리 — AI 에이전트가 자연어로 데이터 접근 가능',
    '• 벡터DB 9.3GB+ (CS이력 + 매뉴얼 + 기술문서) — 지식 검색 기반 구축 완료',
], font_size=10)
slide_num(s, 2)


# ═══════════════════════════════════════════
# SLIDE 4: AI 확장 가능 영역
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, 'AI 확장 가능 영역')

areas = [
    ('부적합관리 AI', [
        '• 텍스트 분석 → 원인 자동분류',
        '• 과거 패턴 기반 검사기준 제안',
        '• 재발 방지 자동 추천',
        '현재 데이터: 684건 부적합 보고서',
    ], RED),
    ('품질검사 AI', [
        '• 불량률 예측 모델',
        '• 시리즈별 품질 트렌드 분석',
        '• 체크리스트 자동 생성',
        '현재 데이터: QC 체크리스트 벡터DB',
    ], BLUE),
    ('일정관리 AI', [
        '• ML 납기지연 예측',
        '• 병목공정 자동 감지',
        '• 스케줄 자동 재조정',
        '현재 데이터: 57 프로젝트, 99 수주',
    ], GREEN),
    ('설비 모니터링 AI', [
        '• 센서 데이터 이상 예측',
        '• 서보/모터 파라미터 최적화',
        '• 예지보전 알림',
        '현재 데이터: hardware 스키마 15 테이블',
    ], ORANGE),
]
for i, (title, items, clr) in enumerate(areas):
    col = i % 2
    row = i // 2
    x = Cm(2) + col * Cm(12.5)
    y = Cm(3) + row * Cm(5.2)
    add_rect(s, x, y, Cm(11.5), Cm(4.8), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.5), y + Cm(0.2), Cm(10.5), Cm(0.5), title,
             font_size=13, bold=True, color=clr)
    add_multiline(s, x + Cm(0.5), y + Cm(0.9), Cm(10.5), Cm(3.5), items, font_size=10)
slide_num(s, 3)


# ═══════════════════════════════════════════
# SLIDE 5: 업무효율성 증대 효과
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '업무효율성 증대 효과')

eff_table = [
    ['업무 영역', '기존 방식', 'AI 적용', '효과'],
    ['CS 응답', '수 시간 ~ 1일', '수 분 이내', 'RAG 즉시 검색 응답'],
    ['부적합 등록', '수동 입력 5~10분', '자동 30초', '슬랙 대화 자동 파싱'],
    ['데이터 조회', 'DBA에게 요청 → 대기', '자연어 즉시 조회', 'db-manager 자동 처리'],
    ['문서 작성', '수일 소요', '수 시간', 'AI 초안 + 다국어 번역'],
    ['일정 관리', '수동 엑셀 추적', '자동 알림 + 예측', '지연 사전 감지'],
    ['품질 분석', '월간 수동 집계', '실시간 자동 분석', '트렌드 즉시 파악'],
]
add_table(s, eff_table, Cm(2), Cm(3), Cm(23), [Cm(4), Cm(5), Cm(5), Cm(9)], row_h=Cm(0.6))

# 비용 절감
add_rect(s, Cm(2), Cm(8.5), Cm(23), Cm(5), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(2.5), Cm(8.7), Cm(22), Cm(0.5), '비용 절감 효과', font_size=14, bold=True, color=GREEN)

cost_table = [
    ['항목', '현재 비용', '절감 방안', '절감액'],
    ['Jira 라이선스', '59명 유료', 'AI 팀 활용으로 축소', '상당'],
    ['ERP 유지보수', '연간 계약', 'MES 자체 기능 대체', '상당'],
    ['AI 라이선스', '개별 구독 분산', '통합 (Rovo + Claude Team)', '효율화'],
]
add_table(s, cost_table, Cm(2.5), Cm(9.5), Cm(22), [Cm(5), Cm(5), Cm(6.5), Cm(5.5)])

add_text(s, Cm(2.5), Cm(12.5), Cm(22), Cm(0.8),
         '연간 약 2,000만원 절감 예상 (Jira 해지 + ERP 유지보수 해지 + AI 라이선스 통합)',
         font_size=13, bold=True, color=RED)
slide_num(s, 4)


# ═══════════════════════════════════════════
# SLIDE 6: 슬랙 비전
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '슬랙 비전: 직원 대화 → 데이터')

add_text(s, Cm(2), Cm(3.2), Cm(23), Cm(0.8),
         '각 슬랙 채널에 AI 팀원을 배치하여, 직원의 자연스러운 대화가 구조화된 데이터로 자동 전환됩니다.',
         font_size=13, color=BODY)

# 흐름도
flow = [
    ('직원 대화\n(슬랙 채널)', BLUE),
    ('AI 자동 수집\n(라우팅+분류)', ORANGE),
    ('정형 데이터\n(구조화)', GREEN),
    ('MES DB\n(축적)', RED),
]
for i, (label, clr) in enumerate(flow):
    x = Cm(2) + i * Cm(6.2)
    add_rect(s, x, Cm(4.5), Cm(5.2), Cm(2), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.2), Cm(4.6), Cm(4.8), Cm(1.8), label,
             font_size=11, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        add_text(s, x + Cm(5.2), Cm(5), Cm(1), Cm(1), '→',
                 font_size=18, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)

# ERP 연동 확장
add_rect(s, Cm(2), Cm(7.5), Cm(23), Cm(6), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(7.7), Cm(22), Cm(0.5), 'ERP 연동 AI 확장 계획', font_size=13, bold=True, color=BLUE)

erp_items = [
    ('원가계산 자동화', '자재비+인건비+경비 자동 집계\nERP 원가 데이터 실시간 연동', RED),
    ('재고 수요 예측', '과거 출고 패턴 ML 분석\n적정 재고량 자동 산출', GREEN),
    ('공급망 리스크', '납기 지연 이력 분석\n대체 공급처 자동 추천', ORANGE),
]
for i, (title, desc, clr) in enumerate(erp_items):
    x = Cm(3) + i * Cm(7.5)
    add_rect(s, x, Cm(8.5), Cm(6.5), Cm(4.2), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, x + Cm(0.3), Cm(8.7), Cm(5.9), Cm(0.5), title,
             font_size=11, bold=True, color=clr)
    add_text(s, x + Cm(0.3), Cm(9.4), Cm(5.9), Cm(2.5), desc,
             font_size=9, color=BODY)
slide_num(s, 5)


# ═══════════════════════════════════════════
# SLIDE 7: 로드맵
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '로드맵')

# Phase 1 — 완료
add_rect(s, Cm(2), Cm(3), Cm(23), Cm(3), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(2.5), Cm(3.2), Cm(22), Cm(0.5), 'Phase 1 (현재) — ✅ 완료', font_size=14, bold=True, color=GREEN)
add_multiline(s, Cm(2.5), Cm(3.9), Cm(10), Cm(2), [
    '• CS AI 챗봇 구축 (RAG 9.3GB+)',
    '• 슬랙 채널별 AI 라우팅',
    '• MES 119 페이지 / 73 API',
], font_size=11)
add_multiline(s, Cm(14), Cm(3.9), Cm(10), Cm(2), [
    '• 15명 에이전트 체계 확립',
    '• 부적합 자동 등록',
    '• ERP 읽기 연동',
], font_size=11)

# Phase 2
add_rect(s, Cm(2), Cm(6.5), Cm(23), Cm(3), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(6.7), Cm(22), Cm(0.5), 'Phase 2 (Q2 2026) — 품질·일정 AI', font_size=14, bold=True, color=BLUE)
add_multiline(s, Cm(2.5), Cm(7.4), Cm(10), Cm(2), [
    '• 부적합 원인 자동분류',
    '• 불량률 예측 모델',
    '• 체크리스트 자동생성',
], font_size=11)
add_multiline(s, Cm(14), Cm(7.4), Cm(10), Cm(2), [
    '• ML 납기지연 예측',
    '• 병목공정 감지',
    '• 원가계산 자동화',
], font_size=11)

# Phase 3
add_rect(s, Cm(2), Cm(10), Cm(23), Cm(3.5), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(2.5), Cm(10.2), Cm(22), Cm(0.5), 'Phase 3 (Q3~Q4 2026) — 글로벌·예지보전', font_size=14, bold=True, color=ORANGE)
add_multiline(s, Cm(2.5), Cm(10.9), Cm(10), Cm(2), [
    '• 설비 예지보전 시스템',
    '• 센서 이상 탐지 AI',
    '• 서보/모터 파라미터 최적화',
], font_size=11)
add_multiline(s, Cm(14), Cm(10.9), Cm(10), Cm(2), [
    '• 다국어 글로벌 CS (7개국)',
    '• 전사 AI 통합 플랫폼',
    '• 경영진 AI 대시보드',
], font_size=11)
slide_num(s, 6)


# ═══════════════════════════════════════════
# SLIDE 8: 엔딩
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)


# ═══ 저장 ═══
prs.save(OUTPUT_PATH)
print(f'PPTX 저장 완료: {OUTPUT_PATH}')
print(f'슬라이드 수: {len(prs.slides)}')
