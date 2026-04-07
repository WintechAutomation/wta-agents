"""
2026년 기술경영계획 월별/부서별 PPT 분석 → HTML 리포트 생성 (텍스트 + 이미지)
"""
import sys, os, re, hashlib
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import pptx.exc

BASE = Path('C:/MES/wta-agents/data/plan')
OUTPUT = Path('C:/MES/wta-agents/reports/MAX/2026-monthly-plan-analysis.html')
IMAGES_DIR = Path('C:/MES/wta-agents/reports/MAX/2026-monthly-plan-images')
OUTPUT.parent.mkdir(parents=True, exist_ok=True)
IMAGES_DIR.mkdir(parents=True, exist_ok=True)

MONTHS = [
    ('2026-01', '기술경영계획자료 2026년 1월', '1월'),
    ('2026-02', '기술경영계획자료 2026년 2월', '2월'),
    ('2026-03', '기술경영계획자료 2026년 3월', '3월'),
]

# 부서 우선순위 매핑 (파일명 키워드 → 부서명)
DEPT_MAP = [
    ('영업',           '영업'),
    ('전략기획',        '전략기획'),
    ('품질CS',         '품질CS그룹'),
    ('품질',           '품질CS그룹'),
    ('생산관리',        '생산관리'),
    ('구매',           '구매'),
    ('HAM팀(Press)',   '기구설계HAM(Press)'),
    ('HAM팀',          '기구설계HAM'),
    ('HIM팀',          '기구설계HIM'),
    ('제작팀',          '제작'),
    ('제어팀',          '제어'),
    ('소프트웨어팀',     '소프트웨어'),
    ('소프트웨어',       '소프트웨어'),
    ('비전팀',          '비전'),
    ('광학기술센터',     '광학기술센터'),
    ('연삭팀',          '연삭'),
    ('관리팀',          '관리'),
]

def get_dept(filename):
    for key, dept in DEPT_MAP:
        if key in filename:
            return dept
    return filename.split('_')[-1].replace('.pptx','').replace('.ppt','')

def safe_dir_name(s):
    return re.sub(r'[<>:"/\\|?*\s]', '_', s)

def extract_images_from_pptx(path, month_key, dept, max_images=15):
    """PPT에서 이미지 추출, 저장 후 상대 경로 목록 반환 (최대 max_images개)"""
    dept_safe = safe_dir_name(dept)
    out_dir = IMAGES_DIR / month_key / dept_safe
    out_dir.mkdir(parents=True, exist_ok=True)
    saved = []
    try:
        prs = Presentation(str(path))
    except Exception:
        return saved

    img_count = 0
    seen_hashes = set()
    for slide_idx, slide in enumerate(prs.slides):
        if img_count >= max_images:
            break
        try:
            for shape in slide.shapes:
                if img_count >= max_images:
                    break
                # PICTURE 타입(13)만 추출
                if shape.shape_type != 13:
                    continue
                try:
                    img = shape.image
                    blob = img.blob
                    # 중복 스킵
                    img_hash = hashlib.md5(blob[:512]).hexdigest()
                    if img_hash in seen_hashes:
                        continue
                    seen_hashes.add(img_hash)
                    # 너무 작은 이미지 스킵 (아이콘, 로고 등 2KB 미만)
                    if len(blob) < 2048:
                        continue
                    ext = (img.ext or 'png').lower()
                    if ext not in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff'):
                        ext = 'png'
                    fname = f'slide{slide_idx+1:02d}_{img_count+1:02d}.{ext}'
                    fpath = out_dir / fname
                    fpath.write_bytes(blob)
                    rel = f'2026-monthly-plan-images/{month_key}/{dept_safe}/{fname}'
                    saved.append({'path': rel, 'slide': slide_idx+1})
                    img_count += 1
                except Exception:
                    continue
        except Exception:
            continue
    return saved


def extract_text_from_pptx(path):
    """PPT 파일에서 슬라이드별 텍스트 추출"""
    try:
        prs = Presentation(str(path))
    except Exception as e:
        print(f'  [오류] {path.name}: {e}')
        return []

    slides_text = []
    for i, slide in enumerate(prs.slides):
        try:
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t and len(t) > 1:
                            slide_lines.append(t)
            if slide_lines:
                slides_text.append({'slide': i+1, 'lines': slide_lines})
        except Exception:
            continue
    return slides_text

def extract_key_info(slides_text, dept):
    """슬라이드 텍스트에서 핵심 과제/이슈 추출"""
    all_lines = []
    for s in slides_text:
        all_lines.extend(s['lines'])

    # 제목 슬라이드 제외 키워드
    SKIP_WORDS = ['EXIT', 'NEXT', 'CONTENTS', '작성일자', '작성자', '참조)', 'INDEX', 'PAGE']

    # 과제/주제 관련 패턴
    TASK_PATTERNS = ['과제', '목표', '계획', '개발', '개선', '구축', '추진', '적용', '검토',
                     '완료', '진행', '이슈', '현안', '문제', '해결', '분석']

    tasks = []
    issues = []
    kpis = []
    content_slides = []

    # CONTENTS 슬라이드에서 목차 추출
    for s in slides_text[:3]:
        for line in s['lines']:
            if any(w in line for w in SKIP_WORDS):
                continue
            if re.search(r'\d+\.?\s+.{5,}', line) and len(line) < 60:
                content_slides.append(line)

    # 전체 텍스트에서 과제/이슈 추출
    for line in all_lines:
        if any(w in line for w in SKIP_WORDS):
            continue
        if len(line) < 4 or len(line) > 150:
            continue

        # KPI (숫자+단위 포함)
        if re.search(r'\d+[억만%건대개호]', line) and len(line) < 80:
            kpis.append(line)
        # 이슈/문제
        elif any(w in line for w in ['이슈', '문제', '지연', '불량', '오류', '알람', '실패', '미달', '부족', '차질']):
            issues.append(line)
        # 과제/계획
        elif any(w in line for w in TASK_PATTERNS) and len(line) < 100:
            tasks.append(line)

    # 목차 항목 우선, 없으면 과제 목록
    main_tasks = content_slides[:8] if content_slides else tasks[:8]

    return {
        'tasks': main_tasks,
        'issues': issues[:5],
        'kpis': kpis[:5],
    }

# ---- 데이터 수집 ----
report_data = {}  # {month_key: {dept: {tasks, issues, kpis}}}

for month_key, folder_name, month_label in MONTHS:
    folder = BASE / folder_name
    if not folder.exists():
        print(f'[SKIP] {folder_name} 없음')
        continue

    print(f'\n=== {month_label} ===')
    report_data[month_key] = {'label': month_label, 'depts': {}}

    ppt_files = sorted(f for f in folder.iterdir() if f.suffix.lower() in ('.pptx', '.ppt'))
    # 크기 제한 (300MB 초과 파일 스킵 — 연삭팀 등 초대형)
    for f in ppt_files:
        size_mb = f.stat().st_size / (1024*1024)
        if size_mb > 200:
            dept = get_dept(f.name)
            print(f'  [스킵 {size_mb:.0f}MB] {f.name}')
            report_data[month_key]['depts'][dept] = {
                'tasks': [f'※ 파일 용량 초과({size_mb:.0f}MB)로 자동 추출 불가 — 직접 확인 필요'],
                'issues': [], 'kpis': [], 'filename': f.name
            }
            continue

        dept = get_dept(f.name)
        print(f'  [{size_mb:.0f}MB] {f.name} → {dept}')
        slides = extract_text_from_pptx(f)
        info = extract_key_info(slides, dept)
        info['filename'] = f.name
        info['slides'] = len(slides)
        # 이미지 추출
        images = extract_images_from_pptx(f, month_key, dept, max_images=15)
        info['images'] = images
        print(f'    이미지: {len(images)}개 추출')
        report_data[month_key]['depts'][dept] = info

print('\n데이터 수집 완료. HTML 생성 중...')

# ---- HTML 생성 ----
DEPT_ORDER = ['영업', '전략기획', '품질CS그룹', '생산관리', '구매',
              '기구설계HAM(Press)', '기구설계HAM', '기구설계HIM',
              '제작', '제어', '소프트웨어', '비전', '광학기술센터', '연삭', '관리']

DEPT_COLORS = {
    '영업': '#4472C4', '전략기획': '#5B9BD5', '품질CS그룹': '#ED7D31',
    '생산관리': '#A5A5A5', '구매': '#FFC000', '기구설계HAM(Press)': '#70AD47',
    '기구설계HAM': '#70AD47', '기구설계HIM': '#255E91', '제작': '#9E480E',
    '제어': '#843C0C', '소프트웨어': '#264478', '비전': '#806000',
    '광학기술센터': '#375623', '연삭': '#7030A0', '관리': '#808080',
}

def dept_card(dept, info, month):
    color = DEPT_COLORS.get(dept, '#4472C4')
    tasks_html = ''.join(f'<li>{t}</li>' for t in info.get('tasks', []) if t)
    issues_html = ''.join(f'<li class="issue">{i}</li>' for i in info.get('issues', []) if i)
    kpis_html = ''.join(f'<li class="kpi">{k}</li>' for k in info.get('kpis', []) if k)

    # 이미지 HTML (썸네일 그리드)
    images = info.get('images', [])
    img_html = ''
    if images:
        thumb_parts = []
        for img in images:
            p = img['path']
            s = img['slide']
            thumb_parts.append(f'<a href="{p}" target="_blank"><img src="{p}" alt="슬라이드{s}" title="슬라이드 {s}" loading="lazy"></a>')
        thumbs = ''.join(thumb_parts)
        img_html = f'<div class="section"><div class="section-title">🖼 슬라이드 이미지 ({len(images)}개)</div><div class="img-grid">{thumbs}</div></div>'

    return f'''
    <div class="dept-card">
      <div class="dept-header" style="background:{color}">
        <span class="dept-name">{dept}</span>
        <span class="dept-file">{info.get("filename","")[:40]}</span>
      </div>
      <div class="dept-body">
        <div class="section">
          <div class="section-title">📋 주요 과제 / 목차</div>
          <ul class="task-list">{tasks_html or "<li class='empty'>추출된 과제 없음</li>"}</ul>
        </div>
        {"" if not info.get("issues") else f'<div class="section"><div class="section-title">⚠️ 이슈 / 현안</div><ul class="task-list">{issues_html}</ul></div>'}
        {"" if not info.get("kpis") else f'<div class="section"><div class="section-title">📊 KPI / 성과지표</div><ul class="task-list">{kpis_html}</ul></div>'}
        {img_html}
      </div>
    </div>'''

month_sections = ''
for month_key, folder_name, month_label in MONTHS:
    if month_key not in report_data:
        continue
    data = report_data[month_key]
    depts_html = ''
    # 정해진 순서대로 출력, 나머지는 뒤에
    seen = set()
    for dept in DEPT_ORDER:
        if dept in data['depts']:
            depts_html += dept_card(dept, data['depts'][dept], month_key)
            seen.add(dept)
    for dept, info in data['depts'].items():
        if dept not in seen:
            depts_html += dept_card(dept, info, month_key)

    dept_count = len(data['depts'])
    month_sections += f'''
  <section class="month-section" id="{month_key}">
    <div class="month-header">
      <h2>2026년 {data["label"]}</h2>
      <span class="badge">{dept_count}개 팀/부서</span>
    </div>
    <div class="dept-grid">
      {depts_html}
    </div>
  </section>'''

html = f'''<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>2026년 기술경영계획 월별/부서별 과제 현황 — (주)윈텍오토메이션</title>
<style>
* {{ box-sizing: border-box; margin: 0; padding: 0; }}
body {{ font-family: "맑은 고딕", "Malgun Gothic", sans-serif; background: #f0f4f8; color: #1a1a2e; }}

/* 헤더 */
.top-header {{ background: #4472C4; color: white; padding: 24px 32px; }}
.top-header h1 {{ font-size: 1.6rem; font-weight: 700; margin-bottom: 6px; }}
.top-header p {{ font-size: 0.9rem; opacity: 0.85; }}

/* 네비게이션 */
.nav-bar {{ background: white; border-bottom: 2px solid #4472C4; padding: 12px 32px; display: flex; gap: 12px; position: sticky; top: 0; z-index: 100; box-shadow: 0 2px 8px rgba(0,0,0,0.08); }}
.nav-btn {{ background: none; border: 2px solid #4472C4; color: #4472C4; padding: 6px 18px; border-radius: 20px; cursor: pointer; font-family: "맑은 고딕", sans-serif; font-size: 0.9rem; font-weight: 600; text-decoration: none; }}
.nav-btn:hover {{ background: #4472C4; color: white; }}

/* 월별 섹션 */
.month-section {{ padding: 32px; max-width: 1600px; margin: 0 auto; }}
.month-header {{ display: flex; align-items: center; gap: 14px; margin-bottom: 20px; padding-bottom: 12px; border-bottom: 3px solid #4472C4; }}
.month-header h2 {{ font-size: 1.4rem; color: #4472C4; font-weight: 700; }}
.badge {{ background: #4472C4; color: white; padding: 3px 12px; border-radius: 12px; font-size: 0.82rem; }}

/* 부서 그리드 */
.dept-grid {{ display: grid; grid-template-columns: repeat(auto-fill, minmax(340px, 1fr)); gap: 16px; }}

/* 부서 카드 */
.dept-card {{ background: white; border-radius: 10px; overflow: hidden; box-shadow: 0 2px 8px rgba(0,0,0,0.08); border: 1px solid #e8eaf0; }}
.dept-header {{ padding: 12px 16px; display: flex; flex-direction: column; gap: 3px; }}
.dept-name {{ font-size: 1.05rem; font-weight: 700; color: white; }}
.dept-file {{ font-size: 0.72rem; color: rgba(255,255,255,0.8); }}
.dept-body {{ padding: 14px 16px; }}
.section {{ margin-bottom: 12px; }}
.section:last-child {{ margin-bottom: 0; }}
.section-title {{ font-size: 0.8rem; font-weight: 700; color: #555; margin-bottom: 6px; text-transform: uppercase; letter-spacing: 0.3px; }}

/* 리스트 */
.task-list {{ padding-left: 16px; font-size: 0.85rem; line-height: 1.7; color: #333; }}
.task-list li {{ margin-bottom: 2px; }}
.task-list li.issue {{ color: #c0392b; }}
.task-list li.kpi {{ color: #1a7a3e; font-weight: 600; }}
.task-list li.empty {{ color: #aaa; list-style: none; margin-left: -16px; font-style: italic; }}

/* 구분선 */
.month-divider {{ border: none; border-top: 1px solid #dde3ee; margin: 0 32px; }}

/* 이미지 그리드 */
.img-grid {{ display: flex; flex-wrap: wrap; gap: 6px; margin-top: 6px; }}
.img-grid a {{ display: block; }}
.img-grid img {{ width: 120px; height: 80px; object-fit: cover; border: 1px solid #ddd; border-radius: 4px; cursor: pointer; transition: transform 0.2s; }}
.img-grid img:hover {{ transform: scale(1.05); box-shadow: 0 4px 12px rgba(0,0,0,0.2); }}

/* 반응형 */
@media (max-width: 768px) {{
  .dept-grid {{ grid-template-columns: 1fr; }}
  .month-section {{ padding: 16px; }}
  .top-header h1 {{ font-size: 1.2rem; }}
}}
</style>
</head>
<body>

<div class="top-header">
  <h1>2026년 기술경영계획 월별/부서별 과제 현황</h1>
  <p>(주)윈텍오토메이션 · 분석 기준: 2026년 1월~3월 · python-pptx 텍스트 자동 추출</p>
</div>

<div class="nav-bar">
  <a class="nav-btn" href="#2026-01">1월</a>
  <a class="nav-btn" href="#2026-02">2월</a>
  <a class="nav-btn" href="#2026-03">3월</a>
</div>

{month_sections}

<div style="text-align:center; padding: 24px; color:#888; font-size:0.8rem;">
  생성: 2026-04-05 · db-manager 자동 분석 · 원본 PPT에서 텍스트 자동 추출 (이미지/그래프 내 텍스트 미포함)
</div>
</body>
</html>'''

OUTPUT.write_text(html, encoding='utf-8')
print(f'\n✅ HTML 저장 완료: {OUTPUT}')
print(f'   파일 크기: {OUTPUT.stat().st_size//1024}KB')
