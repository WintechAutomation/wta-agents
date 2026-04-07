"""slack-proposal.html 5슬라이드 -> PPTX 생성.
Template.pptx 배경 이미지 사용, WTA 공식 템플릿."""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Template.pptx 사용 (배경 이미지 포함)
TEMPLATE_PATH = 'C:/MES/wta-agents/data/uploads/fa85f7cd-7ce6-4fe8-9a13-c9fc4d77b56d/Template.pptx'
OUTPUT_PATH = 'C:/MES/frontend/public/slack-proposal.pptx'

prs = Presentation(TEMPLATE_PATH)

# 색상
RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x15, 0x65, 0xC0)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
DARK = RGBColor(0x22, 0x22, 0x22)
BODY = RGBColor(0x44, 0x44, 0x44)
GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = '맑은 고딕'

# 레이아웃 인덱스 확인
def find_layout(name_hint):
    for layout in prs.slide_layouts:
        if name_hint in layout.name:
            return layout
    return prs.slide_layouts[0]

cover_layout = find_layout('제목')
content_layout = find_layout('제목 및 내용')
blank_layout = find_layout('빈')

# ── 헬퍼 ──
def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=BODY, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.font.bold = bold
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_table(slide, rows_data, left, top, width, col_widths, row_h=Cm(0.55)):
    nrows, ncols = len(rows_data), len(rows_data[0])
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, row_h * nrows)
    tbl = shape.table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x44, 0x54, 0x6A)
                else:
                    p.font.color.rgb = BODY
    return shape

def slide_title(slide, text):
    add_text(slide, Cm(5), Cm(1), Cm(20), Cm(1.5), text,
             font_size=28, bold=True, color=RED)

def slide_num(slide, num):
    add_text(slide, Cm(25), Cm(14.5), Cm(2), Cm(0.5), f'{num:02d}',
             font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

# 기존 슬라이드 삭제 (템플릿에 포함된 샘플)
while len(prs.slides) > 0:
    sld_id = prs.slides._sldIdLst[0]
    rId = sld_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') or sld_id.get('r:id')
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sld_id)


# ═══════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════
s = prs.slides.add_slide(cover_layout)
add_text(s, Cm(2), Cm(4), Cm(18), Cm(3),
         '업무 메신저 개선 제안\nSlack 도입',
         font_size=36, bold=True, color=DARK)
add_text(s, Cm(2), Cm(7.5), Cm(18), Cm(1),
         '채널 기반 소통 · AI 연동 · 협업 효율 향상',
         font_size=18, color=GRAY)
add_text(s, Cm(2), Cm(9.5), Cm(18), Cm(1.5),
         '(주)윈텍오토메이션 생산관리팀 (AI운영팀)\n2026년 4월',
         font_size=13, color=GRAY)


# ═══════════════════════════════════════════
# SLIDE 2: 현황 및 도입 배경
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '현황 및 도입 배경')

# 현재 환경 박스
add_rect(s, Cm(2), Cm(3), Cm(23), Cm(1.5), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(3.2), Cm(22), Cm(1),
         '현재 환경: 다우오피스 그룹웨어 메신저를 사내 업무 메신저로 사용 중',
         font_size=13, bold=True, color=BLUE)

# 현재 한계점 (좌)
add_rect(s, Cm(2), Cm(5), Cm(11), Cm(6.5), RGBColor(0xFF, 0xF5, 0xF5), RED)
add_text(s, Cm(2.5), Cm(5.2), Cm(10), Cm(0.6), '현재 한계점', font_size=14, bold=True, color=RED)
add_multiline(s, Cm(2.5), Cm(6), Cm(10), Cm(5), [
    '▸ 채널 분류 불가 — 주제별 대화 구분 없이',
    '   1:1/그룹 채팅만 가능',
    '▸ 외부 연동 미지원 — API, 봇, 웹훅 등',
    '   자동화 연결 불가',
    '▸ 검색 기능 제한 — 과거 대화 탐색 및',
    '   파일 추적이 어려움',
], font_size=11, color=RED)

# Slack 도입 목적 (우)
add_rect(s, Cm(14), Cm(5), Cm(11), Cm(6.5), RGBColor(0xF0, 0xFF, 0xF0), GREEN)
add_text(s, Cm(14.5), Cm(5.2), Cm(10), Cm(0.6), 'Slack 도입 목적', font_size=14, bold=True, color=GREEN)
add_multiline(s, Cm(14.5), Cm(6), Cm(10), Cm(5), [
    '▸ 업무 채널화 — 부서·주제별 채널로',
    '   정보 구조화',
    '▸ AI 연동 기반 확보 — 봇·자동 알림·',
    '   워크플로 즉시 구현',
    '▸ 협업 효율 향상 — 검색, 스레드,',
    '   앱 연동으로 생산성 극대화',
], font_size=11, color=GREEN)
slide_num(s, 1)


# ═══════════════════════════════════════════
# SLIDE 3: 비용 비교 및 Slack 무료 플랜
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '비용 비교 및 Slack 무료 플랜')

cost_data = [
    ['항목', 'Slack Free (무료)', 'Slack Pro (참고)'],
    ['비용', '무료', '월 $7.25/인'],
    ['메시지 보존', '최근 90일', '무제한'],
    ['채널', '무제한', '무제한'],
    ['앱/봇 연동', '최대 10개', '무제한'],
    ['통화', '1:1 음성/영상', '그룹 허들'],
    ['파일 저장', '제한 없음 (90일)', '제한 없음'],
]
add_table(s, cost_data, Cm(2), Cm(3.2), Cm(23), [Cm(5), Cm(9), Cm(9)])

# 핵심 포인트 (좌)
add_rect(s, Cm(2), Cm(8.5), Cm(11), Cm(4), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(2.5), Cm(8.7), Cm(10), Cm(0.6), '핵심 포인트', font_size=13, bold=True, color=GREEN)
add_multiline(s, Cm(2.5), Cm(9.5), Cm(10), Cm(2.5), [
    '추가 비용 0원 — 무료 플랜으로',
    '채널 기반 소통 + AI 봇 연동이 모두 가능',
    '다우오피스 대비 추가 비용 없이 기능 대폭 확장',
], font_size=11, color=BODY)

# 업그레이드 시점 (우)
add_rect(s, Cm(14), Cm(8.5), Cm(11), Cm(4), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(14.5), Cm(8.7), Cm(10), Cm(0.6), '업그레이드 시점', font_size=13, bold=True, color=ORANGE)
add_multiline(s, Cm(14.5), Cm(9.5), Cm(10), Cm(2.5), [
    '90일 이상 메시지 보관 또는',
    '10개 초과 앱 연동이 필요해지면 Pro 검토',
    '현 단계에서는 무료 플랜으로 충분',
], font_size=11, color=BODY)
slide_num(s, 2)


# ═══════════════════════════════════════════
# SLIDE 4: 달라지는 점 & 기대 효과
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '달라지는 점 & 기대 효과')

# 기존 (좌)
add_rect(s, Cm(2), Cm(3), Cm(11), Cm(5), RGBColor(0xFF, 0xF5, 0xF5), RED)
add_text(s, Cm(2.5), Cm(3.2), Cm(10), Cm(0.6), '기존 (다우오피스)', font_size=13, bold=True, color=RED)
add_multiline(s, Cm(2.5), Cm(4), Cm(10), Cm(3.5), [
    '▸ 1:1 / 그룹 채팅 위주',
    '▸ 파일 첨부 후 추적 어려움',
    '▸ 키워드 검색 한계',
    '▸ 외부 시스템 연동 불가',
], font_size=12, color=RED)

# 변경 후 (우)
add_rect(s, Cm(14), Cm(3), Cm(11), Cm(5), RGBColor(0xF0, 0xFF, 0xF0), GREEN)
add_text(s, Cm(14.5), Cm(3.2), Cm(10), Cm(0.6), '변경 후 (Slack)', font_size=13, bold=True, color=GREEN)
add_multiline(s, Cm(14.5), Cm(4), Cm(10), Cm(3.5), [
    '▸ 채널 기반 주제별 소통',
    '▸ 스레드로 맥락 유지',
    '▸ 강력한 검색 + 필터링',
    '▸ API/봇/웹훅 자유 연동',
], font_size=12, color=GREEN)

# 기대 효과 3개 박스
effects = [
    ('정보 체계화', '부서별·주제별 채널로\n대화가 자동 분류되어\n정보 유실 방지', RED),
    ('AI 봇 연동', '자동 알림, CS 자동응답,\n품질 이슈 알림 등\nAI 에이전트 즉시 배치', GREEN),
    ('외부 서비스 연동', 'Jira, MES, 캘린더 등\n업무 도구와\n실시간 연결', BLUE),
]
for i, (title, desc, clr) in enumerate(effects):
    x = Cm(2) + i * Cm(8.3)
    bg = RGBColor(0xFF, 0xEB, 0xEE) if clr == RED else RGBColor(0xE8, 0xF5, 0xE9) if clr == GREEN else RGBColor(0xE3, 0xF2, 0xFD)
    add_rect(s, x, Cm(8.5), Cm(7.5), Cm(3.5), bg, clr)
    add_text(s, x + Cm(0.5), Cm(8.7), Cm(6.5), Cm(0.6), title, font_size=12, bold=True, color=clr)
    add_multiline(s, x + Cm(0.5), Cm(9.5), Cm(6.5), Cm(2.3), desc.split('\n'), font_size=10, color=BODY)

# 전환 방식 안내
add_rect(s, Cm(2), Cm(12.5), Cm(23), Cm(1.2), RGBColor(0xFF, 0xF8, 0xE1), ORANGE)
add_text(s, Cm(2.5), Cm(12.6), Cm(22), Cm(0.8),
         '전환 방식: 기존 다우오피스 메신저와 병행 운영 후 단계적 전환 (업무 중단 없음)',
         font_size=12, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
slide_num(s, 3)


# ═══════════════════════════════════════════
# SLIDE 5: 엔딩
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
# 엔딩 슬라이드 — 배경 이미지만


# ═══ 저장 ═══
prs.save(OUTPUT_PATH)
print(f'PPTX 저장 완료: {OUTPUT_PATH}')
print(f'슬라이드 수: {len(prs.slides)}')
