"""system-optimization.html 4슬라이드 → PPTX 생성.
Template.pptx 의존 제거, 동일 슬라이드 크기(28.22x15.88cm) 사용."""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 빈 프레젠테이션 (Template.pptx 의존 제거)
prs = Presentation()
# Template.pptx와 동일한 슬라이드 크기
prs.slide_width = 10160000   # 28.22 cm
prs.slide_height = 5715000   # 15.88 cm

# 색상 (gen_pptx.py 동일)
NAVY = RGBColor(0x44, 0x54, 0x6A)
BLUE = RGBColor(0x44, 0x72, 0xC4)
LBLUE = RGBColor(0x5B, 0x9B, 0xD5)
GREEN = RGBColor(0x70, 0xAD, 0x47)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
RED = RGBColor(0xD4, 0x20, 0x27)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0x99)
DARK = RGBColor(0x33, 0x33, 0x33)
CARD_BG = RGBColor(0xF5, 0xF7, 0xFC)
FONT = '맑은 고딕'

# ── 헬퍼 (gen_pptx.py 동일 시그니처) ──

def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_bar(slide, left, top, width, height, fill_color):
    """직사각형 바."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_table(slide, rows_data, left, top, width, col_widths, row_h=Cm(0.45)):
    """테이블 생성 — 첫 행 헤더(NAVY)."""
    nrows, ncols = len(rows_data), len(rows_data[0])
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, row_h * nrows)
    tbl = shape.table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = NAVY
                else:
                    p.font.color.rgb = DARK
    return shape

def add_footer(slide, text):
    add_text(slide, Cm(1), Cm(13.8), Cm(20), Cm(0.5), text, font_size=9, color=GRAY)


# Blank layout (빈 프레젠테이션 기본)
layout = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════
# 슬라이드 1: 현행 월 IT 시스템 비용
# (HTML 슬라이드1: AI구독, IT서비스, 보안, 인프라 — 2x2)
# ═══════════════════════════════════════════════════
s1 = prs.slides.add_slide(layout)
add_text(s1, Cm(4), Cm(0.8), Cm(20), Cm(1.5),
         '현행 월 IT 시스템 비용', font_size=32, bold=True, color=DARK)
add_text(s1, Cm(4), Cm(2.1), Cm(20), Cm(0.8),
         '현재 운영 중인 시스템별 월 지출 현황', font_size=16, color=GRAY)

# 2x2 그리드 (카드 크기를 슬라이드에 맞게)
cw = Cm(12.2)   # 카드 너비
ch = Cm(4.6)    # 카드 높이
gap = Cm(0.4)
sx = Cm(1.2)
sy = Cm(3.2)

# AI 구독 (좌상)
add_rect(s1, sx, sy, cw, ch, CARD_BG)
add_bar(s1, sx, sy, Cm(0.2), ch, PURPLE)
add_text(s1, sx + Cm(0.5), sy + Cm(0.15), cw, Cm(0.5),
         'AI 구독', font_size=12, bold=True, color=PURPLE)
ai_tbl = [
    ['서비스', '수량', '월 비용'],
    ['OpenAI', '13개', '$260 (약 36만원)'],
    ['Gemini', '18개', '$432 (약 60만원)'],
    ['Claude Code MAX', '2개', '$400 (약 55만원)'],
    ['Claude Team', '6개', '$180 (약 25만원)'],
    ['AI 소계', '', '$1,272/월 (약 175만원)'],
]
add_table(s1, ai_tbl, sx + Cm(0.4), sy + Cm(0.7), cw - Cm(0.6),
          [Cm(4.2), Cm(2.2), Cm(5)], Cm(0.38))

# IT 서비스 (우상)
x2 = sx + cw + gap
add_rect(s1, x2, sy, cw, ch, CARD_BG)
add_bar(s1, x2, sy, Cm(0.2), ch, BLUE)
add_text(s1, x2 + Cm(0.5), sy + Cm(0.15), cw, Cm(0.5),
         'IT 서비스', font_size=12, bold=True, color=BLUE)
it_tbl = [
    ['서비스', '수량', '월 비용'],
    ['Jira', '84명', '라이선스 기준'],
    ['Confluence + Rovo', '전 직원', '현행 유지'],
    ['다우오피스 그룹웨어', '전사', '현행 유지'],
    ['ERP 유지보수', '1건', '50만원'],
]
add_table(s1, it_tbl, x2 + Cm(0.4), sy + Cm(0.7), cw - Cm(0.6),
          [Cm(4.2), Cm(2.2), Cm(5)], Cm(0.38))

# 보안 (좌하)
y2 = sy + ch + gap
add_rect(s1, sx, y2, cw, ch, CARD_BG)
add_bar(s1, sx, y2, Cm(0.2), ch, ORANGE)
add_text(s1, sx + Cm(0.5), y2 + Cm(0.15), cw, Cm(0.5),
         '보안', font_size=12, bold=True, color=ORANGE)
sec_tbl = [
    ['항목', '수량', '월 비용'],
    ['DLP + 백신 + DRM + SW자산관리', '85EA', '84만원'],
    ['UTM (통합 위협 관리)', '1EA', '16만원'],
    ['NAC (Genian)', '', '26.4만원'],
    ['보안 소계', '', '126.4만원'],
]
add_table(s1, sec_tbl, sx + Cm(0.4), y2 + Cm(0.7), cw - Cm(0.6),
          [Cm(5), Cm(2), Cm(4.4)], Cm(0.38))

# 인프라 (우하)
add_rect(s1, x2, y2, cw, ch, CARD_BG)
add_bar(s1, x2, y2, Cm(0.2), ch, GREEN)
add_text(s1, x2 + Cm(0.5), y2 + Cm(0.15), cw, Cm(0.5),
         '인프라', font_size=12, bold=True, color=GREEN)
infra_tbl = [
    ['항목', '월 비용'],
    ['AWS (EC2+Storage+SES)', '~20만원 (예상)'],
    ['Cloudflare / SSL', '무료'],
    ['Supabase / NAS 1·2차', '무료'],
]
add_table(s1, infra_tbl, x2 + Cm(0.4), y2 + Cm(0.7), cw - Cm(0.6),
          [Cm(6.5), Cm(5)], Cm(0.38))

# 총합 바 (gen_pptx.py 스타일)
total_y = Cm(12.6)
add_rect(s1, sx, total_y, Cm(24), Cm(0.9), NAVY)
add_text(s1, Cm(2), total_y + Cm(0.1), Cm(12), Cm(0.6),
         '월 IT 비용 합계', font_size=13, bold=True, color=WHITE)
add_text(s1, Cm(14), total_y + Cm(0.05), Cm(10), Cm(0.8),
         '약 400만원+/월', font_size=20, bold=True, color=GOLD, alignment=PP_ALIGN.RIGHT)

add_footer(s1, 'WTA · IT 시스템 최적화')


# ═══════════════════════════════════════════════════
# 슬라이드 2: 최적화 방안
# (HTML 슬라이드2: 3가지 핵심 변경)
# ═══════════════════════════════════════════════════
s2 = prs.slides.add_slide(layout)
add_text(s2, Cm(4), Cm(0.8), Cm(20), Cm(1.5),
         '최적화 방안', font_size=32, bold=True, color=DARK)
add_text(s2, Cm(4), Cm(2.1), Cm(20), Cm(0.8),
         '비용 절감을 위한 3가지 핵심 변경', font_size=16, color=GRAY)

# gen_pptx.py의 카드 레이아웃 재사용
cards = [
    {'label': '절감 ①', 'name': 'Jira 라이선스 축소',
     'detail': '84명 → 24명 (60명 해지)\n유지: SW설계·광학·비전\n해지: 제어·기구·구매·품질 등\n대체: MES 이슈트래커 (이미 운영 중) · Confluence 현행 유지',
     'amount': '780', 'color': BLUE},
    {'label': '절감 ②', 'name': 'ERP 유지보수 해지',
     'detail': '월 50만원 → 0원\n현재 개발 의뢰 없이 유지보수비만 지출\n대체: DB 자체관리 + MES 고도화',
     'amount': '600', 'color': LBLUE},
    {'label': '절감 ③', 'name': 'AI 구독 변화',
     'detail': 'OpenAI 유료 라이선스 축소 (13개 → 해지)\nGemini, Claude 전환\nAtlassian Jira/Confluence(Rovo) 활용',
     'amount': '430', 'color': GREEN},
]

card_w = Cm(7.8)
card_h = Cm(6.5)
card_y = Cm(3.2)
card_gap = Cm(0.6)
card_sx = Cm(1.2)

for i, c in enumerate(cards):
    x = card_sx + i * (card_w + card_gap)
    add_rect(s2, x, card_y, card_w, card_h, CARD_BG)
    add_bar(s2, x, card_y, Cm(0.25), card_h, c['color'])
    add_text(s2, x + Cm(0.6), card_y + Cm(0.3), card_w - Cm(1), Cm(0.5),
             c['label'], font_size=10, color=GRAY, bold=True)
    add_text(s2, x + Cm(0.6), card_y + Cm(0.8), card_w - Cm(1), Cm(0.7),
             c['name'], font_size=16, bold=True, color=DARK)
    add_text(s2, x + Cm(0.6), card_y + Cm(1.7), card_w - Cm(1), Cm(3),
             c['detail'], font_size=10, color=GRAY)
    add_text(s2, x + Cm(0.6), card_y + Cm(5), card_w - Cm(1), Cm(1.2),
             c['amount'] + '만원/년', font_size=28, bold=True, color=c['color'])

# 총합 바
total_y = Cm(10.2)
add_rect(s2, card_sx, total_y, Cm(24), Cm(3.2), NAVY)
add_text(s2, Cm(2), total_y + Cm(0.4), Cm(14), Cm(0.8),
         '3개 항목 합산 연간 총 절감 예상액', font_size=16, bold=True, color=WHITE)
add_text(s2, Cm(2), total_y + Cm(1.3), Cm(14), Cm(0.6),
         '780 + 600 + 430 = 1,810만원', font_size=12, color=RGBColor(0xAA, 0xAA, 0xBB))
add_text(s2, Cm(16), total_y + Cm(0.2), Cm(8), Cm(2),
         '1,810만원', font_size=40, bold=True, color=GOLD, alignment=PP_ALIGN.RIGHT)
add_text(s2, Cm(16), total_y + Cm(2.2), Cm(8), Cm(0.6),
         '월 환산 약 151만원 절감', font_size=11,
         color=RGBColor(0xAA, 0xAA, 0xBB), alignment=PP_ALIGN.RIGHT)

add_footer(s2, 'WTA · IT 시스템 최적화')


# ═══════════════════════════════════════════════════
# 슬라이드 3: 부서별 협의 필요사항
# (HTML 슬라이드3: 3개 카드)
# ═══════════════════════════════════════════════════
s3 = prs.slides.add_slide(layout)
add_text(s3, Cm(4), Cm(0.8), Cm(20), Cm(1.5),
         '부서별 협의 필요사항', font_size=32, bold=True, color=DARK)
add_text(s3, Cm(4), Cm(2.1), Cm(20), Cm(0.8),
         '변경 전 확인·합의가 필요한 항목', font_size=16, color=GRAY)

consult = [
    {
        'dept': '개발부서 (SW·비전·광학)',
        'title': 'Jira 라이선스 유지 범위 확정',
        'items': '▸ 현재 24명 유지 기준 (SW설계·광학·비전)\n▸ MES 이슈트래커로 대체 가능 여부 확인\n▸ Confluence는 전 직원 유지 — 영향 없음\n▸ 다우오피스 기준 인원 vs Jira 계정 동기화',
        'status': '협의 진행 중', 'color': BLUE,
    },
    {
        'dept': '경영지원 / 회계',
        'title': 'ERP 유지보수 계약 해지',
        'items': '▸ 현재 월 50만원 유지보수 계약 상태\n▸ 최근 개발 의뢰 이력 확인\n▸ DB 백업·자체관리 체계 확인\n▸ admin-agent 자동화 백업 완료 후 해지',
        'status': '협의 필요', 'color': ORANGE,
    },
    {
        'dept': '전 부서',
        'title': 'AI 구독 전환',
        'items': '▸ OpenAI 사용 부서별 현황 파악\n▸ Gemini + Claude 전환 워크플로우 검증\n▸ Rovo AI 전사 활용 교육 계획\n▸ 전환 기간 중 병행 운영 방안',
        'status': '즉시 가능', 'color': GREEN,
    },
]

cc_w = Cm(7.8)
cc_h = Cm(9.5)
cc_y = Cm(3.2)
cc_gap = Cm(0.6)

for i, c in enumerate(consult):
    cx = Cm(1.2) + i * (cc_w + cc_gap)
    add_rect(s3, cx, cc_y, cc_w, cc_h, CARD_BG)
    # 상단 컬러 바
    add_bar(s3, cx, cc_y, cc_w, Cm(0.2), c['color'])
    # 부서명
    add_text(s3, cx + Cm(0.5), cc_y + Cm(0.4), cc_w - Cm(1), Cm(0.5),
             c['dept'], font_size=9, bold=True, color=GRAY)
    # 제목
    add_text(s3, cx + Cm(0.5), cc_y + Cm(1), cc_w - Cm(1), Cm(0.7),
             c['title'], font_size=14, bold=True, color=DARK)
    # 항목
    add_text(s3, cx + Cm(0.5), cc_y + Cm(2), cc_w - Cm(1), Cm(5.5),
             c['items'], font_size=10, color=RGBColor(0x66, 0x66, 0x66))
    # 상태 배지
    add_text(s3, cx + Cm(0.5), cc_y + Cm(8.2), cc_w - Cm(1), Cm(0.6),
             c['status'], font_size=10, bold=True, color=c['color'])

add_footer(s3, 'WTA · IT 시스템 최적화')


# ═══════════════════════════════════════════════════
# 슬라이드 4: 변경 후 비용 절감 효과
# (HTML 슬라이드4: 현행 vs 변경 후 비교)
# ═══════════════════════════════════════════════════
s4 = prs.slides.add_slide(layout)
add_text(s4, Cm(4), Cm(0.8), Cm(20), Cm(1.5),
         '변경 후 비용 절감 효과', font_size=32, bold=True, color=DARK)
add_text(s4, Cm(4), Cm(2.1), Cm(20), Cm(0.8),
         '현행 대비 변경 후 월 비용 비교', font_size=16, color=GRAY)

# 현행 (좌) — gen_pptx.py의 좌표 체계 사용
col_w = Cm(10)
ey = Cm(3.2)

add_rect(s4, Cm(1.2), ey, col_w, Cm(7.5), CARD_BG)
add_bar(s4, Cm(1.2), ey, col_w, Cm(0.2), NAVY)
add_text(s4, Cm(1.7), ey + Cm(0.4), col_w - Cm(1), Cm(0.7),
         '현행 (월)', font_size=14, bold=True, color=NAVY)

before = [
    ('Jira 84명', '65만원'),
    ('ERP 유지보수', '50만원'),
    ('AI 구독 (4종)', '175만원'),
    ('보안 (5종)', '126.4만원'),
    ('인프라 (AWS 등)', '~20만원'),
]
for j, (name, cost) in enumerate(before):
    ry = ey + Cm(1.3) + j * Cm(1.1)
    add_text(s4, Cm(1.7), ry, Cm(5), Cm(0.5), name, font_size=11, color=DARK)
    add_text(s4, Cm(7), ry, Cm(3.8), Cm(0.5), cost, font_size=12,
             bold=True, color=NAVY, alignment=PP_ALIGN.RIGHT)

# 화살표
add_text(s4, Cm(11.5), ey + Cm(2.5), Cm(2.5), Cm(2),
         '→', font_size=44, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

# 변경 후 (우)
ax = Cm(14.2)
add_rect(s4, ax, ey, col_w, Cm(7.5), CARD_BG)
add_bar(s4, ax, ey, col_w, Cm(0.2), BLUE)
add_text(s4, ax + Cm(0.5), ey + Cm(0.4), col_w - Cm(1), Cm(0.7),
         '변경 후 (월)', font_size=14, bold=True, color=BLUE)

after = [
    ('Jira 24명', '20만원', BLUE, False),
    ('ERP 유지보수', '0원', GRAY, True),
    ('AI 구독 (3종)', '140만원', BLUE, False),
    ('보안 (5종)', '126.4만원', BLUE, False),
    ('인프라 (AWS 등)', '~20만원', BLUE, False),
]
for j, (name, cost, clr, strike) in enumerate(after):
    ry = ey + Cm(1.3) + j * Cm(1.1)
    nc = RGBColor(0xBB, 0xBB, 0xBB) if strike else DARK
    cc = RGBColor(0xBB, 0xBB, 0xBB) if strike else clr
    add_text(s4, ax + Cm(0.5), ry, Cm(5), Cm(0.5), name, font_size=11, color=nc)
    add_text(s4, ax + Cm(5.5), ry, Cm(4), Cm(0.5), cost, font_size=12,
             bold=True, color=cc, alignment=PP_ALIGN.RIGHT)

# 총 절감 효과 바
save_y = Cm(11.5)
add_rect(s4, Cm(1.2), save_y, Cm(24), Cm(2), NAVY)
add_text(s4, Cm(2), save_y + Cm(0.3), Cm(14), Cm(0.6),
         '연간 총 절감 예상', font_size=14, bold=True, color=WHITE)
add_text(s4, Cm(2), save_y + Cm(1), Cm(14), Cm(0.5),
         'Jira 780 + ERP 600 + AI 430', font_size=10, color=RGBColor(0xAA, 0xAA, 0xBB))
add_text(s4, Cm(16), save_y + Cm(0.1), Cm(8), Cm(1.8),
         '약 1,810만원', font_size=36, bold=True, color=GOLD, alignment=PP_ALIGN.RIGHT)

add_footer(s4, 'WTA · IT 시스템 최적화 · 2026-03-30')


# ═══ 저장 ═══
output_path = 'C:/MES/frontend/public/system-optimization.pptx'
prs.save(output_path)
print(f'PPT 저장 완료: {output_path}')
print(f'슬라이드 수: {len(prs.slides)}')
