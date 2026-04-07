"""security-slide.html 13슬라이드 → PPTX 생성.
Template.pptx 배경 이미지 사용, WTA 공식 템플릿."""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Template.pptx 사용 (배경 이미지 포함)
TEMPLATE_PATH = 'C:/MES/wta-agents/data/uploads/fa85f7cd-7ce6-4fe8-9a13-c9fc4d77b56d/Template.pptx'
OUTPUT_PATH = 'C:/MES/frontend/public/security-slide.pptx'

prs = Presentation(TEMPLATE_PATH)

# 색상
RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x15, 0x65, 0xC0)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
DARK = RGBColor(0x22, 0x22, 0x22)
BODY = RGBColor(0x44, 0x44, 0x44)
GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = '맑은 고딕'

# 레이아웃 인덱스 확인
def find_layout(name_hint):
    for i, layout in enumerate(prs.slide_layouts):
        if name_hint in layout.name:
            return layout
    return prs.slide_layouts[0]

cover_layout = find_layout('제목')
content_layout = find_layout('제목 및 내용')
blank_layout = find_layout('빈')

# ── 헬퍼 ──
def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=BODY, bold=False):
    """여러 줄 텍스트 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.font.bold = bold
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_table(slide, rows_data, left, top, width, col_widths, row_h=Cm(0.55)):
    nrows, ncols = len(rows_data), len(rows_data[0])
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, row_h * nrows)
    tbl = shape.table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x44, 0x54, 0x6A)
                else:
                    p.font.color.rgb = BODY
    return shape

def slide_title(slide, text):
    """내용 슬라이드 제목 (빨간색, 왼쪽 상단)"""
    add_text(slide, Cm(5), Cm(1), Cm(20), Cm(1.5), text,
             font_size=28, bold=True, color=RED)

def slide_num(slide, num):
    add_text(slide, Cm(25), Cm(14.5), Cm(2), Cm(0.5), f'{num:02d}',
             font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

# 기존 슬라이드 삭제 (템플릿에 포함된 샘플)
while len(prs.slides) > 0:
    sld_id = prs.slides._sldIdLst[0]
    rId = sld_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') or sld_id.get('r:id')
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sld_id)


# ═══════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════
s = prs.slides.add_slide(cover_layout)
add_text(s, Cm(2), Cm(4), Cm(18), Cm(3),
         'AI 멀티에이전트 시스템\n데이터 보안 보고',
         font_size=36, bold=True, color=DARK)
add_text(s, Cm(2), Cm(7.5), Cm(18), Cm(1),
         '기술 데이터 보호 체계 및 외부 유출 방지 대책',
         font_size=18, color=GRAY)
add_text(s, Cm(2), Cm(9.5), Cm(18), Cm(1.5),
         '(주)윈텍오토메이션 생산관리팀 (AI운영팀)\n2026년 4월 2일',
         font_size=13, color=GRAY)


# ═══════════════════════════════════════════
# SLIDE 2: 핵심 요약
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '핵심 요약')
add_text(s, Cm(2), Cm(3.2), Cm(23), Cm(1),
         '"AI를 쓰면 우리 기술 데이터가 밖으로 새어나가는 것 아닌가?"',
         font_size=18, bold=True, color=DARK)

# 2개 박스
add_rect(s, Cm(2), Cm(5), Cm(11), Cm(3), RGBColor(0xF0, 0xFF, 0xF0), GREEN)
add_text(s, Cm(2.5), Cm(5.2), Cm(10), Cm(0.6), '✅ AI가 우리 데이터를 학습하지 않습니다', font_size=14, bold=True, color=GREEN)
add_text(s, Cm(2.5), Cm(6), Cm(10), Cm(1.5), '유료 API는 계약상 고객 데이터를 AI 모델\n훈련에 사용하지 않습니다. 법적 보장.', font_size=11, color=BODY)

add_rect(s, Cm(14), Cm(5), Cm(11), Cm(3), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(14.5), Cm(5.2), Cm(10), Cm(0.6), '🔐 모든 전송은 암호화됩니다', font_size=14, bold=True, color=BLUE)
add_text(s, Cm(14.5), Cm(6), Cm(10), Cm(1.5), '인터넷 뱅킹과 동일한 수준의\nTLS 1.3 암호화로 보호됩니다.', font_size=11, color=BODY)

# 3개 큰 숫자
nums = [('0건', GREEN, 'AI 학습에 사용된 WTA 데이터'), ('30일', BLUE, '이후 전송 기록 자동 삭제'), ('100%', RED, '사내 서버에서 운영')]
for i, (val, clr, lbl) in enumerate(nums):
    x = Cm(2) + i * Cm(8.3)
    add_text(s, x, Cm(9), Cm(7), Cm(2), val, font_size=40, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_text(s, x, Cm(11.2), Cm(7), Cm(0.5), lbl, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)
slide_num(s, 1)


# ═══════════════════════════════════════════
# SLIDE 3: 동작 방식
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, 'AI는 어떻게 동작하는가')
add_multiline(s, Cm(2), Cm(3.2), Cm(23), Cm(2), [
    'AI 에이전트는 우리 회사 서버 안에서 동작합니다.',
    '외부 AI 서비스(Claude)는 "두뇌" 역할만 합니다 — 질문을 보내면 답변을 돌려줍니다.',
], font_size=14)

add_rect(s, Cm(2), Cm(6), Cm(23), Cm(6.5), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_text(s, Cm(2.5), Cm(6.2), Cm(22), Cm(0.6), '💡 비유하자면', font_size=14, bold=True, color=BLUE)
add_multiline(s, Cm(2.5), Cm(7), Cm(22), Cm(5), [
    'AI 에이전트 = 우리 회사에 상주하는 직원',
    'Claude API = 외부 전문가에게 전화로 자문을 구하는 것',
    '',
    '직원이 전문가에게 전화할 때 필요한 내용만 말하고,',
    '전문가는 답변 후 통화 내용을 30일 뒤 자동으로 폐기합니다.',
    '전문가가 우리 통화 내용으로 다른 고객을 가르치지 않습니다.',
], font_size=12)
slide_num(s, 2)


# ═══════════════════════════════════════════
# SLIDE 4: 시스템 구성도
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, 'WTA AI 시스템 전체 구조')

# 사내 서버 영역
add_rect(s, Cm(2), Cm(3.5), Cm(23), Cm(6), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(2.5), Cm(3.7), Cm(22), Cm(0.5), '🏢 WTA 사내 서버 (외부 접근 차단)', font_size=10, bold=True, color=GREEN)

agents = ['👑 MAX', '📊 DB매니저', '🛠️ CS담당', '💻 개발팀', '🔧 크래프터',
          '🔍 NC관리', '🔬 출하품질', '💰 영업팀', '📐 설계팀', '📅 일정관리']
for i, name in enumerate(agents):
    x = Cm(3) + (i % 5) * Cm(4.2)
    y = Cm(4.5) + (i // 5) * Cm(1.2)
    add_rect(s, x, y, Cm(3.8), Cm(0.9), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
    add_text(s, x + Cm(0.2), y + Cm(0.1), Cm(3.4), Cm(0.7), name, font_size=9, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)

dbs = ['🗄️ MES DB', '🗄️ ERP DB (읽기전용)', '📁 기술문서']
for i, db in enumerate(dbs):
    x = Cm(4) + i * Cm(7)
    add_rect(s, x, Cm(7.2), Cm(6), Cm(0.9), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
    add_text(s, x + Cm(0.2), Cm(7.3), Cm(5.6), Cm(0.7), db, font_size=9, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

add_text(s, Cm(2), Cm(8.3), Cm(23), Cm(0.5), '🔒 모든 데이터는 사내에 저장 — 외부 반출 없음',
         font_size=10, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# 화살표
add_text(s, Cm(8), Cm(9.5), Cm(11), Cm(0.8), '⬆️⬇️ 암호화된 질문/답변만 오감 (TLS 1.3)',
         font_size=11, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

# 외부 서비스
add_rect(s, Cm(2), Cm(10.8), Cm(23), Cm(2.5), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(2.5), Cm(11), Cm(22), Cm(0.5), '☁️ 외부 클라우드 서비스', font_size=10, bold=True, color=ORANGE)
ext = ['🧠 Claude AI (미국)', '💬 Slack', '📋 Jira']
for i, name in enumerate(ext):
    x = Cm(4) + i * Cm(7)
    add_rect(s, x, Cm(11.8), Cm(6), Cm(0.9), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
    add_text(s, x + Cm(0.2), Cm(11.9), Cm(5.6), Cm(0.7), name, font_size=9, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
slide_num(s, 3)


# ═══════════════════════════════════════════
# SLIDE 5: 데이터 흐름
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '어떤 데이터가 외부로 나가는가?')
data_flow = [
    ['전송 경로', '전송 데이터', '암호화', 'AI 학습', '보존'],
    ['→ Claude AI', '에이전트 질문/분석 텍스트', 'TLS 1.3', '사용 안함', '30일 삭제'],
    ['→ Slack', '채널 대화 메시지', 'TLS 1.2+', '사용 안함', '슬랙 정책'],
    ['→ Jira', '이슈/문서 내용', 'TLS 1.2+', '사용 안함', '직접 관리'],
]
add_table(s, data_flow, Cm(2), Cm(3.5), Cm(23), [Cm(4), Cm(7), Cm(3), Cm(4), Cm(5)])

add_rect(s, Cm(2), Cm(8), Cm(23), Cm(4.5), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text(s, Cm(2.5), Cm(8.2), Cm(22), Cm(0.5), '📌 중요한 점', font_size=14, bold=True, color=GREEN)
add_multiline(s, Cm(2.5), Cm(9), Cm(22), Cm(3), [
    'DB 원본 데이터, 도면 파일, CAD 파일은 외부로 전송되지 않습니다.',
    '에이전트가 분석에 필요한 텍스트 요약본만 Claude에 질문으로 보내고, 답변을 받아옵니다.',
    '원본 데이터는 항상 사내 서버에만 존재합니다.',
], font_size=12, bold=True)
slide_num(s, 4)


# ═══════════════════════════════════════════
# SLIDE 6: 핵심 Q&A
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '핵심 질문과 답변')
qas = [
    ('Q1. AI가 우리 기술을 배워서 경쟁사에 알려주지 않나?', 'NO — 유료 API는 고객 데이터를 모델 학습에 사용하지 않음. 법적 의무.', GREEN),
    ('Q2. 전송 중에 해커가 데이터를 훔칠 수 있나?', 'NO — 인터넷 뱅킹과 동일한 TLS 1.3 암호화 적용. 제3자 열람 불가능.', GREEN),
    ('Q3. Anthropic 직원이 우리 데이터를 볼 수 있나?', '일상 업무에서 열람 안 함. 악용 조사 목적에 한해 제한적 접근만 가능.', BLUE),
    ('Q4. 만약 Anthropic이 해킹당하면?', 'SOC 2 Type II + AES-256 + 30일 자동 삭제로 노출 범위 매우 제한적.', BLUE),
]
for i, (q, a, clr) in enumerate(qas):
    y = Cm(3.5) + i * Cm(2.7)
    add_rect(s, Cm(2), y, Cm(23), Cm(2.3), RGBColor(0xF8, 0xF8, 0xF8), clr)
    add_text(s, Cm(2.5), y + Cm(0.2), Cm(22), Cm(0.6), q, font_size=12, bold=True, color=DARK)
    add_text(s, Cm(2.5), y + Cm(1), Cm(22), Cm(1), a, font_size=11, color=BODY)
slide_num(s, 5)


# ═══════════════════════════════════════════
# SLIDE 7: 무료 vs 유료
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '왜 유료 API를 사용하는가?')
add_text(s, Cm(2), Cm(3.2), Cm(23), Cm(0.8),
         '같은 AI라도 무료 웹 채팅과 유료 API는 데이터 취급 방식이 완전히 다릅니다.',
         font_size=14, color=BODY)

# 무료 (좌)
add_rect(s, Cm(2), Cm(4.5), Cm(11), Cm(6), RGBColor(0xFF, 0xF5, 0xF5), RED)
add_text(s, Cm(2.5), Cm(4.7), Cm(10), Cm(0.6), '✗ 무료 AI 채팅 (일반 사용)', font_size=13, bold=True, color=RED)
add_multiline(s, Cm(2.5), Cm(5.5), Cm(10), Cm(4.5), [
    '✗ 대화 내용이 AI 학습에 사용될 수 있음',
    '✗ 데이터 삭제 절차가 복잡함',
    '✗ 기업 보안 계약(DPA) 불가',
    '✗ 감사 로그 없음',
    '✗ 개인 단위 — 통제 불가',
], font_size=11, color=RED)

# 유료 (우)
add_rect(s, Cm(14), Cm(4.5), Cm(11), Cm(6), RGBColor(0xF0, 0xFF, 0xF0), GREEN)
add_text(s, Cm(14.5), Cm(4.7), Cm(10), Cm(0.6), '✓ 유료 API (WTA가 사용 중)', font_size=13, bold=True, color=GREEN)
add_multiline(s, Cm(14.5), Cm(5.5), Cm(10), Cm(4.5), [
    '✓ 학습에 절대 사용되지 않음 (법적 보장)',
    '✓ 30일 후 자동 삭제 + 즉시 삭제 가능',
    '✓ 기업 보안 계약(DPA) 체결 가능',
    '✓ API 사용 감사 로그 제공',
    '✓ 회사 관리 — 접근 통제 가능',
], font_size=11, color=GREEN)

# 경고
add_rect(s, Cm(2), Cm(11), Cm(23), Cm(2), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text(s, Cm(2.5), Cm(11.2), Cm(22), Cm(0.5), '⚠️ 오히려 위험한 것은', font_size=12, bold=True, color=ORANGE)
add_text(s, Cm(2.5), Cm(11.8), Cm(22), Cm(1), '직원이 개인 무료 AI에 기술 데이터를 직접 입력하는 것이 훨씬 큰 리스크. 회사 시스템으로 통제 가능.',
         font_size=11, color=BODY)
slide_num(s, 6)


# ═══════════════════════════════════════════
# SLIDE 8: Anthropic 보안 수준
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, 'Anthropic(Claude)은 얼마나 안전한가?')
certs = [('🏅', 'SOC 2 Type II', '정보보안 감사 최고 인증'),
         ('🔐', 'AES-256', '군사 등급 암호화'),
         ('🗑️', '30일 자동 삭제', 'API 기록 영구 삭제'),
         ('📄', 'DPA 체결 가능', '법적 보호 계약')]
for i, (icon, title, desc) in enumerate(certs):
    x = Cm(2) + i * Cm(6.2)
    add_rect(s, x, Cm(3.5), Cm(5.6), Cm(3), RGBColor(0xF8, 0xF8, 0xF8))
    add_text(s, x, Cm(3.7), Cm(5.6), Cm(1), icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Cm(0.3), Cm(4.8), Cm(5), Cm(0.5), title, font_size=12, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Cm(0.3), Cm(5.5), Cm(5), Cm(0.5), desc, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

comp = [
    ['항목', 'Claude (WTA 사용 중)', 'GPT (OpenAI)', 'Gemini (Google)'],
    ['API 학습 제외', '✓ 기본 적용', '✓ 기본 적용', '✓ 유료만'],
    ['데이터 보존', '30일', '30일', '18개월*'],
    ['SOC 2 / DPA', '✓ Type II / ✓', '✓ / ✓', '✓ / ✓'],
]
add_table(s, comp, Cm(2), Cm(7.5), Cm(23), [Cm(5), Cm(7), Cm(5.5), Cm(5.5)])
slide_num(s, 7)


# ═══════════════════════════════════════════
# SLIDE 9: 현재 보호 조치
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, 'WTA가 이미 적용한 보안 장치들')
measures = [
    ('✅ 유료 API만 사용', '학습 제외 보장되는 유료 API만 사용'),
    ('✅ 100% 사내 서버 운영', '모든 에이전트 사내 서버, 외부 호스팅 없음'),
    ('✅ DB 외부 차단', 'MES·ERP DB 사내망 전용, 인터넷 접근 불가'),
    ('✅ ERP 읽기 전용', 'AI가 ERP 데이터 변경 불가'),
    ('✅ 슬랙 대화 사내 백업', '채널 대화를 사내 서버에 별도 저장'),
    ('✅ 업무별 권한 분리', '각 에이전트는 자기 업무 범위만 접근'),
]
for i, (title, desc) in enumerate(measures):
    col = i % 2
    row = i // 2
    x = Cm(2) + col * Cm(12.5)
    y = Cm(3.5) + row * Cm(3)
    add_rect(s, x, y, Cm(11.5), Cm(2.5), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
    add_text(s, x + Cm(0.5), y + Cm(0.3), Cm(10.5), Cm(0.6), title, font_size=12, bold=True, color=GREEN)
    add_text(s, x + Cm(0.5), y + Cm(1.2), Cm(10.5), Cm(1), desc, font_size=11, color=BODY)
slide_num(s, 8)


# ═══════════════════════════════════════════
# SLIDE 10: 강화 계획
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '추가 보안 강화 계획')
add_text(s, Cm(2), Cm(3.3), Cm(10), Cm(0.5), '단기 (2주 내)', font_size=14, bold=True, color=RED)
short = [
    ['#', '항목', '내용'],
    ['1', '민감 데이터 자동 필터', '도면번호·고객코드 AI 전송 전 자동 마스킹'],
    ['2', '전송 감사 로그', 'AI 전송 데이터 전체 기록'],
    ['3', '기밀 문서 차단', 'CAD·도면·BOM 자동 차단'],
]
add_table(s, short, Cm(2), Cm(4), Cm(23), [Cm(2), Cm(7), Cm(14)])

add_text(s, Cm(2), Cm(7.5), Cm(10), Cm(0.5), '중장기 로드맵', font_size=14, bold=True, color=RED)
road = [
    ['시기', '항목'],
    ['Q2', 'DPA 정식 체결 — Anthropic 법적 데이터 보호 계약'],
    ['Q2~Q3', '사내 전용 AI 검토 — 최고 기밀은 사내 AI로만 처리'],
    ['Q3', '분기별 보안 점검 — 정기 자동 점검 체계'],
    ['Q4', 'ISO 27001 준비 — 국제 정보보안 인증'],
]
add_table(s, road, Cm(2), Cm(8.2), Cm(23), [Cm(4), Cm(19)])
slide_num(s, 9)


# ═══════════════════════════════════════════
# SLIDE 11: 리스크 분석
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '남아 있는 위험과 대비')
risks = [
    ['우려 사항', '가능성', '대비 방안'],
    ['Anthropic 서버 해킹으로 API 기록 유출', '극히 낮음', 'SOC2, AES-256, 30일 자동 삭제. DPA 법적 보호.'],
    ['직원 실수로 기밀 자료 AI에 입력', '중간', '민감 데이터 자동 필터링. 직원 교육 + 가이드라인.'],
    ['사내 서버 외부 공격', '낮음', 'DB 외부 차단, 방화벽. 에이전트 간 인증 추가 예정.'],
    ['Anthropic 정책 변경', '극히 낮음', 'DPA 체결 시 법적 구속력. 사전 통지 의무.'],
]
add_table(s, risks, Cm(2), Cm(3.5), Cm(23), [Cm(8), Cm(3), Cm(12)])
slide_num(s, 10)


# ═══════════════════════════════════════════
# SLIDE 12: 종합 결론
# ═══════════════════════════════════════════
s = prs.slides.add_slide(content_layout)
slide_title(s, '종합 결론')

conclusions = [
    ('1. 기술 데이터가 AI 학습에 사용되지 않습니다.', '유료 API 계약상 법적 보장', GREEN),
    ('2. 데이터 전송은 은행 수준으로 보호됩니다.', 'TLS 1.3 + 30일 자동 삭제', BLUE),
    ('3. 직원의 개인 AI 사용이 더 큰 위험입니다.', '회사 시스템 통제가 더 안전', RED),
]
for i, (title, desc, clr) in enumerate(conclusions):
    y = Cm(3.5) + i * Cm(2)
    add_rect(s, Cm(2), y, Cm(23), Cm(1.6), RGBColor(0xE8, 0xF5, 0xE9) if clr == GREEN else RGBColor(0xE3, 0xF2, 0xFD) if clr == BLUE else RGBColor(0xFF, 0xEB, 0xEE))
    add_text(s, Cm(2.5), y + Cm(0.2), Cm(22), Cm(0.6), title, font_size=14, bold=True, color=clr)
    add_text(s, Cm(2.5), y + Cm(0.9), Cm(22), Cm(0.5), desc, font_size=11, color=GRAY)

# 결정 사항 박스
add_rect(s, Cm(2), Cm(10), Cm(11), Cm(3), RGBColor(0xFF, 0xEB, 0xEE), RED)
add_text(s, Cm(2.5), Cm(10.2), Cm(10), Cm(0.5), '📄 DPA 체결 추진', font_size=12, bold=True, color=RED)
add_text(s, Cm(2.5), Cm(10.9), Cm(10), Cm(1.5), 'Anthropic 법적 데이터 보호 계약\n(무료, API 구독 포함)', font_size=11, color=BODY)

add_rect(s, Cm(14), Cm(10), Cm(11), Cm(3), RGBColor(0xFF, 0xEB, 0xEE), RED)
add_text(s, Cm(14.5), Cm(10.2), Cm(10), Cm(0.5), '🖥️ 사내 전용 AI 검토', font_size=12, bold=True, color=RED)
add_text(s, Cm(14.5), Cm(10.9), Cm(10), Cm(1.5), '최고 기밀은 사내 AI로 처리\n(서버 투자 필요)', font_size=11, color=BODY)
slide_num(s, 11)


# ═══════════════════════════════════════════
# SLIDE 13: 엔딩
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
# 엔딩 슬라이드 — 배경 이미지만


# ═══ 저장 ═══
prs.save(OUTPUT_PATH)
print(f'PPTX 저장 완료: {OUTPUT_PATH}')
print(f'슬라이드 수: {len(prs.slides)}')
