"""
WTA AI 에이전트 시스템 운영 현황 — slide-test.pptx 생성
slide-test.html과 동일한 내용으로 PPTX 생성
출력: C:/MES/frontend/public/slide-test.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ASSETS = "C:/MES/wta-agents/config/ci/assets"
OUT = "C:/MES/frontend/public/slide-test.pptx"

# ── 색상 상수 ──
C_DARK   = RGBColor(0x22, 0x22, 0x22)
C_NAVY   = RGBColor(0x44, 0x54, 0x6A)
C_RED    = RGBColor(0xD4, 0x20, 0x27)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x88, 0x88, 0x88)
C_BLUE   = RGBColor(0x44, 0x72, 0xC4)
C_LBLUE  = RGBColor(0x5B, 0x9B, 0xD5)
C_ORANGE = RGBColor(0xED, 0x7D, 0x31)
C_GREEN  = RGBColor(0x70, 0xAD, 0x47)
C_LGRAY  = RGBColor(0xF5, 0xF7, 0xFC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def blank_layout(prs):
    return prs.slide_layouts[6]  # 완전 빈 레이아웃

def add_bg(slide, img_name):
    """배경 이미지 전체 슬라이드"""
    path = os.path.join(ASSETS, img_name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, 0, 0, W, H)

def add_footer(slide, text, slide_num, total):
    """하단 네이비 푸터 바"""
    fh = Inches(0.28)
    fy = H - fh
    # 푸터 배경
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, 0, fy, W, fh)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_NAVY
    shape.line.fill.background()
    # 푸터 텍스트
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    # 슬라이드 번호 (우측)
    nb = slide.shapes.add_textbox(W - Inches(1.2), fy, Inches(1.1), fh)
    nf = nb.text_frame
    np_ = nf.paragraphs[0]
    np_.alignment = PP_ALIGN.RIGHT
    nr = np_.add_run()
    nr.text = f"{slide_num} / {total}"
    nr.font.size = Pt(9)
    nr.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

def add_header(slide, title, subtitle=None, is_title_slide=False):
    """슬라이드 헤더 (제목 + 부제목)"""
    top = Inches(1.0) if is_title_slide else Inches(0.25)
    tb = slide.shapes.add_textbox(Inches(0.6), top, W - Inches(1.2), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.color.rgb = C_DARK
    run.font.size = Pt(32) if is_title_slide else Pt(24)

    if subtitle:
        from pptx.oxml.ns import qn
        from lxml import etree
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14) if is_title_slide else Pt(12)
        r2.font.color.rgb = C_RED if is_title_slide else C_GRAY

TOTAL = 8
FOOTER_TEXT = "WTA · AI 에이전트 시스템 운영 현황"

# ══════════════════════════════════════════════════
# 슬라이드 1: 표지
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
add_bg(sl, "bg-title.jpeg")
add_header(sl, "WTA AI 에이전트 시스템\n운영 현황 보고",
           "2026년 3월 · MAX 오케스트레이터 팀 · 비공개", is_title_slide=True)
add_footer(sl, FOOTER_TEXT, 1, TOTAL)

# ══════════════════════════════════════════════════
# 슬라이드 2: 목차
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
add_bg(sl, "bg-slide-header.jpeg")
add_header(sl, "목차", "Table of Contents")

toc_items = [
    ("1", "시스템 개요", "에이전트 구성 및 역할 분담", "3"),
    ("2", "운영 성과", "2026년 1분기 주요 지표 및 처리 현황", "4"),
    ("3", "도입 방안 비교", "AI 라이선스 통합 옵션 검토", "5"),
    ("4", "검토 의견", "팀별 장단점 분석", "6"),
    ("5", "향후 계획", "Phase 3 확장 로드맵", "7"),
    ("6", "결론 및 권고사항", "", "8"),
]

y = Inches(1.3)
row_h = Inches(0.72)
for num, title, sub, page in toc_items:
    # 번호 원
    circ = sl.shapes.add_shape(9, Inches(0.5), y + Inches(0.08), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = C_RED
    circ.line.fill.background()
    ctf = circ.text_frame
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = num
    cr.font.bold = True
    cr.font.color.rgb = C_WHITE
    cr.font.size = Pt(13)
    # 제목
    ttb = sl.shapes.add_textbox(Inches(1.15), y, Inches(10.5), row_h)
    ttf = ttb.text_frame
    ttf.word_wrap = False
    tp = ttf.paragraphs[0]
    tr = tp.add_run()
    tr.text = title + (f"  {sub}" if sub else "")
    tr.font.size = Pt(14)
    tr.font.bold = True
    tr.font.color.rgb = C_DARK
    # 페이지
    ptb = sl.shapes.add_textbox(W - Inches(1.0), y + Inches(0.1), Inches(0.8), Inches(0.4))
    ptf = ptb.text_frame
    pp_ = ptf.paragraphs[0]
    pp_.alignment = PP_ALIGN.RIGHT
    pr = pp_.add_run()
    pr.text = page
    pr.font.size = Pt(13)
    pr.font.bold = True
    pr.font.color.rgb = C_NAVY
    y += row_h

add_footer(sl, FOOTER_TEXT, 2, TOTAL)

# ══════════════════════════════════════════════════
# 슬라이드 3: 시스템 개요 (2단 레이아웃)
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
add_bg(sl, "bg-slide-header.jpeg")
add_header(sl, "시스템 개요", "에이전트 구성 및 역할 분담")

# 왼쪽 패널
lbox = sl.shapes.add_shape(1, Inches(0.4), Inches(1.4), Inches(7.8), Inches(5.5))
lbox.fill.solid(); lbox.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFC)
lbox.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
ltb = sl.shapes.add_textbox(Inches(0.55), Inches(1.5), Inches(7.5), Inches(5.2))
ltf = ltb.text_frame; ltf.word_wrap = True
lh = ltf.paragraphs[0]; lhr = lh.add_run()
lhr.text = "운영 중 에이전트 (6명)"; lhr.font.bold = True; lhr.font.size = Pt(13); lhr.font.color.rgb = C_NAVY
agents = [
    "MAX (오케스트레이터) — 요청 분류 · 팀원 위임 · 보고 검토",
    "crafter — MES 백엔드(Go/Gin) · API · 인프라",
    "db-manager — DB 조회 · ERP 연동 · 벡터DB 관리",
    "nc-manager — 부적합 이력 · 품질 분석 · NC 보고",
    "qa-agent — 출하검사 · 체크리스트 · QA 검증",
    "cs-agent — 대외 고객 응대 · CS 매뉴얼 RAG",
    "issue-manager — Jira 티켓 · 납기지연 · 크리티컬 이슈",
]
for ag in agents:
    p = ltf.add_paragraph(); p.space_before = Pt(4)
    r = p.add_run(); r.text = "• " + ag; r.font.size = Pt(11); r.font.color.rgb = C_DARK

# 오른쪽 패널
rbox = sl.shapes.add_shape(1, Inches(8.5), Inches(1.4), Inches(4.4), Inches(5.5))
rbox.fill.solid(); rbox.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFC)
rbox.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
rtb = sl.shapes.add_textbox(Inches(8.65), Inches(1.5), Inches(4.1), Inches(5.2))
rtf = rtb.text_frame; rtf.word_wrap = True
rh = rtf.paragraphs[0]; rhr = rh.add_run()
rhr.text = "기술 스택"; rhr.font.bold = True; rhr.font.size = Pt(13); rhr.font.color.rgb = C_NAVY
stack = [
    "Claude Code CLI 기반 멀티에이전트",
    "MCP agent-channel 실시간 통신",
    "MES: Go + React (PostgreSQL)",
    "ERP: SQL Server (읽기 전용)",
    "벡터DB: pgvector (Qwen3-8B)",
    "텔레그램 · 슬랙 양방향 연동",
    "대시보드: Flask + SocketIO",
]
for s in stack:
    p = rtf.add_paragraph(); p.space_before = Pt(4)
    r = p.add_run(); r.text = "• " + s; r.font.size = Pt(11); r.font.color.rgb = C_DARK

add_footer(sl, FOOTER_TEXT, 3, TOTAL)

# ══════════════════════════════════════════════════
# 슬라이드 4: 비교 테이블
# ══════════════════════════════════════════════════
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

sl = prs.slides.add_slide(blank_layout(prs))
add_bg(sl, "bg-slide-header.jpeg")
add_header(sl, "AI 라이선스 통합 옵션 비교", "전사 AI 도입 방안 · 비용 절감 시나리오")

rows = 7; cols = 6
table = sl.shapes.add_table(rows, cols,
    Inches(0.3), Inches(1.45), Inches(12.7), Inches(5.3)).table

# 열 너비
widths = [Inches(1.6), Inches(2.2), Inches(1.8), Inches(2.0), Inches(2.2), Inches(2.2)]
for i, w in enumerate(widths): table.columns[i].width = w

# 행 높이
for r in range(rows): table.rows[r].height = Inches(5.3 / rows)

def tc(table, r, c, text, bold=False, bg=None, fg=None, align=PP_ALIGN.CENTER, size=11):
    cell = table.cell(r, c)
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if fg: run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

# 헤더행
tc(table, 0, 0, "구분", bold=True, bg=C_NAVY, fg=C_WHITE)
tc(table, 0, 1, "", bold=True, bg=C_NAVY, fg=C_WHITE)
tc(table, 0, 2, "현행", bold=True, bg=C_NAVY, fg=C_WHITE)
tc(table, 0, 3, "안 A\nRovo + CC Team", bold=True, bg=C_BLUE, fg=C_WHITE, size=10)
tc(table, 0, 4, "안 B\nCC Team 단독", bold=True, bg=C_LBLUE, fg=C_WHITE, size=10)
tc(table, 0, 5, "안 C\nRovo 단독", bold=True, bg=C_ORANGE, fg=C_WHITE, size=10)

data = [
    ("비용(월)", "", "약 450만원", "약 210만원", "약 280만원", "약 190만원",
     None, None, RGBColor(0xF0,0xFD,0xF4), None, None, RGBColor(0xFF,0xF7,0xED)),
    ("전사 AI", "Jira/Conf AI", "Jira 유지보수", "Rovo 통합", "미지원", "Rovo 통합",
     None, None, RGBColor(0xF0,0xFD,0xF4), RGBColor(0xFF,0xF1,0xF2), RGBColor(0xF0,0xFD,0xF4), None),
    ("", "메일/문서 AI", "미운영", "Rovo", "미지원", "Rovo",
     None, None, RGBColor(0xFF,0xF1,0xF2), RGBColor(0xF0,0xFD,0xF4), RGBColor(0xFF,0xF1,0xF2), RGBColor(0xF0,0xFD,0xF4)),
    ("", "일반 직원 AI", "미운영", "Rovo 59명", "CC Team 59명", "Rovo 59명",
     None, None, RGBColor(0xFF,0xF1,0xF2), RGBColor(0xF0,0xFD,0xF4), RGBColor(0xFF,0xF7,0xED), RGBColor(0xF0,0xFD,0xF4)),
    ("개발 AI", "Claude Code", "Pro 개인", "CC Team", "CC Team", "미포함",
     None, None, None, RGBColor(0xF0,0xFD,0xF4), RGBColor(0xF0,0xFD,0xF4), RGBColor(0xFF,0xF1,0xF2)),
    ("", "ERP 유지보수 해지", "유지", "해지 (-600만)", "해지 (-600만)", "해지 (-600만)",
     None, None, RGBColor(0xFF,0xF1,0xF2), RGBColor(0xF0,0xFD,0xF4), RGBColor(0xF0,0xFD,0xF4), RGBColor(0xF0,0xFD,0xF4)),
]

for ri, row_data in enumerate(data):
    texts = row_data[:6]
    bgs = list(row_data[6:])
    for ci, text in enumerate(texts):
        bg = bgs[ci] if ci < len(bgs) else None
        align = PP_ALIGN.LEFT if ci < 2 else PP_ALIGN.CENTER
        tc(table, ri+1, ci, text, bg=bg if bg else RGBColor(0xF5,0xF7,0xFC) if ci < 2 else None,
           align=align, size=10)

add_footer(sl, FOOTER_TEXT, 4, TOTAL)

# ══════════════════════════════════════════════════
# 슬라이드 5: 검토 의견 그리드 (2×2)
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
add_bg(sl, "bg-slide-header.jpeg")
add_header(sl, "도입 방안 검토 의견", "안별 장단점 · 팀 의견 종합")

cards = [
    ("안 A — Rovo + CC Team", C_BLUE,
     ["전사 AI + 개발 AI 동시 커버", "Jira/Confluence Rovo 통합", "연 ~2,000만원 절감 가능"],
     ["두 라이선스 체계 관리 필요", "Rovo 초기 세팅 공수"]),
    ("안 B — CC Team 단독", C_LBLUE,
     ["단일 라이선스 관리 단순", "Claude Code 최대 활용"],
     ["Jira AI 기능 미지원", "일반 직원 AI 경험 제한", "Rovo 대비 협업 기능 부족"]),
    ("안 C — Rovo 단독", C_ORANGE,
     ["월 비용 최저 (~190만원)", "전사 Atlassian 통합 최강", "59명 모두 AI 접근 가능"],
     ["Claude Code 미포함", "개발팀 생산성 저하 우려", "AI 에이전트 연동 제한"]),
    ("권고사항", C_GREEN,
     ["단기: 안 A (Rovo + CC Team) 추진", "ERP 유지보수 조기 해지로 재원 확보", "분기별 ROI 측정 후 조정"],
     ["전환 기간 중 병행 운영 필요", "Rovo 도입 전 Jira 정리 선행 권고"]),
]

positions = [
    (Inches(0.3), Inches(1.4)),
    (Inches(6.75), Inches(1.4)),
    (Inches(0.3), Inches(4.4)),
    (Inches(6.75), Inches(4.4)),
]
cw, ch = Inches(6.2), Inches(2.75)

for (cx, cy), (title, accent, pros, cons) in zip(positions, cards):
    box = sl.shapes.add_shape(1, cx, cy, cw, ch)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFC)
    box.line.color.rgb = accent
    # 상단 강조선
    line = sl.shapes.add_shape(1, cx, cy, cw, Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = accent; line.line.fill.background()
    # 텍스트
    tb = sl.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.1), cw - Inches(0.3), ch - Inches(0.15))
    tf = tb.text_frame; tf.word_wrap = True
    hp = tf.paragraphs[0]; hr = hp.add_run()
    hr.text = title; hr.font.bold = True; hr.font.size = Pt(12); hr.font.color.rgb = C_DARK
    for item in pros:
        p = tf.add_paragraph(); r = p.add_run()
        r.text = "✓ " + item; r.font.size = Pt(10); r.font.color.rgb = C_GREEN
    for item in cons:
        p = tf.add_paragraph(); r = p.add_run()
        r.text = "✗ " + item; r.font.size = Pt(10); r.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)

add_footer(sl, FOOTER_TEXT, 5, TOTAL)

# ══════════════════════════════════════════════════
# 슬라이드 6: Phase 3 로드맵 (2단)
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
add_bg(sl, "bg-slide-header.jpeg")
add_header(sl, "Phase 3 확장 로드맵", "2026년 2분기 ~ 3분기")

q2_items = [
    "슬랙 채널별 AI 팀원 배치 완료",
    "dev-agent 온라인 전환",
    "sales-agent 활성화",
    "MES 품질검사 모듈 Go 마이그레이션",
    "CS RAG 부품 매뉴얼 1,000건 임베딩",
]
q3_items = [
    "전사 AI 데이터 채널 확보",
    "schedule-agent 납기 자동 알림",
    "MES 모바일 앱 (React Native)",
    "Jira 해지 → Rovo 이슈 관리 전환",
    "AI ROI 중간 보고",
]

for items, left, title in [
    (q2_items, Inches(0.4), "2분기 (4~6월) 목표"),
    (q3_items, Inches(6.9), "3분기 (7~9월) 목표"),
]:
    box = sl.shapes.add_shape(1, left, Inches(1.4), Inches(6.0), Inches(5.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFC)
    box.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    tb = sl.shapes.add_textbox(left + Inches(0.2), Inches(1.5), Inches(5.7), Inches(5.3))
    tf = tb.text_frame; tf.word_wrap = True
    hp = tf.paragraphs[0]; hr = hp.add_run()
    hr.text = title; hr.font.bold = True; hr.font.size = Pt(13); hr.font.color.rgb = C_NAVY
    for item in items:
        p = tf.add_paragraph(); p.space_before = Pt(5)
        r = p.add_run(); r.text = "• " + item; r.font.size = Pt(11); r.font.color.rgb = C_DARK

add_footer(sl, FOOTER_TEXT, 6, TOTAL)

# ══════════════════════════════════════════════════
# 슬라이드 7: 요약 박스 (결론)
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
add_bg(sl, "bg-slide-header.jpeg")
add_header(sl, "결론 및 권고사항", "의사결정 요약")

sbox = sl.shapes.add_shape(1, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.5))
sbox.fill.solid(); sbox.fill.fore_color.rgb = RGBColor(0x44, 0x54, 0x6A)
sbox.line.fill.background()

stb = sl.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.1))
stf = stb.text_frame; stf.word_wrap = True

sh = stf.paragraphs[0]; shr = sh.add_run()
shr.text = "핵심 결론"; shr.font.bold = True; shr.font.size = Pt(22); shr.font.color.rgb = C_WHITE

summary_lines = [
    "WTA AI 에이전트 시스템은 6명 상시 운영 체계로 안정화되었으며, 품질·CS·개발·인프라 영역에서 24시간 자동화를 실현하고 있습니다.",
    "",
    "라이선스 통합(안 A: Rovo + CC Team)을 통해 연 ~2,000만원 절감이 가능하며, ERP 유지보수 해지와 병행 추진을 권고합니다.",
    "",
    "Phase 3에서는 슬랙 채널 AI 배치 및 전사 데이터 채널 구축을 통해 현실 직원과 AI 에이전트 간 데이터 흐름을 완성합니다.",
]
for line in summary_lines:
    p = stf.add_paragraph(); p.space_before = Pt(6)
    r = p.add_run(); r.text = line
    r.font.size = Pt(13); r.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)

add_footer(sl, FOOTER_TEXT, 7, TOTAL)

# ══════════════════════════════════════════════════
# 슬라이드 8: 마지막 (bg-ending.jpeg)
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
add_bg(sl, "bg-ending.jpeg")

# 저장
prs.save(OUT)
print(f"저장 완료: {OUT}")
