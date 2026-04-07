# -*- coding: utf-8 -*-
"""범용 WTA 슬라이드 HTML → PPTX 변환기

사용법:
  PPTX_HTML_PATH=<html> PPTX_OUT_PATH=<out.pptx> python gen_pptx.py
  또는: python gen_pptx.py <html_path> [out_path]

공통 구조:
  .slide.cover        → 표지 (h1, .sub, .meta)
  .slide.content-slide → 내용 (.s-header > .s-title, .s-body 내 p/ul/table/div.box/div.vs/div.big 등)
  .slide.ending       → 엔딩 (빈 슬라이드, 배경만)
"""
import sys
import os
import re

sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from bs4 import BeautifulSoup, Tag

# ── 경로 설정 ──
if len(sys.argv) >= 2:
    html_path = sys.argv[1]
    out_path = sys.argv[2] if len(sys.argv) >= 3 else html_path.rsplit('.', 1)[0] + '.pptx'
else:
    html_path = os.environ.get('PPTX_HTML_PATH', r'C:\MES\wta-agents\reports\MAX\slide-cs-chatbot.html')
    out_path = os.environ.get('PPTX_OUT_PATH', html_path.rsplit('.', 1)[0] + '.pptx')

img_dir = os.environ.get('PPTX_IMG_DIR', r'C:\MES\wta-agents\reports\MAX\template-images')

# ── 색상 ──
WTA_RED = RGBColor(0xCC, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x88, 0x88, 0x88)
DARK = RGBColor(0x22, 0x22, 0x22)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
BLUE = RGBColor(0x15, 0x65, 0xC0)
BODY_COLOR = RGBColor(0x44, 0x44, 0x44)

COLOR_MAP = {
    'red': WTA_RED,
    'green': GREEN,
    'orange': ORANGE,
    'blue': BLUE,
}

# ── PPTX 초기화 ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 레이아웃 영역 상수
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
CONTENT_W = Inches(11.7)
TITLE_TOP = Inches(0.35)
BODY_TOP = Inches(1.3)


# ── 헬퍼 함수 ──
def add_bg_image(slide, img_name):
    """슬라이드 배경에 이미지 추가"""
    img_path = os.path.join(img_dir, img_name)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, 0, 0, SLIDE_W, SLIDE_H)


def set_run(run, text, size=14, bold=False, color=BLACK, font_name='맑은 고딕'):
    """Run 속성 설정"""
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    """텍스트박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = '맑은 고딕'
    p.alignment = alignment
    return txBox


def add_rich_paragraph(tf, el, font_size=14, color=BODY_COLOR, is_first=False):
    """HTML 요소를 rich text paragraph로 변환 (strong, .hl, .safe, .accent 처리)"""
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
    p.space_after = Pt(3)

    # 인라인 자식이 없으면 단순 텍스트
    if not el.find_all(['strong', 'span', 'em', 'b', 'a']):
        run = p.add_run()
        set_run(run, el.get_text(strip=True), size=font_size, color=color)
        return p

    for child in el.children:
        if isinstance(child, str):
            text = child.strip()
            if text:
                run = p.add_run()
                set_run(run, text, size=font_size, color=color)
        elif isinstance(child, Tag):
            text = child.get_text(strip=True)
            if not text:
                continue
            classes = child.get('class', [])
            run = p.add_run()
            if child.name in ('strong', 'b') or 'hl' in classes:
                set_run(run, text, size=font_size, bold=True, color=WTA_RED)
            elif 'safe' in classes:
                set_run(run, text, size=font_size, color=GREEN)
            elif 'accent' in classes:
                set_run(run, text, size=font_size, bold=True, color=WTA_RED)
            elif child.name == 'em':
                set_run(run, text, size=font_size, color=GRAY)
            else:
                set_run(run, text, size=font_size, color=color)
    return p


def add_table(slide, left, top, width, rows_data):
    """표 추가"""
    n_rows = len(rows_data)
    n_cols = max(len(r) for r in rows_data) if rows_data else 0
    if n_rows == 0 or n_cols == 0:
        return None
    row_h = Inches(0.35)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows)
    table = table_shape.table

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            if c >= n_cols:
                break
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11 if r == 0 else 12)
                para.font.name = '맑은 고딕'
                para.font.color.rgb = RGBColor(0x55, 0x55, 0x55) if r == 0 else BLACK
                para.font.bold = (r == 0)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    return table_shape


def add_box(slide, left, top, width, height, title, lines, accent_color=BLUE):
    """박스 (왼쪽 액센트선 + 제목 + 내용)"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)

    border = slide.shapes.add_shape(1, left, top, Inches(0.06), height)
    border.fill.solid()
    border.fill.fore_color.rgb = accent_color
    border.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = '맑은 고딕'

    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = BLACK
        p.font.name = '맑은 고딕'
        p.space_before = Pt(2)
    return shape


def detect_box_color(el):
    """div.box의 색상 클래스 감지"""
    classes = el.get('class', [])
    for c in classes:
        if c in COLOR_MAP:
            return COLOR_MAP[c]
    return BLUE


def extract_text_lines(el):
    """요소 내 텍스트를 줄 단위로 추출"""
    lines = []
    for child in el.find_all(['p', 'li'], recursive=True):
        text = child.get_text(strip=True)
        if text:
            prefix = '• ' if child.name == 'li' else ''
            lines.append(prefix + text)
    if not lines:
        text = el.get_text(strip=True)
        if text:
            lines = [text]
    return lines


def parse_html_table(table_el):
    """HTML <table>에서 행 데이터 추출"""
    rows = []
    for tr in table_el.find_all('tr'):
        cells = [td.get_text(strip=True) for td in tr.find_all(['th', 'td'])]
        if cells:
            rows.append(cells)
    return rows


# ── 슬라이드 렌더 함수 ──
def render_cover(slide_el, pptx_slide):
    """표지 슬라이드 렌더링"""
    add_bg_image(pptx_slide, 'image2.jpeg')

    h1 = slide_el.find('h1')
    sub = slide_el.find(class_='sub')
    meta = slide_el.find(class_='meta')

    if h1:
        title_text = h1.get_text(strip=True)
        # 긴 제목은 줄바꿈 처리
        if len(title_text) > 25:
            mid = len(title_text) // 2
            # 공백이나 구분자 근처에서 줄바꿈
            break_pos = title_text.rfind(' ', 0, mid + 10)
            if break_pos < mid - 10:
                break_pos = mid
            title_text = title_text[:break_pos].strip() + '\n' + title_text[break_pos:].strip()
        add_textbox(pptx_slide, Inches(1.2), Inches(1.8), Inches(9), Inches(1.5),
                    title_text, font_size=36, bold=True, color=BLACK)

    if sub:
        add_textbox(pptx_slide, Inches(1.2), Inches(3.5), Inches(9), Inches(0.7),
                    sub.get_text(strip=True), font_size=18, color=RGBColor(0x66, 0x66, 0x66))

    if meta:
        meta_lines = [line.strip() for line in meta.get_text().split('\n') if line.strip()]
        txBox = pptx_slide.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(7), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(meta_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.font.name = '맑은 고딕'


def render_ending(slide_el, pptx_slide):
    """엔딩 슬라이드 렌더링 (배경만)"""
    add_bg_image(pptx_slide, 'image5.jpeg')


def render_content(slide_el, pptx_slide):
    """내용 슬라이드 범용 렌더링"""
    add_bg_image(pptx_slide, 'image1.jpeg')

    # 제목
    title_el = slide_el.find(class_='s-title')
    title_text = title_el.get_text(strip=True) if title_el else ''
    if title_text:
        add_textbox(pptx_slide, Inches(2.5), TITLE_TOP, Inches(9), Inches(0.7),
                    title_text, font_size=25, bold=True, color=BLACK)

    # 본문 (s-body)
    body = slide_el.find(class_='s-body')
    if not body:
        return

    # body 자식 요소들을 순회하며 렌더링
    cursor_y = BODY_TOP
    children = [c for c in body.find_all(recursive=False) if isinstance(c, Tag)]

    for child in children:
        classes = child.get('class', [])
        tag = child.name

        # .title-line (부제목 영역)
        if 'title-line' in classes:
            lead = child.find(class_='lead')
            date = child.find(class_='date')
            if lead:
                add_textbox(pptx_slide, MARGIN_L, cursor_y, CONTENT_W, Inches(0.4),
                            lead.get_text(strip=True), font_size=15, bold=True, color=DARK)
                cursor_y += Inches(0.4)
            if date:
                add_textbox(pptx_slide, MARGIN_L, cursor_y, CONTENT_W, Inches(0.3),
                            date.get_text(strip=True), font_size=11, color=GRAY)
                cursor_y += Inches(0.35)

        # <p> 텍스트
        elif tag == 'p':
            text = child.get_text(strip=True)
            if text:
                # 긴 텍스트는 작은 폰트
                fsize = 13 if len(text) > 150 else 14
                height = Inches(0.3 + (len(text) // 100) * 0.2)
                txBox = pptx_slide.shapes.add_textbox(MARGIN_L, cursor_y, CONTENT_W, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                add_rich_paragraph(tf, child, font_size=fsize, is_first=True)
                cursor_y += height + Inches(0.05)

        # <ul>/<ol> 리스트
        elif tag in ('ul', 'ol'):
            items = child.find_all('li', recursive=False)
            height = Inches(0.25 * len(items) + 0.1)
            txBox = pptx_slide.shapes.add_textbox(MARGIN_L, cursor_y, CONTENT_W, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, li in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = '• ' + li.get_text(strip=True)
                p.font.size = Pt(13)
                p.font.color.rgb = BODY_COLOR
                p.font.name = '맑은 고딕'
                p.space_after = Pt(2)
            cursor_y += height + Inches(0.05)

        # <table>
        elif tag == 'table':
            rows_data = parse_html_table(child)
            if rows_data:
                table_h = Inches(0.32 * len(rows_data))
                add_table(pptx_slide, MARGIN_L, cursor_y, CONTENT_W, rows_data)
                cursor_y += table_h + Inches(0.1)

        # div.box
        elif 'box' in classes:
            accent = detect_box_color(child)
            h4 = child.find('h4')
            box_title = h4.get_text(strip=True) if h4 else ''
            lines = extract_text_lines(child)
            # h4 텍스트를 lines에서 제거
            if box_title and lines and lines[0] == box_title:
                lines = lines[1:]
            box_h = Inches(0.35 + 0.22 * len(lines))
            box_h = min(box_h, Inches(2.0))
            add_box(pptx_slide, MARGIN_L, cursor_y, CONTENT_W, box_h,
                    box_title, lines, accent)
            cursor_y += box_h + Inches(0.1)

        # div.vs (Before/After 비교)
        elif 'vs' in classes:
            cols = child.find_all(class_=re.compile(r'vs-col'))
            if len(cols) >= 2:
                half_w = Inches(5.7)
                for ci, col in enumerate(cols[:2]):
                    is_bad = 'bad' in col.get('class', [])
                    accent = RGBColor(0xC6, 0x28, 0x28) if is_bad else GREEN
                    h4 = col.find('h4')
                    col_title = h4.get_text(strip=True) if h4 else ('Before' if is_bad else 'After')
                    col_lines = extract_text_lines(col)
                    if col_title and col_lines and col_lines[0] == col_title:
                        col_lines = col_lines[1:]
                    box_h = Inches(0.4 + 0.22 * len(col_lines))
                    box_h = min(box_h, Inches(2.5))
                    x = MARGIN_L if ci == 0 else MARGIN_L + half_w + Inches(0.3)
                    add_box(pptx_slide, x, cursor_y, half_w, box_h,
                            col_title, col_lines, accent)
                vs_h = box_h
                cursor_y += vs_h + Inches(0.1)

        # div.big (큰 숫자)
        elif 'big' in classes:
            num_el = child.find(class_='num')
            label_el = child.find(class_='label')
            if num_el:
                num_classes = num_el.get('class', [])
                num_color = WTA_RED
                for c in num_classes:
                    if c in COLOR_MAP:
                        num_color = COLOR_MAP[c]
                add_textbox(pptx_slide, MARGIN_L, cursor_y, CONTENT_W, Inches(0.7),
                            num_el.get_text(strip=True), font_size=40, bold=True,
                            color=num_color, alignment=PP_ALIGN.CENTER)
                cursor_y += Inches(0.7)
            if label_el:
                add_textbox(pptx_slide, MARGIN_L, cursor_y, CONTENT_W, Inches(0.3),
                            label_el.get_text(strip=True), font_size=13,
                            color=GRAY, alignment=PP_ALIGN.CENTER)
                cursor_y += Inches(0.35)

        # div.g2, div.g3, div.g4 (그리드 → 박스 나열)
        elif any(g in classes for g in ('g2', 'g3', 'g4')):
            grid_children = [gc for gc in child.find_all(recursive=False) if isinstance(gc, Tag)]
            n_cols_grid = 2
            if 'g3' in classes:
                n_cols_grid = 3
            elif 'g4' in classes:
                n_cols_grid = 4
            col_w = Inches(11.7 / n_cols_grid - 0.15)
            for gi, gc in enumerate(grid_children):
                col_idx = gi % n_cols_grid
                row_idx = gi // n_cols_grid
                x = MARGIN_L + Inches((11.7 / n_cols_grid) * col_idx)
                y = cursor_y + Inches(row_idx * 1.6)

                gc_classes = gc.get('class', [])
                if 'big' in gc_classes:
                    # 큰 숫자 그리드 아이템
                    num_el = gc.find(class_='num')
                    label_el = gc.find(class_='label')
                    if num_el:
                        num_color = WTA_RED
                        for c in num_el.get('class', []):
                            if c in COLOR_MAP:
                                num_color = COLOR_MAP[c]
                        add_textbox(pptx_slide, x, y, col_w, Inches(0.6),
                                    num_el.get_text(strip=True), font_size=36, bold=True,
                                    color=num_color, alignment=PP_ALIGN.CENTER)
                    if label_el:
                        add_textbox(pptx_slide, x, y + Inches(0.6), col_w, Inches(0.3),
                                    label_el.get_text(strip=True), font_size=12,
                                    color=GRAY, alignment=PP_ALIGN.CENTER)
                elif 'box' in gc_classes:
                    accent = detect_box_color(gc)
                    h4 = gc.find('h4')
                    gc_title = h4.get_text(strip=True) if h4 else ''
                    gc_lines = extract_text_lines(gc)
                    if gc_title and gc_lines and gc_lines[0] == gc_title:
                        gc_lines = gc_lines[1:]
                    add_box(pptx_slide, x, y, col_w, Inches(1.4), gc_title, gc_lines, accent)
                else:
                    # 일반 그리드 아이템 → 텍스트 박스
                    gc_title_el = gc.find(['h4', 'h3', 'strong'])
                    gc_title = gc_title_el.get_text(strip=True) if gc_title_el else ''
                    gc_lines = extract_text_lines(gc)
                    if gc_title and gc_lines and gc_lines[0] == gc_title:
                        gc_lines = gc_lines[1:]
                    add_box(pptx_slide, x, y, col_w, Inches(1.4), gc_title, gc_lines, BLUE)

            # cursor_y 업데이트
            total_rows = (len(grid_children) + n_cols_grid - 1) // n_cols_grid
            cursor_y += Inches(total_rows * 1.6 + 0.1)

        # div.flow (흐름도)
        elif 'flow' in classes:
            nodes = child.find_all(class_=re.compile(r'd-node|flow-node|f-node'))
            if nodes:
                n = len(nodes)
                node_w = Inches(min(2.5, 11.0 / n))
                gap = Inches(0.2)
                for ni, node in enumerate(nodes):
                    x = MARGIN_L + (node_w + gap) * ni
                    node_title_el = node.find(['strong', 'h4', 'b'])
                    node_title = node_title_el.get_text(strip=True) if node_title_el else ''
                    node_lines = extract_text_lines(node)
                    if node_title and node_lines and node_lines[0] == node_title:
                        node_lines = node_lines[1:]
                    # 노드 박스
                    shape = pptx_slide.shapes.add_shape(1, x, cursor_y, node_w, Inches(0.9))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
                    shape.line.color.rgb = BLUE
                    shape.line.width = Pt(1.5)
                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.08)
                    p = tf.paragraphs[0]
                    p.text = node_title
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = DARK
                    p.font.name = '맑은 고딕'
                    p.alignment = PP_ALIGN.CENTER
                    for line in node_lines[:2]:
                        p2 = tf.add_paragraph()
                        p2.text = line
                        p2.font.size = Pt(9)
                        p2.font.color.rgb = BLACK
                        p2.font.name = '맑은 고딕'
                        p2.alignment = PP_ALIGN.CENTER
                    # 화살표
                    if ni < n - 1:
                        add_textbox(pptx_slide, x + node_w, cursor_y + Inches(0.2),
                                    gap, Inches(0.4), '→', font_size=20, bold=True,
                                    color=WTA_RED, alignment=PP_ALIGN.CENTER)
                cursor_y += Inches(1.05)
            else:
                # 플로우 텍스트만
                text = child.get_text(strip=True)
                if text:
                    add_textbox(pptx_slide, MARGIN_L, cursor_y, CONTENT_W, Inches(0.4),
                                text[:200], font_size=12, color=DARK)
                    cursor_y += Inches(0.45)

        # div.sample-grid 또는 기타 이미지 그리드 → 텍스트 설명으로 변환
        elif 'sample-grid' in classes:
            items = child.find_all(class_=re.compile(r'sample-item|grid-item'))
            for gi, item in enumerate(items[:6]):
                caption = item.find(class_=re.compile(r'caption|label|title'))
                if caption:
                    add_textbox(pptx_slide, MARGIN_L + Inches((gi % 3) * 4),
                                cursor_y + Inches((gi // 3) * 0.35),
                                Inches(3.8), Inches(0.3),
                                caption.get_text(strip=True), font_size=11, color=DARK)
            rows = (min(len(items), 6) + 2) // 3
            cursor_y += Inches(rows * 0.4 + 0.1)

        # div.diagram (도식)
        elif 'diagram' in classes:
            nodes = child.find_all(class_='d-node')
            if nodes:
                for ni, node in enumerate(nodes[:8]):
                    text = node.get_text(strip=True)
                    x = MARGIN_L + Inches((ni % 4) * 3)
                    y = cursor_y + Inches((ni // 4) * 0.5)
                    add_textbox(pptx_slide, x, y, Inches(2.8), Inches(0.4),
                                text[:60], font_size=11, color=DARK)
                rows = (min(len(nodes), 8) + 3) // 4
                cursor_y += Inches(rows * 0.55)

        # 기타 div → 텍스트 추출 시도
        else:
            text = child.get_text(strip=True)
            if text and len(text) > 5:
                fsize = 12 if len(text) > 200 else 13
                height = Inches(0.3 + min(len(text) // 120, 3) * 0.2)
                add_textbox(pptx_slide, MARGIN_L, cursor_y, CONTENT_W, height,
                            text[:500], font_size=fsize, color=BODY_COLOR)
                cursor_y += height + Inches(0.05)


# ── 메인 변환 ──
print(f"HTML: {html_path}")
print(f"OUT:  {out_path}")

with open(html_path, 'r', encoding='utf-8') as f:
    soup = BeautifulSoup(f.read(), 'html.parser')

slides = soup.find_all('div', class_='slide')
print(f"슬라이드 {len(slides)}장 감지")

for i, slide_el in enumerate(slides):
    classes = slide_el.get('class', [])
    pptx_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    if 'cover' in classes:
        render_cover(slide_el, pptx_slide)
        print(f"  [{i+1}] 표지 렌더링 완료")
    elif 'ending' in classes:
        render_ending(slide_el, pptx_slide)
        print(f"  [{i+1}] 엔딩 렌더링 완료")
    elif 'content-slide' in classes:
        title_el = slide_el.find(class_='s-title')
        title = title_el.get_text(strip=True) if title_el else '(제목 없음)'
        render_content(slide_el, pptx_slide)
        print(f"  [{i+1}] 내용: {title}")
    else:
        # 미분류 → 내용 슬라이드로 시도
        render_content(slide_el, pptx_slide)
        print(f"  [{i+1}] 기타 슬라이드")

prs.save(out_path)
size_kb = os.path.getsize(out_path) / 1024
print(f"\nPPTX 저장 완료: {out_path} ({size_kb:.0f}KB)")
